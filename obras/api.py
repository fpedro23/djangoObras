import json
from datetime import *

from django.db.models import Q, Sum
from django.http import HttpResponse
from oauth2_provider.models import AccessToken
from oauth2_provider.views import ProtectedResourceView
from django.db.models import Count

from obras.BuscarObras import BuscarObras
from obras.models import Obra, Estado, Dependencia, Impacto, TipoClasificacion, TipoInversion, TipoObra, Inaugurador, \
    InstanciaEjecutora, get_subdependencias_as_list_flat, Municipio
from obras.views import get_array_or_none
import pytz

from pptx import Presentation
from pptx.util import Inches
from pptx.util import Pt
from pptx.dml.color import RGBColor

try:
    import cStringIO as StringIO
except ImportError:
    import StringIO
from xlsxwriter.workbook import Workbook
from django.core.servers.basehttp import FileWrapper
from django.http import StreamingHttpResponse


def from_utc_to_local(utc_dt):
    local_tz = pytz.timezone('America/Mexico_City')
    local_dt = utc_dt.replace(tzinfo=pytz.utc).astimezone(local_tz)
    return local_tz.normalize(local_dt)


def get_usuario_for_token(token):
    if token:
        return AccessToken.objects.get(token=token).user.usuario
    else:
        return None


class HoraUltimaActualizacion(ProtectedResourceView):
    def get(self, request):
        json_response = {}
        dependencia = Dependencia.objects.all().order_by('fecha_ultima_modificacion').last()

        if dependencia is not None:
            date = from_utc_to_local(dependencia.fecha_ultima_modificacion)
        else:
            date = datetime.now()

        if date.day >= 10:
            json_response['dia'] = str(date.day)
        else:
            json_response['dia'] = "0" + str(date.day)
        if date.month >= 10:
            json_response['mes'] = str(date.month)
        else:
            json_response['mes'] = "0" + str(date.month)
        json_response['ano'] = str(date.year)

        time = date.time()
        if time.hour >= 10:
            json_response['hora'] = str(time.hour)
        else:
            json_response['hora'] = "0" + str(time.hour)
        if time.minute >= 10:
            json_response['minuto'] = str(time.minute)
        else:
            json_response['minuto'] = "0" + str(time.minute)
        if time.second >= 10:
            json_response['segundo'] = str(time.second)
        else:
            json_response['segundo'] = "0" + str(time.second)
        return HttpResponse(json.dumps(json_response), 'application/json')


class HoraEndpoint(ProtectedResourceView):
    def get(self, request):
        json_response = {}
        dependencia = Dependencia.objects.all().order_by('fecha_ultima_modificacion').last()

        if dependencia is not None:
            date = from_utc_to_local(dependencia.fecha_ultima_modificacion)
        else:
            date = datetime.now()

        if date.day >= 10:
            json_response['dia'] = str(date.day)
        else:
            json_response['dia'] = "0" + str(date.day)
        if date.month >= 10:
            json_response['mes'] = str(date.month)
        else:
            json_response['mes'] = "0" + str(date.month)
        json_response['ano'] = str(date.year)

        time = date.time()
        if time.hour >= 10:
            json_response['hora'] = str(time.hour)
        else:
            json_response['hora'] = "0" + str(time.hour)
        if time.minute >= 10:
            json_response['minuto'] = str(time.minute)
        else:
            json_response['minuto'] = "0" + str(time.minute)
        if time.second >= 10:
            json_response['segundo'] = str(time.second)
        else:
            json_response['segundo'] = "0" + str(time.second)
        return HttpResponse(json.dumps(json_response), 'application/json')


class ObrasIniciadasPptxEndpoint(ProtectedResourceView):
    def get(self, request):

        usuario = get_usuario_for_token(request.GET.get('access_token'))

        query = Q(fechaInicio__lte=datetime.now().date()) & Q(tipoObra_id=1)
        if not (usuario.rol == 'SA'):
            subdependencias = get_subdependencias_as_list_flat(usuario.dependencia.all())
            query = query & (Q(dependencia__in=subdependencias) | Q(subdependencia__in=subdependencias))
        obras = Obra.objects.filter(query)

        contador = 1
        json_map = {}
        json_map['obras'] = []
        for obra in obras.values('id', 'identificador_unico', 'estado__nombreEstado', 'denominacion'):
            contador = contador + 1
            json_map['obras'].append(obra)



        output = StringIO.StringIO()
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        shapes = slide.shapes
        shapes.title.text = 'Obras Iniciadas'
        rows = 20
        cols = 3
        left = Inches(0.921)
        top = Inches(1.2)
        width = Inches(6.0)
        height = Inches(0.8)
        table = shapes.add_table(rows, cols, left, top, width, height).table
        # set column widths
        table.columns[0].width = Inches(2.0)
        table.columns[1].width = Inches(2.0)
        table.columns[2].width = Inches(4.0)
        indice = 1

        for obra in json_map['obras']:
            if indice == 20:
                indice = 1
                slide = prs.slides.add_slide(prs.slide_layouts[5])
                shapes = slide.shapes
                shapes.title.text = 'Resultados'

                rows = 20
                cols = 3
                left = Inches(0.921)
                top = Inches(1.2)
                width = Inches(6.0)
                height = Inches(0.8)

                table = shapes.add_table(rows, cols, left, top, width, height).table

                # set column widths
                table.columns[0].width = Inches(2.0)
                table.columns[1].width = Inches(2.0)
                table.columns[2].width = Inches(4.0)

            for x in range(0, 3):
                cell = table.rows[0].cells[x]
                paragraph = cell.textframe.paragraphs[0]
                paragraph.font.size = Pt(12)
                paragraph.font.name = 'Arial Black'
                paragraph.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

            for x in range(0, 3):
                cell = table.rows[indice].cells[x]
                paragraph = cell.textframe.paragraphs[0]
                paragraph.font.size = Pt(8)
                paragraph.font.name = 'Arial'
                paragraph.font.color.rgb = RGBColor(0x0B, 0x0B, 0x0B)

            # write column headings
            table.cell(0, 0).text = 'Identificador'
            table.cell(0, 1).text = 'Estado'
            table.cell(0, 2).text = 'Denominacion'

            # write body cells
            table.cell(indice, 0).text = obra['identificador_unico']
            table.cell(indice, 1).text = obra['estado__nombreEstado']
            table.cell(indice, 2).text = obra['denominacion']
            indice += 1
            # for obra in json_map['obras']:
            #    table.cell(i, 0).text = obra['identificador_unico']
            #    table.cell(i, 1).text = obra['estado__nombreEstado']
            #    table.cell(i, 2).text = obra['denominacion']
            #    i+=1

        prs.save(output)
        response = StreamingHttpResponse(FileWrapper(output),
                                         content_type='application/vnd.openxmlformats-officedocument.presentationml.presentation')
        response['Content-Disposition'] = 'attachment; filename="ObrasIniciadas.pptx"'
        response['Content-Length'] = output.tell()

        output.seek(0)

        return response


class ObrasVencidasPptxEndpoint(ProtectedResourceView):
    def get(self, request):
        usuario = get_usuario_for_token(request.GET.get('access_token'))

        query = Q(fechaTermino__lte=datetime.now().date()) & Q(tipoObra_id=2)
        if not usuario.rol == 'SA':
            subdependencias = get_subdependencias_as_list_flat(usuario.dependencia.all())
            query = query & (Q(dependencia__in=subdependencias) | Q(subdependencia__in=subdependencias))
        obras = Obra.objects.filter(query)

        contador = 1
        json_map = {}
        json_map['obras'] = []
        for obra in obras.values('id', 'identificador_unico', 'estado__nombreEstado', 'denominacion'):
            contador = contador + 1
            json_map['obras'].append(obra)

        output = StringIO.StringIO()
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        shapes = slide.shapes
        shapes.title.text = 'Obras Vencidas'
        rows = 20
        cols = 3
        left = Inches(0.921)
        top = Inches(1.2)
        width = Inches(6.0)
        height = Inches(0.8)
        table = shapes.add_table(rows, cols, left, top, width, height).table
        # set column widths
        table.columns[0].width = Inches(2.0)
        table.columns[1].width = Inches(2.0)
        table.columns[2].width = Inches(4.0)
        indice = 1

        for obra in json_map['obras']:
            if indice == 20:
                indice = 1
                slide = prs.slides.add_slide(prs.slide_layouts[5])
                shapes = slide.shapes
                shapes.title.text = 'Resultados'

                rows = 20
                cols = 3
                left = Inches(0.921)
                top = Inches(1.2)
                width = Inches(6.0)
                height = Inches(0.8)

                table = shapes.add_table(rows, cols, left, top, width, height).table
                # set column widths
                table.columns[0].width = Inches(2.0)
                table.columns[1].width = Inches(2.0)
                table.columns[2].width = Inches(4.0)

            for x in range(0, 3):
                cell = table.rows[0].cells[x]
                paragraph = cell.textframe.paragraphs[0]
                paragraph.font.size = Pt(12)
                paragraph.font.name = 'Arial Black'
                paragraph.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

            for x in range(0, 3):
                cell = table.rows[indice].cells[x]
                paragraph = cell.textframe.paragraphs[0]
                paragraph.font.size = Pt(8)
                paragraph.font.name = 'Arial'
                paragraph.font.color.rgb = RGBColor(0x0B, 0x0B, 0x0B)

            # write column headings
            table.cell(0, 0).text = 'Identificador'
            table.cell(0, 1).text = 'Estado'
            table.cell(0, 2).text = 'Denominacion'

            # write body cells
            table.cell(indice, 0).text = obra['identificador_unico']
            table.cell(indice, 1).text = obra['estado__nombreEstado']
            table.cell(indice, 2).text = obra['denominacion']
            indice += 1
            # for obra in json_map['obras']:
            #    table.cell(i, 0).text = obra['identificador_unico']
            #    table.cell(i, 1).text = obra['estado__nombreEstado']
            #    table.cell(i, 2).text = obra['denominacion']
            #    i+=1

        prs.save(output)
        response = StreamingHttpResponse(FileWrapper(output),
                                         content_type='application/vnd.openxmlformats-officedocument.presentationml.presentation')
        response['Content-Disposition'] = 'attachment; filename="ObrasVencidas.pptx"'
        response['Content-Length'] = output.tell()

        output.seek(0)

        return response


class ReporteNoTrabajoPptxEndpoint(ProtectedResourceView):
    def get(self, request):
        comp_date = datetime.now() - timedelta(15)
        print "****"
        print comp_date

        dependencias = Dependencia.objects.filter(
            Q(obraoprograma='O') &
            Q(fecha_ultima_modificacion__lt=comp_date))

        contador = 1
        json_map = {}
        json_map['obras'] = []
        for obra in dependencias.values('nombreDependencia', 'fecha_ultima_modificacion'):
            contador = contador + 1
            json_map['obras'].append(obra)

        output = StringIO.StringIO()
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        shapes = slide.shapes
        shapes.title.text = 'Bitacora por Dependencias'
        rows = 17
        cols = 2
        left = Inches(0.921)
        top = Inches(1.2)
        width = Inches(6.0)
        height = Inches(0.8)
        table = shapes.add_table(rows, cols, left, top, width, height).table
        # set column widths
        table.columns[0].width = Inches(2.0)
        table.columns[1].width = Inches(2.0)
        indice = 1

        for x in range(0, 2):
            cell = table.rows[0].cells[x]
            paragraph = cell.textframe.paragraphs[0]
            paragraph.font.size = Pt(12)
            paragraph.font.name = 'Arial Black'
            paragraph.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

        # write column headings
        table.cell(0, 0).text = 'Dependencia'
        table.cell(0, 1).text = 'Fecha ultima modificacion'

        for obra in json_map['obras']:
            for x in range(0, 2):
                cell = table.rows[indice].cells[x]
                paragraph = cell.textframe.paragraphs[0]
                paragraph.font.size = Pt(8)
                paragraph.font.name = 'Arial'
                paragraph.font.color.rgb = RGBColor(0x0B, 0x0B, 0x0B)
            # write body cells
            table.cell(indice, 0).text = obra['nombreDependencia']
            table.cell(indice, 1).text = str(obra['fecha_ultima_modificacion'])

            indice += 1

        prs.save(output)
        response = StreamingHttpResponse(FileWrapper(output),
                                         content_type='application/vnd.openxmlformats-officedocument.presentationml.presentation')
        response['Content-Disposition'] = 'attachment; filename="BitacoraDependencias.pptx"'
        response['Content-Length'] = output.tell()

        output.seek(0)

        return response


class ObrasIniciadasEndpoint(ProtectedResourceView):
    def get(self, request):
        usuario = get_usuario_for_token(request.GET.get('access_token'))

        query = Q(fechaInicio__lte=datetime.now().date()) & Q(tipoObra_id=1)
        if not (usuario.rol == 'SA'):
            subdependencias = get_subdependencias_as_list_flat(usuario.dependencia.all())
            query = query & (Q(dependencia__in=subdependencias) | Q(subdependencia__in=subdependencias))
        obras = Obra.objects.filter(query)

        json_ans = '['
        for obra in obras.values('id', 'identificador_unico', 'estado__nombreEstado', 'denominacion'):
            json_ans += '{"id":"'
            json_ans += str(obra['id'])

            json_ans += '","identificador_unico":"'
            json_ans += obra['identificador_unico']

            json_ans += '","estado__nombreEstado":"'
            json_ans += obra['estado__nombreEstado']

            json_ans += '","denominacion":"'
            json_ans += obra['denominacion']

            json_ans += '"},'
        json_ans = json_ans[:-1]
        json_ans += ']'

        return HttpResponse(json_ans, 'application/json')


class ObrasVencidasEndpoint(ProtectedResourceView):
    def get(self, request):
        usuario = get_usuario_for_token(request.GET.get('access_token'))

        query = Q(fechaTermino__lte=datetime.now().date()) & Q(tipoObra_id=2)
        if not usuario.rol == 'SA':
            subdependencias = get_subdependencias_as_list_flat(usuario.dependencia.all())
            query = query & (Q(dependencia__in=subdependencias) | Q(subdependencia__in=subdependencias))
        obras = Obra.objects.filter(query)

        the_list = []
        for obra in obras.values('id', 'identificador_unico', 'estado__nombreEstado', 'denominacion'):
            the_list.append(obra)

        return HttpResponse(json.dumps(the_list), 'application/json')


class ObrasForDependenciaEndpoint(ProtectedResourceView):
    def get(self, request):
        usuario = get_usuario_for_token(request.GET.get('access_token'))

        if usuario.rol == 'SA':
            obras = Obra.objects.all()
        else:
            obras = Obra.filter(dependencia__in=get_subdependencias_as_list_flat(usuario.dependencia.all()))

        the_list = []
        for obra in obras.values('id', 'identificador_unico', 'estado__nombreEstado', 'denominacion'):
            the_list.append(obra)

        return HttpResponse(json.dumps(the_list), 'application/json')


class SubDependenciasiPadEndpoint(ProtectedResourceView):
    def get(self, request):
        token = request.GET.get('access_token')
        token_model = AccessToken.objects.get(token=token)

        if token_model.user.usuario.rol == 'SA':
            dicts = map(lambda dependencia: dependencia.to_serializable_dict(), Dependencia.objects.filter(
                Q(obraoprograma='O') & Q(dependienteDe__isnull=False)
            )
                        )

        elif token_model.user.usuario.rol == 'AD' or token_model.user.usuario.rol == 'FU' or token_model.user.usuario.rol == 'UD':
            dicts = map(lambda dependencia: dependencia.to_serializable_dict(), Dependencia.objects.filter(
                Q(id__in=token_model.user.usuario.dependencia.all()) |
                Q(dependienteDe__in=token_model.user.usuario.dependencia.all()))
                        & Q(dependienteDe__isnull=False)
                        )
        else:
            dicts = map(lambda dependencia: dependencia.to_serializable_dict(), Dependencia.objects.filter(
                Q(id__in=token_model.user.usuario.subdependencia.all()))
                        & Q(dependienteDe__isnull=False)
                        )

        return HttpResponse(json.dumps(dicts), 'application/json')


class DependenciasEndpoint(ProtectedResourceView):
    def get(self, request):
        token = request.GET.get('access_token')
        token_model = AccessToken.objects.get(token=token)

        if token_model.user.usuario.rol == 'SA':
            dicts = map(lambda dependencia: dependencia.to_serializable_dict(), Dependencia.objects.filter(
                Q(obraoprograma='O')
                & Q(dependienteDe__isnull=True)
            ))
        elif token_model.user.usuario.rol == 'AD' or token_model.user.usuario.rol == 'FU' or token_model.user.usuario.rol == 'UD':
            dicts = map(lambda dependencia: dependencia.to_serializable_dict(), Dependencia.objects.filter(
                Q(id__in=token_model.user.usuario.dependencia.all()) |
                Q(dependienteDe__in=token_model.user.usuario.dependencia.all())
                & Q(dependienteDe__isnull=True)
            ))
        else:
            dicts = map(lambda dependencia: dependencia.to_serializable_dict(), Dependencia.objects.filter(
                Q(id__in=token_model.user.usuario.subdependencia.all())
                & Q(dependienteDe__isnull=True)
            ))

        return HttpResponse(json.dumps(dicts), 'application/json')


class ImpactosEndpoint(ProtectedResourceView):
    def get(self, request):
        return HttpResponse(json.dumps(map(lambda impacto: impacto.to_serializable_dict(), Impacto.objects.all())),
                            'application/json')


class EstadosEndpoint(ProtectedResourceView):
    def get(self, request):
        return HttpResponse(json.dumps(map(lambda estado: estado.to_serializable_dict(), Estado.objects.all())),
                            'application/json')


class MunicipiosForEstadosEndpoint(ProtectedResourceView):
    def get(self, request):
        estado_ids = get_array_or_none(request.GET.get('estados'))
        all_estados = False

        if estado_ids is None:
            all_estados = True
        else:
            for estado_id in estado_ids:
                if estado_id == 33 or estado_id == 34:
                    all_estados = True
                    break

        if all_estados:
            municipios = Municipio.objects.order_by('nombreMunicipio').all()
        else:
            municipios = Municipio.objects.filter(estado_id__in=estado_ids).order_by('nombreMunicipio').all()


        the_list = []
        for municipio in municipios.values('id', 'nombreMunicipio'):
            the_list.append(municipio)

        return HttpResponse(json.dumps(the_list), 'application/json')


class DependenciasTreeEndpoint(ProtectedResourceView):
    def get(self, request):
        usuario = get_usuario_for_token(request.GET.get('access_token'))

        if usuario.rol == 'SA':
            dependencias = Dependencia.objects.all()
        else:
            dependencias = usuario.dependencia.all()

        ans = []
        for dep in dependencias:
            ans.append(dep.get_tree())

        return HttpResponse(json.dumps(ans), 'application/json')


class DependenciasIdEndpoint(ProtectedResourceView):
    def get(self, request):
        usuario = get_usuario_for_token(request.GET.get('access_token'))
        idDep = request.GET.get('id')

        dependencias = Dependencia.objects.filter(id=idDep).all()

        ans = []
        for dep in dependencias:
            ans.append(dep.get_tree())

        return HttpResponse(json.dumps(ans), 'application/json')


class SubependenciasFlatEndpoint(ProtectedResourceView):
    def get(self, request):
        usuario = get_usuario_for_token(request.GET.get('access_token'))

        if usuario.rol == 'SA':
            dependencias = Dependencia.objects.all()
        else:
            dependencias = usuario.dependencia.all()

        ans = map(lambda dependencia: dependencia.to_serializable_dict(), dependencias)
        for dependencia in dependencias:
            subdeps = dependencia.get_subdeps_flat()
            if subdeps:
                ans.extend(map(lambda dep: dep.to_serializable_dict(), subdeps))

        return HttpResponse(json.dumps(ans), 'application/json')


class Subdependencias_forId_Endpoint(ProtectedResourceView):
    def get(self, request):
        usuario = get_usuario_for_token(request.GET.get('access_token'))
        iddependencias = get_array_or_none(request.GET.get('dependencia'))
        ans = []
        if iddependencias is not None:
            if usuario.rol == 'SA':
                dependencias = Dependencia.objects.filter(
                    Q(id__in=iddependencias)
                )
            else:
                dependencias = usuario.dependencia.filter(
                    Q(id__in=iddependencias)
                )

            subdependencias = Dependencia.objects.filter(dependienteDe__in=dependencias).all()
            if subdependencias is not None and subdependencias.count() > 0:
                ans = map(lambda dep: dep.to_serializable_dict(), subdependencias)
        else:
            if usuario.rol == 'SA':
                ans = map(lambda dep: dep.to_serializable_dict(),
                          Dependencia.objects.filter(dependienteDe__isnull=False).all())
            else:
                ans = map(lambda dep: dep.to_serializable_dict(), usuario.subdependencia.all())

        return HttpResponse(json.dumps(ans), 'application/json')


class ClasificacionEndpoint(ProtectedResourceView):
    def get(self, request):
        if request.GET.get('id', False):
            clasificaciones = TipoClasificacion.objects.filter(subclasificacionDe_id=request.GET.get('id'))
        else:
            clasificaciones = TipoClasificacion.objects.filter(subclasificacionDe_id__isnull=True)

        return HttpResponse(
            json.dumps(map(lambda clasificacion: clasificacion.to_serializable_dict(), clasificaciones)),
            "application/json")


class InversionEndpoint(ProtectedResourceView):
    def get(self, request):
        return HttpResponse(
            json.dumps(map(lambda inversion: inversion.to_serializable_dict(), TipoInversion.objects.all())),
            'application/json')


class InstanciaEjecutoraEndpoint(ProtectedResourceView):
    def get(self, request):
        return HttpResponse(
            json.dumps(map(lambda instancia: instancia.to_serializable_dict(), InstanciaEjecutora.objects.all())),
            'application/json')


class TipoDeObraEndpoint(ProtectedResourceView):
    def get(self, request):
        return HttpResponse(json.dumps(map(lambda tipo: tipo.to_serializable_dict(), TipoObra.objects.all())),
                            'application/json')


class IdUnicoEndpoint(ProtectedResourceView):
    def get(self, request):
        identificador_unico = request.GET.get('identificador_unico', None)
        usuario = AccessToken.objects.get(token=request.GET.get('access_token')).user.usuario
        obra_id = None
        obra = None
        if identificador_unico is not None:
            obra = Obra.objects.filter(identificador_unico=identificador_unico)
            if obra and obra.count() > 0:
                obra_id = obra.first().id
                obra = obra.first()
            else:
                obra = None
        if obra_id is None:
            if identificador_unico is not None:
                return HttpResponse(json.dumps({'id': obra_id, 'error': 'No existe la obra con identificador ' + identificador_unico}), 'application/json')
            else:
                return HttpResponse(json.dumps({'id': obra_id, 'error': 'Debes proporcionar un identificador'}), 'application/json')
        elif obra is None:
            return HttpResponse(json.dumps({'id': obra_id, 'error': 'No existe la obra con identificador ' + identificador_unico}), 'application/json')
        elif usuario.rol == 'US' and obra.subdependencia is None:
            return HttpResponse(json.dumps({'id': obra_id, 'error': 'Privilegios insuficientes'}), 'application/json')
	elif usuario.rol == 'US' and (usuario.subdependencia is None or usuario.subdependencia.count() == 0):
            return HttpResponse(json.dumps({'id': obra_id, 'error': 'Privilegios insuficientes'}), 'application/json')
	elif usuario.rol == 'US' and usuario.subdependencia is not None and obra.subdependencia is not None and not usuario.subdependencia.filter(id=obra.subdependencia.id).exists():
            return HttpResponse(json.dumps({'id': obra_id, 'error': 'Privilegios insuficientes'}), 'application/json')
        else:
            return HttpResponse(json.dumps({'id': obra_id, 'error': None}), 'application/json')


class IdUnicoEndpointIpad(ProtectedResourceView):
    def get(self, request):
        identificador_unico = request.GET.get('identificador_unico', None)
        return HttpResponse(
            json.dumps(map(lambda obra: obra.to_serializable_dict(),
                           Obra.objects.filter(identificador_unico=identificador_unico))),
            'application/json')


class BuscadorEndpoint(ProtectedResourceView):
    def get(self, request):

        user = AccessToken.objects.get(token=request.GET.get('access_token')).user

        buscador = BuscarObras(
            idtipoobra=get_array_or_none(request.GET.get('tipoDeObra')),
            iddependencias=get_array_or_none(request.GET.get('dependencia')),
            subdependencias=get_array_or_none(request.GET.get('subdependencias')),
            estados=get_array_or_none(request.GET.get('estado')),
            clasificaciones=get_array_or_none(request.GET.get('clasificacion')),
            inversiones=get_array_or_none(request.GET.get('tipoDeInversion')),
            inauguradores=get_array_or_none(request.GET.get('inaugurador')),
            impactos=get_array_or_none(request.GET.get('impacto')),
            inaugurada=request.GET.get('inaugurada', None),
            inversion_minima=request.GET.get('inversionMinima', None),
            inversion_maxima=request.GET.get('inversionMaxima', None),
            fecha_inicio_primera=request.GET.get('fechaInicio', None),
            fecha_inicio_segunda=request.GET.get('fechaInicio', None),
            fecha_fin_primera=request.GET.get('fechaFin', None),
            fecha_fin_segunda=request.GET.get('fechaFinSegunda', None),
            denominacion=request.GET.get('denominacion', None),
            instancia_ejecutora=get_array_or_none(request.GET.get('instanciaEjecutora')),
            limite_min=request.GET.get("limiteMin"),
            limite_max=request.GET.get("limiteMax"),
            busqueda_rapida=request.GET.get("busquedaRapida", None),
            id_obra=request.GET.get("idObra", None),
            susceptible_inauguracion=request.GET.get("susceptible", None),
            subclasificacion=get_array_or_none(request.GET.get('subclasificacion')),
            municipios=get_array_or_none(request.GET.get('municipios')),
        )

        buscador.filtrar_dependencias(user)
        resultados = buscador.buscar()

        json_map = {}
        json_map['reporte_dependencia'] = []
        for reporte in resultados['reporte_dependencia']:
            map = {}
            map['dependencia'] = Dependencia.objects.get(
                nombreDependencia=reporte['dependencia__nombreDependencia']).to_serializable_dict()
            map['numero_obras'] = reporte['numero_obras']
            if reporte['sumatotal'] is None:
                map['sumatotal'] = 0
            else:
                map['sumatotal'] = float(reporte['sumatotal'])
            json_map['reporte_dependencia'].append(map)

        json_map['reporte_subdependencia'] = []
        for reporteSub in resultados['reporte_subdependencia']:
            map = {}
            if reporteSub['subdependencia__nombreDependencia'] is not None:
                map['subdependencia'] = Dependencia.objects.get(
                    nombreDependencia=reporteSub['subdependencia__nombreDependencia']).to_serializable_dict()
            else:
                map['subdependencia'] = Dependencia.objects.get(
                    nombreDependencia=reporteSub['dependencia__nombreDependencia']).to_serializable_dict()
            map['numero_obras'] = reporteSub['numero_obras']
            if reporteSub['sumatotal'] is None:
                map['sumatotal'] = 0
            else:
                map['sumatotal'] = float(reporteSub['sumatotal'])
            json_map['reporte_subdependencia'].append(map)

        json_map['obras'] = []
        for obra in resultados['obras'].values('id', 'identificador_unico', 'estado__nombreEstado', 'denominacion',
                                               'latitud', 'longitud', 'dependencia__imagenDependencia'):
            json_map['obras'].append(obra)

        json_map['geolocalizacion_obras'] = []
        for obra in resultados['geolocalizacion_obras'].values('identificador_unico', 'denominacion', 'inversionTotal',
                                                               'latitud', 'longitud'):
            json_map['geolocalizacion_obras'].append(obra)

        json_map['reporte_estado'] = []
        for reporte_estado in resultados['reporte_estado']:
            map = {}
            if reporte_estado['sumatotal'] is None:
                map['sumatotal'] = 0.0

            else:
                map['sumatotal'] = float(reporte_estado['sumatotal'])
            map['estado'] = Estado.objects.get(
                nombreEstado=reporte_estado['estado__nombreEstado']).to_serializable_dict()
            map['numeroObras'] = reporte_estado['numero_obras']

            json_map['reporte_estado'].append(map)

        json_map['reporte_general'] = []
        map = {}
        total = resultados['reporte_general']['total_invertido']['inversionTotal__sum']
        if total is None:
            total = 0.0
        else:
            total = float(total)
        map['total_invertido'] = total

        map['obras_totales'] = resultados['reporte_general']['obras_totales']
        json_map['reporte_general'].append(map)

        return HttpResponse(json.dumps(json_map), 'application/json')


class PptxEndpoint(ProtectedResourceView):
    def get(self, request):

        user = AccessToken.objects.get(token=request.GET.get('access_token')).user

        buscador = BuscarObras(
            idtipoobra=get_array_or_none(request.GET.get('tipoDeObra')),
            iddependencias=get_array_or_none(request.GET.get('dependencia')),
            subdependencias=get_array_or_none(request.GET.get('subdependencias')),
            estados=get_array_or_none(request.GET.get('estado')),
            clasificaciones=get_array_or_none(request.GET.get('clasificacion')),
            inversiones=get_array_or_none(request.GET.get('tipoDeInversion')),
            inauguradores=get_array_or_none(request.GET.get('inaugurador')),
            impactos=get_array_or_none(request.GET.get('impacto')),
            inaugurada=request.GET.get('inaugurada', None),
            inversion_minima=request.GET.get('inversionMinima', None),
            inversion_maxima=request.GET.get('inversionMaxima', None),
            fecha_inicio_primera=request.GET.get('fechaInicio', None),
            fecha_inicio_segunda=request.GET.get('fechaInicio', None),
            fecha_fin_primera=request.GET.get('fechaFin', None),
            fecha_fin_segunda=request.GET.get('fechaFinSegunda', None),
            denominacion=request.GET.get('denominacion', None),
            instancia_ejecutora=get_array_or_none(request.GET.get('instanciaEjecutora')),
            limite_min=request.GET.get("limiteMin"),
            limite_max=request.GET.get("limiteMax"),
            busqueda_rapida=request.GET.get("busquedaRapida", None),
            id_obra=request.GET.get("idObra", None),
            susceptible_inauguracion=request.GET.get("susceptible", None),
            subclasificacion=get_array_or_none(request.GET.get('subclasificacion')),
            municipios=get_array_or_none(request.GET.get('municipios'))
        )

        buscador.filtrar_dependencias(user)
        resultados = buscador.buscar()

        json_map = {}

        json_map['obras'] = []
        for obra in resultados['obras'].values('id', 'identificador_unico', 'estado__nombreEstado', 'denominacion',
                                               'latitud', 'longitud', 'dependencia__imagenDependencia'):
            json_map['obras'].append(obra)

        output = StringIO.StringIO()
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        shapes = slide.shapes
        shapes.title.text = 'Resultados'

        renglones = resultados['reporte_general']['obras_totales'] + 1
        if renglones < 22:
            rows = renglones
        else:
            rows = 22
        cols = 3
        left = Inches(0.921)
        top = Inches(1.2)
        width = Inches(6.0)
        height = Inches(0.8)

        table = shapes.add_table(rows, cols, left, top, width, height).table

        # set column widths
        table.columns[0].width = Inches(2.0)
        table.columns[1].width = Inches(2.0)
        table.columns[2].width = Inches(4.0)

        # write column headings
        table.cell(0, 0).text = 'Identificador'
        table.cell(0, 1).text = 'Estado'
        table.cell(0, 2).text = 'Denominacion'

        # write body cells
        indice = 1
        for obra in json_map['obras']:

            if indice == 22:
                indice = 1
                slide = prs.slides.add_slide(prs.slide_layouts[5])
                shapes = slide.shapes
                shapes.title.text = 'Resultados'

                rows = 22
                cols = 3
                left = Inches(0.921)
                top = Inches(1.2)
                width = Inches(6.0)
                height = Inches(0.8)

                table = shapes.add_table(rows, cols, left, top, width, height).table
                # set column widths
                table.columns[0].width = Inches(2.0)
                table.columns[1].width = Inches(2.0)
                table.columns[2].width = Inches(4.0)

            # write column headings
            for x in range(0, 3):
                cell = table.rows[0].cells[x]
                paragraph = cell.textframe.paragraphs[0]
                paragraph.font.size = Pt(12)
                paragraph.font.name = 'Arial Black'
                paragraph.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

            for x in range(0, 3):
                cell = table.rows[indice].cells[x]
                paragraph = cell.textframe.paragraphs[0]
                paragraph.font.size = Pt(8)
                paragraph.font.name = 'Arial'
                paragraph.font.color.rgb = RGBColor(0x0B, 0x0B, 0x0B)

            table.cell(0, 0).text = 'Identificador'
            table.cell(0, 1).text = 'Estado'
            table.cell(0, 2).text = 'Denominacion'

            # write body cells
            table.cell(indice, 0).text = obra['identificador_unico']
            table.cell(indice, 1).text = obra['estado__nombreEstado']
            table.cell(indice, 2).text = obra['denominacion']
            indice += 1

        prs.save(output)
        response = StreamingHttpResponse(FileWrapper(output),
                                         content_type='application/vnd.openxmlformats-officedocument.presentationml.presentation')
        response['Content-Disposition'] = 'attachment; filename="resultado_obras.pptx"'
        response['Content-Length'] = output.tell()

        output.seek(0)

        return response


class PptxReporteEndpoint(ProtectedResourceView):
    def get(self, request):

        user = AccessToken.objects.get(token=request.GET.get('access_token')).user
        tipoReporte = request.GET.get("tipoReporte", None)
        buscador = BuscarObras(
            idtipoobra=get_array_or_none(request.GET.get('tipoDeObra')),
            iddependencias=get_array_or_none(request.GET.get('dependencia')),
            subdependencias=get_array_or_none(request.GET.get('subdependencias')),
            estados=get_array_or_none(request.GET.get('estado')),
            clasificaciones=get_array_or_none(request.GET.get('clasificacion')),
            inversiones=get_array_or_none(request.GET.get('tipoDeInversion')),
            inauguradores=get_array_or_none(request.GET.get('inaugurador')),
            impactos=get_array_or_none(request.GET.get('impacto')),
            inaugurada=request.GET.get('inaugurada', None),
            inversion_minima=request.GET.get('inversionMinima', None),
            inversion_maxima=request.GET.get('inversionMaxima', None),
            fecha_inicio_primera=request.GET.get('fechaInicio', None),
            fecha_inicio_segunda=request.GET.get('fechaInicio', None),
            fecha_fin_primera=request.GET.get('fechaFin', None),
            fecha_fin_segunda=request.GET.get('fechaFinSegunda', None),
            denominacion=request.GET.get('denominacion', None),
            instancia_ejecutora=get_array_or_none(request.GET.get('instanciaEjecutora')),
            limite_min=request.GET.get("limiteMin"),
            limite_max=request.GET.get("limiteMax"),
            busqueda_rapida=request.GET.get("busquedaRapida", None),
            id_obra=request.GET.get("idObra", None),
            susceptible_inauguracion=request.GET.get("susceptible", None),
            subclasificacion=get_array_or_none(request.GET.get('subclasificacion')),
            municipios=get_array_or_none(request.GET.get('municipios'))
        )

        buscador.filtrar_dependencias(user)
        resultados = buscador.buscar()

        json_map = {}

        json_map['reporte_dependencia'] = []
        for reporte in resultados['reporte_dependencia']:
            map = {}
            map['dependencia'] = Dependencia.objects.get(
                nombreDependencia=reporte['dependencia__nombreDependencia']).to_serializable_dict()
            map['numero_obras'] = reporte['numero_obras']
            if reporte['sumatotal'] is None:
                map['sumatotal'] = 0
            else:
                map['sumatotal'] = float(reporte['sumatotal'])
            json_map['reporte_dependencia'].append(map)

        json_map['reporte_subdependencia'] = []
        for reporteSub in resultados['reporte_subdependencia']:
            map = {}
            if reporteSub['subdependencia__nombreDependencia'] is not None:
                map['subdependencia'] = Dependencia.objects.get(
                    nombreDependencia=reporteSub['subdependencia__nombreDependencia']).to_serializable_dict()
            else:
                map['subdependencia'] = Dependencia.objects.get(
                    nombreDependencia=reporteSub['dependencia__nombreDependencia']).to_serializable_dict()
            map['numero_obras'] = reporteSub['numero_obras']
            if reporteSub['sumatotal'] is None:
                map['sumatotal'] = 0
            else:
                map['sumatotal'] = float(reporteSub['sumatotal'])
            json_map['reporte_subdependencia'].append(map)

        json_map['reporte_estado'] = []
        for reporte_estado in resultados['reporte_estado']:
            map = {}
            if reporte_estado['sumatotal'] is None:
                map['sumatotal'] = 0.0

            else:
                map['sumatotal'] = float(reporte_estado['sumatotal'])
            map['estado'] = Estado.objects.get(
                nombreEstado=reporte_estado['estado__nombreEstado']).to_serializable_dict()
            map['numeroObras'] = reporte_estado['numero_obras']

            json_map['reporte_estado'].append(map)

        json_map['reporte_general'] = []
        map = {}
        total = resultados['reporte_general']['total_invertido']['inversionTotal__sum']
        if total is None:
            total = 0.0
        else:
            total = float(total)
        map['total_invertido'] = total

        map['obras_totales'] = resultados['reporte_general']['obras_totales']
        json_map['reporte_general'].append(map)

        output = StringIO.StringIO()
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        shapes = slide.shapes
        shapes.title.text = 'Reporte'

        if tipoReporte == 'Dependencia':


            rows = 17
            cols = 3
            left = Inches(0.921)
            top = Inches(1.2)
            width = Inches(6.0)
            height = Inches(0.8)

            table = shapes.add_table(rows, cols, left, top, width, height).table
            # set column widths
            table.columns[0].width = Inches(3.0)
            table.columns[1].width = Inches(2.0)
            table.columns[2].width = Inches(2.0)

            for x in range(0, 3):
                cell = table.rows[0].cells[x]
                paragraph = cell.textframe.paragraphs[0]
                paragraph.font.size = Pt(12)
                paragraph.font.name = 'Arial Black'
                paragraph.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

            # write column headings
            table.cell(0, 0).text = 'TipoInversion'
            table.cell(0, 1).text = 'No. de Obras '
            table.cell(0, 2).text = 'Monto'

            # write body cells
            i = 1
            for obra in json_map['reporte_dependencia']:
                for x in range(0, 3):
                    cell = table.rows[i].cells[x]
                    paragraph = cell.textframe.paragraphs[0]
                    paragraph.font.size = Pt(8)
                    paragraph.font.name = 'Arial'
                    paragraph.font.color.rgb = RGBColor(0x0B, 0x0B, 0x0B)

                table.cell(i, 0).text = obra['dependencia']['nombreDependencia']
                table.cell(i, 1).text = str(obra['numero_obras'])
                table.cell(i, 2).text = str(obra['sumatotal'])
                i += 1

        if tipoReporte == 'Subdependencia':
            rows = 17
            cols = 3
            left = Inches(0.921)
            top = Inches(1.2)
            width = Inches(6.0)
            height = Inches(0.8)

            table = shapes.add_table(rows, cols, left, top, width, height).table
            # set column widths
            table.columns[0].width = Inches(3.0)
            table.columns[1].width = Inches(2.0)
            table.columns[2].width = Inches(2.0)

            for x in range(0, 3):
                cell = table.rows[0].cells[x]
                paragraph = cell.textframe.paragraphs[0]
                paragraph.font.size = Pt(12)
                paragraph.font.name = 'Arial Black'
                paragraph.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

            # write column headings
            table.cell(0, 0).text = 'TipoInversion'
            table.cell(0, 1).text = 'No. de Obras '
            table.cell(0, 2).text = 'Monto'

            # write body cells
            i = 1
            for obra in json_map['reporte_subdependencia']:
                for x in range(0, 3):
                    cell = table.rows[i].cells[x]
                    paragraph = cell.textframe.paragraphs[0]
                    paragraph.font.size = Pt(8)
                    paragraph.font.name = 'Arial'
                    paragraph.font.color.rgb = RGBColor(0x0B, 0x0B, 0x0B)
                table.cell(i, 0).text = obra['subdependencia']['nombreDependencia']
                table.cell(i, 1).text = str(obra['numero_obras'])
                table.cell(i, 2).text = str(obra['sumatotal'])
                i += 1

        if tipoReporte == 'Estado':
            rows = 35
            cols = 3
            left = Inches(0.921)
            top = Inches(1.2)
            width = Inches(6.0)
            height = Inches(0.8)

            table = shapes.add_table(rows, cols, left, top, width, height).table
            # set column widths
            table.columns[0].width = Inches(3.0)
            table.columns[1].width = Inches(2.0)
            table.columns[2].width = Inches(2.0)

            for x in range(0, 3):
                cell = table.rows[0].cells[x]
                paragraph = cell.textframe.paragraphs[0]
                paragraph.font.size = Pt(12)
                paragraph.font.name = 'Arial Black'
                paragraph.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

            # write column headings
            table.cell(0, 0).text = 'TipoInversion'
            table.cell(0, 1).text = 'No. de Obras '
            table.cell(0, 2).text = 'Monto'

            # write body cells
            i = 1
            for obra in json_map['reporte_estado']:
                for x in range(0, 3):
                    cell = table.rows[i].cells[x]
                    paragraph = cell.textframe.paragraphs[0]
                    paragraph.font.size = Pt(8)
                    paragraph.font.name = 'Arial'
                    paragraph.font.color.rgb = RGBColor(0x0B, 0x0B, 0x0B)
                table.cell(i, 0).text = obra['estado']['nombreEstado']
                table.cell(i, 1).text = str(obra['numeroObras'])
                table.cell(i, 2).text = str(obra['sumatotal'])
                i += 1

        prs.save(output)
        response = StreamingHttpResponse(FileWrapper(output),
                                         content_type='application/vnd.openxmlformats-officedocument.presentationml.presentation')
        response['Content-Disposition'] = 'attachment; filename="resultado_obras.pptx"'
        response['Content-Length'] = output.tell()

        output.seek(0)

        return response


class ListarEndpoint(ProtectedResourceView):
    def get(self, request):

        user = AccessToken.objects.get(token=request.GET.get('access_token')).user

        idtipoobra = get_array_or_none(request.GET.get('tipoDeObra')),
        iddependencias = get_array_or_none(request.GET.get('dependencia')),
        idsubdependencias = get_array_or_none(request.GET.get('subdependencias')),
        idmunicipios = get_array_or_none(request.GET.get('muicipios')),
        estados = get_array_or_none(request.GET.get('estado')),
        clasificaciones = get_array_or_none(request.GET.get('clasificacion')),
        inversiones = get_array_or_none(request.GET.get('tipoDeInversion')),
        inauguradores = get_array_or_none(request.GET.get('inaugurador')),
        impactos = get_array_or_none(request.GET.get('impacto')),
        inaugurada = request.GET.get('inaugurada', None),
        inversion_minima = request.GET.get('inversionMinima', None),
        inversion_maxima = request.GET.get('inversionMaxima', None),
        fecha_inicio_primera = request.GET.get('fechaInicio', None),
        fecha_inicio_segunda = request.GET.get('fechaInicio', None),
        fecha_fin_primera = request.GET.get('fechaFin', None),
        fecha_fin_segunda = request.GET.get('fechaFinSegunda', None),
        denominacion = request.GET.get('denominacion', None),
        instancia_ejecutora = get_array_or_none(request.GET.get('instanciaEjecutora')),
        busqueda_rapida = request.GET.get("busquedaRapida", None),
        id_obra = request.GET.get("idObra", None),
        susceptible_inauguracion = request.GET.get("susceptible", None),
        subclasificacion = get_array_or_none(request.GET.get('subclasificacion')),

        p_tipoobra = ","
        p_estados = ","
        p_clasificaciones = ","
        p_inversiones = ","
        p_inauguradores = ","
        p_impactos = ","
        p_instancia_ejecutora = ","
        p_subclasificacion = ","

        arreglo_dependencias = []
        arreglo_subdependencias = []
        p_dependencias = ""
        p_subdependencias = ","
        p_municipios = ","
        bandera_AD=0
        bandera_US=0

        if user.usuario.rol == 'SA' and get_array_or_none(request.GET.get('dependencia')) is None:
            p_dependencias = ","

        elif (user.usuario.rol == 'AD' or user.usuario.rol == 'FU') and get_array_or_none(request.GET.get('dependencia')) is  None:

            for dependencia in user.usuario.dependencia.all():
                arreglo_dependencias.append(dependencia.id)

            for subdependencia in user.usuario.subdependencia.all():
                arreglo_dependencias.append(subdependencia.id)

            iddependencias = arreglo_dependencias
            bandera_AD = 1

        elif user.usuario.rol == 'US' and get_array_or_none(request.GET.get('subdependencias')) is None:
            for subdependencia in user.usuario.subdependencia.all():
                arreglo_subdependencias.append(subdependencia.id)

            idsubdependencias = arreglo_subdependencias
            bandera_US = 1





        if iddependencias[0] is not None:
            if bandera_AD == 1:
                for dependencia in iddependencias:
                    p_dependencias += str(dependencia) + ","
            else:
                for dependencia in iddependencias[0]:
                    p_dependencias += str(dependencia) + ","

        if idsubdependencias[0] is not None:
            if bandera_US == 1:
                for subdependencia in idsubdependencias:
                    p_subdependencias += str(subdependencia) + ","
            else:
                for subdependencia in idsubdependencias[0]:
                    p_subdependencias += str(subdependencia) + ","

        if idmunicipios[0] is not None:
            for municipio in idmunicipios[0]:
                p_municipios += str(municipio) + ","

        if idtipoobra[0] is not None:
            p_tipoobra = ""
            for tipoobra in idtipoobra[0]:
                p_tipoobra += str(tipoobra) + ","

        if estados[0] is not None:
            p_estados = ""
            for estado in estados[0]:
                p_estados += str(estado) + ","

        if clasificaciones[0] is not None:
            p_clasificaciones = ""
            for clasificacion in clasificaciones[0]:
                p_clasificaciones += str(clasificacion) + ","

        if inversiones[0] is not None:
            p_inversiones = ""
            for inversion in inversiones[0]:
                p_inversiones += str(inversion) + ","

        if inauguradores[0] is not None:
            p_inauguradores = ""
            for inaugurador in inauguradores[0]:
                p_inauguradores += str(inaugurador) + ","

        if impactos[0] is not None:
            p_impactos = ""
            for impacto in impactos[0]:
                p_impactos += str(impacto) + ","

        if instancia_ejecutora[0] is not None:
            p_instancia_ejecutora = ""
            for instancia in instancia_ejecutora[0]:
                p_instancia_ejecutora += str(instancia) + ","

        p_inaugurada = ""
        p_inversion_minima = 0
        p_inversion_maxima = 999999999.99
        p_fecha_inicio_primera = ""
        p_fecha_inicio_segunda = ""
        p_fecha_fin_primera = ""
        p_fecha_fin_segunda = ""
        p_denominacion = ""

        if inaugurada[0] is None: p_inaugurada = ""
        if inversion_minima[0] is None:
            p_inversion_minima = 0
        else:
            p_inversion_minima = inversion_minima[0]
        if inversion_maxima[0] is None:
            p_inversion_maxima = 999999999.99
        else:
            p_inversion_maxima = inversion_maxima[0]
        if fecha_inicio_primera[0] is None:
            p_fecha_inicio_primera = "1900-01-01"
        else:
            p_fecha_inicio_primera = fecha_inicio_primera[0]
        if fecha_inicio_segunda[0] is None:
            p_fecha_inicio_segunda = "2100-12-31"
        else:
            p_fecha_inicio_segunda = fecha_inicio_segunda[0]
        if fecha_fin_primera[0] is None:
            p_fecha_fin_primera = "1900-01-01"
        else:
            p_fecha_fin_primera = fecha_fin_primera[0]
        if fecha_fin_segunda[0] is None:
            p_fecha_fin_segunda = "2100-12-31"
        else:
            p_fecha_fin_segunda = fecha_fin_segunda[0]
        if denominacion[0] is None:
            p_denominacion = ""
        else:
            p_denominacion = denominacion[0]

        results = Obra.searchList(p_tipoobra[:-1], p_dependencias[:-1], p_subdependencias[:-1], p_municipios[:-1], p_instancia_ejecutora[:-1], p_estados[:-1],
                                  p_inversion_minima, p_inversion_maxima, p_fecha_inicio_primera,
                                  p_fecha_inicio_segunda, p_fecha_fin_primera, p_fecha_fin_segunda, p_impactos[:-1],
                                  p_inauguradores[:-1], p_inversiones[:-1], p_clasificaciones[:-1],
                                  "", "", p_denominacion)

        output = StringIO.StringIO()
        book = Workbook(output)
        sheet = book.add_worksheet('obras')

        if results:
            # Add a bold format to use to highlight cells.
            bold = book.add_format({'bold': True})
            # encabezados
            sheet.write(0, 0, "Tipo de Obra", bold)
            sheet.write(0, 1, "id Unico", bold)
            sheet.write(0, 2, "Dependencia/Organismo", bold)
            sheet.write(0, 3, "Sub Dependencia", bold)
            sheet.write(0, 4, "Estado", bold)
            sheet.write(0, 5, "Denominacion", bold)
            sheet.write(0, 6, "Descripcion", bold)
            sheet.write(0, 7, "Municipio", bold)
            sheet.write(0, 8, "Fecha Inicio", bold)
            sheet.write(0, 9, "Fecha Termino", bold)
            sheet.write(0, 10, "Avance Fisico %", bold)
            sheet.write(0, 11, "F", bold)
            sheet.write(0, 12, "E", bold)
            sheet.write(0, 13, "M", bold)
            sheet.write(0, 14, "S", bold)
            sheet.write(0, 15, "P", bold)
            sheet.write(0, 16, "O", bold)
            sheet.write(0, 17, "Inversion Total", bold)
            sheet.write(0, 18, "Tipo Moneda MDP/MDD", bold)
            sheet.write(0, 19, "Registro de Cartera", bold)        #**************
            sheet.write(0, 20, "Poblacion Objetivo", bold)
            sheet.write(0, 21, "Numero de Beneficiarios", bold)
            sheet.write(0, 22, "Impacto", bold)
            sheet.write(0, 23, "CG", bold)
            sheet.write(0, 24, "PNG", bold)
            sheet.write(0, 25, "PM", bold)
            sheet.write(0, 26, "PNI", bold)
            sheet.write(0, 27, "CNCH", bold)
            sheet.write(0, 28, "OI", bold)
            sheet.write(0, 29, "Senalizacion", bold)
            sheet.write(0, 30, "Observaciones", bold)
            sheet.write(0, 31, "Inaugurado por:", bold)
            sheet.write(0, 32, "Inaugurado:", bold)
            sheet.write(0, 33, "Foto Antes", bold)
            sheet.write(0, 34, "Foto Durante", bold)
            sheet.write(0, 35, "Foto Despues", bold)
            sheet.write(0, 36, "Latitud", bold)
            sheet.write(0, 37, "Longitud", bold)

            i = 1
            for obra in results:
                id_unico = obra[0]
                sheet.write(i, 0, obra[0])
                sheet.write(i, 1, obra[1])
                sheet.write(i, 2, obra[2])
                sheet.write(i, 3, obra[30])
                sheet.write(i, 4, obra[3])
                sheet.write(i, 5, obra[4])
                sheet.write(i, 6, obra[5])
                sheet.write(i, 7, obra[6])
                sheet.write(i, 8, obra[7])
                sheet.write(i, 9, obra[8])
                sheet.write(i, 10, obra[9])

                sheet.write(i, 11, int("0"))  #F
                sheet.write(i, 12, int("0"))  #E
                sheet.write(i, 13, int("0"))  #M
                sheet.write(i, 14, int("0"))  #S
                sheet.write(i, 15, int("0"))  #P
                sheet.write(i, 16, int("0"))  #O

                lista=[0,0,0,0,0,0]
                inversionA=[]
                inversionB=[]

                try:
                    for inv in (obra[11].split(',')):
                        inversionA.append(inv[0])
                except Exception as e:
                    print e

                try:
                    for monto in (obra[12].split(',')):
                        inversionB.append(monto)
                except Exception as e:
                    print e

                if i==37:
                    print i

                if len(inversionB) == 6:
                    for x in inversionA:
                        lista[int(x)-1]=inversionB[int(x)-1]

                    for x in range (0, len(lista)):
                         sheet.write(i, x+11, float(lista[x]))
                else:
                    z=0
                    for x in inversionA:
                        sheet.write(i, int(x)+10, float(inversionB[z]))
                        z+=1

                #for inv in (obra[10].split(',')):
                #    if inv[0] == "F": sheet.write(i, 10, "SI")
                #    if inv[0] == "E": sheet.write(i, 11, "SI")
                #    if inv[0] == "M": sheet.write(i, 12, "SI")
                #    if inv[0] == "S": sheet.write(i, 13, "SI")
                #    if inv[0] == "P": sheet.write(i, 14, "SI")
                #    if inv[0] == "O": sheet.write(i, 15, "SI")

                sheet.write(i, 17, obra[13])
                sheet.write(i, 18, obra[14])
                sheet.write(i, 19, obra[29])
                sheet.write(i, 20, obra[15])
                sheet.write(i, 21, obra[16])
                sheet.write(i, 22, obra[17])

                sheet.write(i, 23, "NO")  #CG
                sheet.write(i, 24, "NO")  #PNG
                sheet.write(i, 25, "NO")  #PM
                sheet.write(i, 26, "NO")  #PNI
                sheet.write(i, 27, "NO")  #CNCH
                sheet.write(i, 28, "NO")  #OI
                if obra[18] is not None:
                    for cla in (obra[18].split(',')):
                        sCla = ''.join(cla)
                        if sCla == 'CG':
                            if obra[19] is not None:
                                for subscla in (obra[19].split(',')):
                                    sSubCla = ''.join(subscla)
                                    if sSubCla[:2] == 'CG': sheet.write(i, 23, sSubCla)
                        if sCla == "PNG": sheet.write(i, 24, "SI")
                        if sCla == "PM": sheet.write(i, 25, "SI")
                        if sCla == "PNI":
                            sheet.write(i, 26, "SI")
                            if obra[19] is not None:
                                for subscla in (obra[19].split(',')):
                                    sSubCla = ''.join(subscla)
                                    if sSubCla[:3] == "PNI": sheet.write(i, 26, sSubCla)
                        if sCla == "CNCH": sheet.write(i, 27, "SI")
                        if sCla == "OI": sheet.write(i, 28, "SI")

                sheet.write(i, 29, obra[20])
                sheet.write(i, 30, obra[21])
                sheet.write(i, 31, obra[22])

                sheet.write(i, 32, "NO")
                if obra[31] == 1:
                    sheet.write(i, 32, "SI")  #inaugurado
                sheet.write(i, 33, obra[24])
                sheet.write(i, 34, obra[25])
                sheet.write(i, 35, obra[26])
                sheet.write(i, 36, obra[27])
                sheet.write(i, 37, obra[28])

                #documentos fuente
                iDoc = 38
                try:
                    for doc in (obra[32].split(',')):
                        sheet.write(0, iDoc, "Documento-"+str(iDoc-37), bold)
                        sheet.write(i, iDoc, doc)
                        iDoc+=1
                except Exception as e:
                    print e

                i += 1
            book.close()

            # construct response

            # output.seek(0)
            #response = HttpResponse(output.read(), content_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet")
            #response['Content-Disposition'] = "attachment; filename=test.xlsx"
        else:
            sheet.write(0, 0,
                        "Los filtros seleccionados no arrojaron informacion alguna sobre las obras, cambie los filtros para una nueva consulta.")
            book.close()

        response = StreamingHttpResponse(FileWrapper(output),
                                         content_type='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet')
        response['Content-Disposition'] = 'attachment; filename="listado_obras.xlsx"'
        response['Content-Length'] = output.tell()

        output.seek(0)

        return response



        # return HttpResponse(json.dumps(json_map), 'application/json')


class InauguradorEndpoint(ProtectedResourceView):
    def get(self, request):
        return HttpResponse(
            json.dumps(map(lambda inaugurador: inaugurador.to_serializable_dict(), Inaugurador.objects.all())),
            'application/json')


class NumeroObrasPendientes(ProtectedResourceView):
    def get(self, request):
        token = request.GET.get('access_token')
        token_model = AccessToken.objects.get(token=token)
        arreglo_dependencias = []
        total_obras = 0

        for dependencia in token_model.user.usuario.dependencia.all():
            arreglo_dependencias.append(dependencia.id)

        if token_model.user.usuario.rol == 'AD' or token_model.user.usuario.rol == 'FU' or token_model.user.usuario.rol == 'UD':
            dependencias = Dependencia.objects.filter(
                Q(id__in=arreglo_dependencias) |
                Q(dependienteDe__id__in=arreglo_dependencias)
            )

            for dependencia in dependencias:
                total_obras = total_obras + dependencia.obra_set.filter(autorizada=False).count()

        ans = {}

        ans['obrasTotalesPendientes'] = total_obras

        return HttpResponse(json.dumps(ans), 'application/json')


class ReporteInicioEndpoint(ProtectedResourceView):
    def rename_estado(self, obra):
        obra['estado'] = obra['estado__nombreEstado']
        del obra['estado__nombreEstado']

    def get(self, request):
        dependencias = get_usuario_for_token(request.GET.get('access_token')).dependencia.all()
        subdependencias = get_usuario_for_token(request.GET.get('access_token')).subdependencia.all()

        if dependencias and dependencias.count() > 0:
            obras = Obra.objects.filter(dependencia__in=dependencias)
        else:
            obras = Obra.objects.all()

        if subdependencias and subdependencias.count() > 0:
            obras = obras.filter(subdependencia__in=subdependencias)

        reporte = {
            'reporte_mapa': {'obras_mapa': {}},
            'reporte_total': {'obras_proceso': {}, 'obras_proyectadas': {}, 'obras_concluidas': {}},
            'reporte2015': {'obras_proceso': {}, 'obras_proyectadas': {}, 'obras_concluidas': {}},
            'reporte2014': {'obras_concluidas': {}},
            'reporte2013': {'obras_concluidas': {}},
            'reporte2012': {'obras_concluidas': {}},
        }

        obras_mapa = obras.exclude(tipoObra_id=4)
        the_list = []
        reporte_estado = obras.values('latitud', 'longitud', 'estado__nombreEstado').annotate(
            numero_obras=Count('estado')).annotate(
            totalinvertido=Sum('inversionTotal'))
        for obra in reporte_estado:
            self.rename_estado(obra)
            the_list.append(obra)
        reporte['reporte_mapa']['obras_mapa']['obras'] = the_list
        reporte['reporte_mapa']['obras_mapa']['total'] = obras_mapa.count()
        reporte['reporte_mapa']['obras_mapa']['inversion_total'] = obras_mapa.aggregate(Sum('inversionTotal'))[
            'inversionTotal__sum']

        # Grafico, obras totales
        obras_totales_proceso = obras.filter(tipoObra_id=2)
        the_list = []
        for obra in obras_totales_proceso.values('latitud', 'longitud', 'estado__nombreEstado'):
            self.rename_estado(obra)
            the_list.append(obra)
        reporte['reporte_total']['obras_proceso']['obras'] = the_list
        reporte['reporte_total']['obras_proceso']['total'] = obras_totales_proceso.count()
        reporte['reporte_total']['obras_proceso']['inversion_total'] = \
            obras_totales_proceso.aggregate(Sum('inversionTotal'))['inversionTotal__sum']

        obras_totales_proyectadas = obras.filter(tipoObra_id=1)
        the_list = []
        for obra in obras_totales_proyectadas.values('latitud', 'longitud', 'estado__nombreEstado'):
            self.rename_estado(obra)
            the_list.append(obra)
        reporte['reporte_total']['obras_proyectadas']['obras'] = the_list
        reporte['reporte_total']['obras_proyectadas']['total'] = obras_totales_proyectadas.count()
        reporte['reporte_total']['obras_proyectadas']['inversion_total'] = \
            obras_totales_proyectadas.aggregate(Sum('inversionTotal'))['inversionTotal__sum']

        obras_totales_concluidas = obras.filter(tipoObra_id=3)
        the_list = []
        for obra in obras_totales_concluidas.values('latitud', 'longitud', 'estado__nombreEstado'):
            self.rename_estado(obra)
            the_list.append(obra)
        reporte['reporte_total']['obras_concluidas']['obras'] = the_list
        reporte['reporte_total']['obras_concluidas']['total'] = obras_totales_concluidas.count()
        reporte['reporte_total']['obras_concluidas']['inversion_total'] = \
            obras_totales_concluidas.aggregate(Sum('inversionTotal'))['inversionTotal__sum']

        # Reportes anuales 2012-2015
        obras2015 = obras.filter(fechaInicio__year=2015)
        obras2015_proceso = obras2015.filter(tipoObra_id=2)
        the_list = []

        for obra in obras2015_proceso.values('latitud', 'longitud', 'estado__nombreEstado').annotate(
                numero_obras=Count('estado')):
            self.rename_estado(obra)
            the_list.append(obra)
        reporte['reporte2015']['obras_proceso']['obras'] = the_list
        reporte['reporte2015']['obras_proceso']['total'] = obras2015_proceso.count()

        obras2015_proyectadas = obras2015.filter(tipoObra_id=1)
        the_list = []
        for obra in obras2015_proyectadas.values('latitud', 'longitud', 'estado__nombreEstado').annotate(
                numero_obras=Count('estado')):
            self.rename_estado(obra)
            the_list.append(obra)
        reporte['reporte2015']['obras_proyectadas']['obras'] = the_list
        reporte['reporte2015']['obras_proyectadas']['total'] = obras2015_proyectadas.count()

        obras2015_concluidas = obras.filter(Q(fechaTermino__year=2015) & Q(tipoObra_id=3))
        the_list = []
        for obra in obras2015_concluidas.values('latitud', 'longitud', 'estado__nombreEstado').annotate(
                numero_obras=Count('estado')):
            self.rename_estado(obra)
            the_list.append(obra)
        reporte['reporte2015']['obras_concluidas']['obras'] = the_list
        reporte['reporte2015']['obras_concluidas']['total'] = obras2015_concluidas.count()

        # Obras Concluidas, barra inferior

        obras2014 = obras.filter(Q(fechaTermino__year=2014) & Q(tipoObra_id=3))
        the_list = []
        for obra in obras2014.values('latitud', 'longitud', 'estado__nombreEstado'):
            self.rename_estado(obra)
            the_list.append(obra)
        reporte['reporte2014']['obras_concluidas']['total'] = obras2014.count()

        obras2013 = obras.filter(Q(fechaTermino__year=2013) & Q(tipoObra_id=3))
        the_list = []
        for obra in obras2013.values('latitud', 'longitud', 'estado__nombreEstado'):
            self.rename_estado(obra)
            the_list.append(obra)
        reporte['reporte2013']['obras_concluidas']['obras'] = the_list
        reporte['reporte2013']['obras_concluidas']['total'] = obras2013.count()

        obras2012 = obras.filter(Q(fechaTermino__year=2012) & Q(tipoObra_id=3))
        the_list = []
        for obra in obras2012.values('latitud', 'longitud', 'estado__nombreEstado'):
            self.rename_estado(obra)
            the_list.append(obra)
        reporte['reporte2012']['obras_concluidas']['obra'] = the_list
        reporte['reporte2012']['obras_concluidas']['total'] = obras2012.count()

        return HttpResponse(json.dumps(reporte), 'application/json')


class ReporteNoTrabajoEndpoint(ProtectedResourceView):
    def get(self, request):
        comp_date = datetime.now() - timedelta(15)
        print "****"
        print comp_date
        dicts = map(lambda dependencia: dependencia.to_serializable_dict(), Dependencia.objects.filter(
            Q(obraoprograma='O') &
            Q(fecha_ultima_modificacion__lt=comp_date)))

        return HttpResponse(json.dumps(dicts), 'application/json')


class ReporteObrasPorAutorizar(ProtectedResourceView):
    def get(self, request):
        dependencias = get_usuario_for_token(request.GET.get('access_token')).dependencia.all()
        subdependencias = get_usuario_for_token(request.GET.get('access_token')).subdependencia.all()

        if dependencias and dependencias.count() > 0:
            if get_usuario_for_token(request.GET.get('access_token')).rol == 'US':
                obras = Obra.objects.filter(
                    Q(subdependencia__in=get_subdependencias_as_list_flat(subdependencias))
                )
            else:
                obras = Obra.objects.filter(
                    Q(dependencia__in=get_subdependencias_as_list_flat(dependencias)) |
                    Q(subdependencia__in=get_subdependencias_as_list_flat(dependencias))
                )
        else:
            obras = Obra.objects.all()

        obras = obras.filter(autorizada=False).values('id', 'identificador_unico', 'estado__nombreEstado',
                                                      'denominacion', 'latitud', 'longitud',
                                                      'dependencia__imagenDependencia')

        the_list = []
        for obra in obras:
            the_list.append(obra)

        return HttpResponse(json.dumps(the_list), 'application/json')


class ReporteImagenesEndpoint(ProtectedResourceView):
    def rename_estado(self, obra):
        obra['estado'] = obra['estado__nombreEstado']
        del obra['estado__nombreEstado']

    def get(self, request):
        obras = Obra.objects.all().exclude(tipoObra_id=4)

        reporte = {
            'reporte_total': {'obras_proceso': {}, 'obras_proyectadas': {}, 'obras_concluidas': {}},
        }

        prs = Presentation('/home/obrasapf/djangoObras/obras/static/ppt/PRESENTACION_OBRAS_APF.pptx')


        # Datos diapositiva 1
        obras_totales_proceso = obras.filter(tipoObra_id=2)
        obras_total_proceso = obras_totales_proceso.count()
        inversion_total_proceso = obras_totales_proceso.aggregate(Sum('inversionTotal'))['inversionTotal__sum']

        obras_totales_proyectadas = obras.filter(tipoObra_id=1)
        obras_total_proyectadas = obras_totales_proyectadas.count()
        inversion_total_proyectadas = obras_totales_proyectadas.aggregate(Sum('inversionTotal'))['inversionTotal__sum']

        obras_totales_concluidas = obras.filter(tipoObra_id=3)
        obras_total_concluidas = obras_totales_concluidas.count()
        inversion_total_concluidas = obras_totales_concluidas.aggregate(Sum('inversionTotal'))['inversionTotal__sum']

        obras_Total = obras_total_proceso + obras_total_proyectadas + obras_total_concluidas
        inversion_total = inversion_total_proceso + inversion_total_proyectadas + inversion_total_concluidas

        prs.slides[0].shapes[13].text_frame.paragraphs[0].font.size = Pt(16)
        prs.slides[0].shapes[13].text_frame.paragraphs[0].font.name = 'Arial Black'
        prs.slides[0].shapes[13].text_frame.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        prs.slides[0].shapes[13].text = '{0:,}'.format(obras_total_concluidas)

        prs.slides[0].shapes[14].text_frame.paragraphs[0].font.size = Pt(16)
        prs.slides[0].shapes[14].text_frame.paragraphs[0].font.name = 'Arial Black'
        prs.slides[0].shapes[14].text_frame.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        prs.slides[0].shapes[14].text = '{0:,}'.format(obras_total_proceso)

        prs.slides[0].shapes[15].text_frame.paragraphs[0].font.size = Pt(16)
        prs.slides[0].shapes[15].text_frame.paragraphs[0].font.name = 'Arial Black'
        prs.slides[0].shapes[15].text_frame.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        prs.slides[0].shapes[15].text = '{0:,}'.format(obras_total_proyectadas)

        prs.slides[0].shapes[11].text_frame.paragraphs[0].font.size = Pt(20)
        prs.slides[0].shapes[11].text_frame.paragraphs[0].font.name = 'Arial Black'
        prs.slides[0].shapes[11].text_frame.paragraphs[0].font.color.rgb = RGBColor(0x40, 0x40, 0x40)
        prs.slides[0].shapes[11].text = '{0:,}'.format(obras_Total)
        prs.slides[0].shapes[12].text_frame.paragraphs[0].font.size = Pt(7)
        prs.slides[0].shapes[12].text_frame.paragraphs[0].font.name = 'Arial Black'
        prs.slides[0].shapes[12].text_frame.paragraphs[0].font.color.rgb = RGBColor(0x40, 0x40, 0x40)
        prs.slides[0].shapes[12].text = '{0:,.2f}'.format(inversion_total)

        # datos para la diapositiva 2

        # Reporte Dependencia
        reporte_dependencia = obras.values('dependencia__id', 'dependencia__nombreDependencia',
                                           'dependencia__orden_secretaria').annotate(
            numero_obras=Count('dependencia')).annotate(sumatotal=Sum('inversionTotal'))
        reporte_dependencia = reporte_dependencia.order_by('dependencia__orden_secretaria')

        json_map_dep = {}
        json_map_dep['reporte_dependencia'] = []
        for obra in reporte_dependencia:
            map = {}
            map['dependencia'] = obra['dependencia__nombreDependencia']
            map['numero_obras'] = obra['numero_obras']
            map['id'] = obra['dependencia__id']
            json_map_dep['reporte_dependencia'].append(map)

        for x in range(20, 35):
            prs.slides[1].shapes[x].text_frame.paragraphs[0].font.size = Pt(10)
            prs.slides[1].shapes[x].text_frame.paragraphs[0].font.name = 'Arial Black'

        for x in range(20, 24):
            prs.slides[1].shapes[x].text_frame.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

        for x in range(24, 26):
            prs.slides[1].shapes[x].text_frame.paragraphs[0].font.color.rgb = RGBColor(0x00, 0xCC, 0x00)
        prs.slides[1].shapes[26].text_frame.paragraphs[0].font.color.rgb = RGBColor(0xCC, 0xCC, 0x00)
        prs.slides[1].shapes[27].text_frame.paragraphs[0].font.color.rgb = RGBColor(0x00, 0xCC, 0x00)
        prs.slides[1].shapes[28].text_frame.paragraphs[0].font.color.rgb = RGBColor(0xCC, 0xCC, 0x00)
        for x in range(29, 31):
            prs.slides[1].shapes[x].text_frame.paragraphs[0].font.color.rgb = RGBColor(0x00, 0xCC, 0x00)
        prs.slides[1].shapes[31].text_frame.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        prs.slides[1].shapes[32].text_frame.paragraphs[0].font.color.rgb = RGBColor(0xCC, 0xCC, 0x00)
        prs.slides[1].shapes[33].text_frame.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        prs.slides[1].shapes[34].text_frame.paragraphs[0].font.color.rgb = RGBColor(0x00, 0xCC, 0x00)

        i = 20
        for obra in json_map_dep['reporte_dependencia']:
            prs.slides[1].shapes[i].text = obra['dependencia'] + '  ' + str('{0:,}'.format(obra['numero_obras']))
            i += 1

        # Datos para la diapositiva 3
        #Reporte Estado
        reporte_estado = obras.exclude(estado_id=33).exclude(estado_id=34).values('estado__nombreEstado',
                                                                                  'estado__id').annotate(
            numero_obras=Count('estado')).annotate(sumatotal=Sum('inversionTotal'))

        reporte_estado = reporte_estado.order_by('estado__id')
        json_map_edo = {}
        json_map_edo['reporte_estado'] = []
        for obra in reporte_estado:
            map = {}
            map['id'] = obra['estado__id']
            map['estado'] = obra['estado__nombreEstado']
            map['numero_obras'] = obra['numero_obras']
            map['total_invertido'] = obra['sumatotal']
            json_map_edo['reporte_estado'].append(map)

        for x in range(7, 39):
            prs.slides[2].shapes[x].text_frame.paragraphs[0].font.size = Pt(9)
            prs.slides[2].shapes[x].text_frame.paragraphs[0].font.name = 'Arial Black'
            prs.slides[2].shapes[x].text_frame.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

        i = 7
        for obra in json_map_edo['reporte_estado']:
            prs.slides[2].shapes[i].text = str('{0:,}'.format(obra['numero_obras']))
            i += 1


        # Datos para las diapositivas 4 a 18
        islide = 3
        for obra in json_map_dep['reporte_dependencia']:
            dependencia_PPC = obras.filter(Q(dependencia_id=obra['id'])).values('dependencia__nombreDependencia',
                                                                                'tipoObra_id').annotate(
                numero_obras=Count('dependencia')).annotate(sumatotal=Sum('inversionTotal'))
            dependencia_PPC = dependencia_PPC.order_by('tipoObra_id').reverse()

            json_map_depPPC = {}
            json_map_depPPC['reporte_dependencia_PPC'] = []
            for obra in dependencia_PPC:
                map = {}
                map['dependencia'] = obra['dependencia__nombreDependencia']
                map['tipo_obra'] = obra['tipoObra_id']
                map['numero_obras'] = obra['numero_obras']
                map['total_invertido'] = obra['sumatotal']
                json_map_depPPC['reporte_dependencia_PPC'].append(map)

            total_obras = 0
            total_invertido = 0
            for obra in json_map_depPPC['reporte_dependencia_PPC']:

                for x in range(14, 20):
                    prs.slides[islide].shapes[x].text_frame.paragraphs[0].font.name = 'Arial Narrow'
                    prs.slides[islide].shapes[x].text_frame.paragraphs[0].font.bold = True
                    prs.slides[islide].shapes[x].text_frame.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                if obra['tipo_obra'] == 3:
                    prs.slides[islide].shapes[14].text_frame.paragraphs[0].font.size = Pt(24)
                    prs.slides[islide].shapes[14].text = str('{0:,}'.format(obra['numero_obras']))
                    prs.slides[islide].shapes[17].text_frame.paragraphs[0].font.size = Pt(10)
                    prs.slides[islide].shapes[17].text = str('{0:,.2f}'.format(obra['total_invertido'])) + ' MDP'
                    total_obras += obra['numero_obras']
                    total_invertido += obra['total_invertido']

                if obra['tipo_obra'] == 2:
                    prs.slides[islide].shapes[15].text_frame.paragraphs[0].font.size = Pt(24)
                    prs.slides[islide].shapes[15].text = str('{0:,}'.format(obra['numero_obras']))
                    prs.slides[islide].shapes[18].text_frame.paragraphs[0].font.size = Pt(10)
                    prs.slides[islide].shapes[18].text = str('{0:,.2f}'.format(obra['total_invertido'])) + ' MDP'
                    total_obras += obra['numero_obras']
                    total_invertido += obra['total_invertido']

                if obra['tipo_obra'] == 1:
                    prs.slides[islide].shapes[16].text_frame.paragraphs[0].font.size = Pt(24)
                    prs.slides[islide].shapes[16].text = str('{0:,}'.format(obra['numero_obras']))
                    prs.slides[islide].shapes[19].text_frame.paragraphs[0].font.size = Pt(10)
                    prs.slides[islide].shapes[19].text = str('{0:,.2f}'.format(obra['total_invertido'])) + ' MDP'
                    total_obras += obra['numero_obras']
                    total_invertido += obra['total_invertido']

            prs.slides[islide].shapes[12].text_frame.paragraphs[0].font.size = Pt(20)
            prs.slides[islide].shapes[12].text_frame.paragraphs[0].font.name = 'Arial Narrow'
            prs.slides[islide].shapes[12].text_frame.paragraphs[0].font.bold = True
            prs.slides[islide].shapes[12].text_frame.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0x80, 0x00)
            prs.slides[islide].shapes[13].text_frame.paragraphs[0].font.size = Pt(10)
            prs.slides[islide].shapes[13].text_frame.paragraphs[0].font.name = 'Arial Narrow'
            prs.slides[islide].shapes[13].text_frame.paragraphs[0].font.bold = True
            prs.slides[islide].shapes[13].text_frame.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0x80, 0x00)

            prs.slides[islide].shapes[12].text = str('{0:,}'.format(total_obras)) + ' obras'
            prs.slides[islide].shapes[13].text = str('{0:,.2f}'.format(total_invertido)) + ' MDP'

            islide += 1
            total_obras = 0
            total_invertido = 0

        # Datos para las diapositivas de la 19 a la 50
        islide = 18
        for estado in json_map_edo['reporte_estado']:
            obras = Obra.objects.filter(Q(estado_id=estado['id'])) \
                .values('dependencia__nombreDependencia', 'dependencia__id').annotate(
                numero_obras=Count('dependencia')).annotate(sumatotal=Sum('inversionTotal'))
            obrasMaximas = obras.order_by('numero_obras').reverse()[:2]

            prs.slides[islide].shapes[11].text_frame.paragraphs[0].font.size = Pt(24)
            prs.slides[islide].shapes[11].text_frame.paragraphs[0].font.name = 'Arial Narrow'
            prs.slides[islide].shapes[11].text_frame.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0x80, 0x00)
            prs.slides[islide].shapes[12].text_frame.paragraphs[0].font.size = Pt(10)
            prs.slides[islide].shapes[12].text_frame.paragraphs[0].font.name = 'Arial Narrow'
            prs.slides[islide].shapes[12].text_frame.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0x80, 0x00)

            for x in range(13, 26):
                prs.slides[islide].shapes[x].text_frame.paragraphs[0].font.name = 'Arial Narrow'
                prs.slides[islide].shapes[x].text_frame.paragraphs[0].font.bold = True
                prs.slides[islide].shapes[x].text_frame.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

            prs.slides[islide].shapes[11].text_frame.paragraphs[0].font.size = Pt(36)
            prs.slides[islide].shapes[11].text = str('{0:,}'.format(estado['numero_obras'])) + ' obras'
            prs.slides[islide].shapes[12].text_frame.paragraphs[0].font.size = Pt(12)
            prs.slides[islide].shapes[12].text = str('{0:,.2f}'.format(estado['total_invertido'])) + ' MDP'

            i = 13
            j = 15
            k = 19
            for obra in obrasMaximas:
                obraMax = Obra.objects.filter(Q(estado_id=estado['id']) & Q(dependencia_id=obra['dependencia__id'])) \
                    .values('tipoObra_id').annotate(numero_obras=Count('dependencia')).annotate(
                    sumatotal=Sum('inversionTotal'))
                obraMax = obraMax.order_by('tipoObra_id').reverse()
                prs.slides[islide].shapes[i].text_frame.paragraphs[0].font.size = Pt(14)
                prs.slides[islide].shapes[i].text = obra['dependencia__nombreDependencia']

                prs.slides[islide].shapes[j].text_frame.paragraphs[0].font.size = Pt(28)
                prs.slides[islide].shapes[j].text = '{0:,}'.format(obra['numero_obras'])
                j += 1
                prs.slides[islide].shapes[j].text_frame.paragraphs[0].font.size = Pt(10)
                prs.slides[islide].shapes[j].text = str('{0:,.2f}'.format(obra['sumatotal'])) + ' MDP'
                j += 1
                i += 1

                prs.slides[islide].shapes[k].text_frame.paragraphs[0].font.size = Pt(18)
                prs.slides[islide].shapes[k + 1].text_frame.paragraphs[0].font.size = Pt(18)
                prs.slides[islide].shapes[k + 2].text_frame.paragraphs[0].font.size = Pt(18)
                prs.slides[islide].shapes[k].text = str('{0:,}'.format(0))
                prs.slides[islide].shapes[k + 1].text = str('{0:,}'.format(0))
                prs.slides[islide].shapes[k + 2].text = str('{0:,}'.format(0))

                for tipo in obraMax:

                    if tipo['tipoObra_id'] == 3:
                        prs.slides[islide].shapes[k].text = str('{0:,}'.format(tipo['numero_obras']))

                    if tipo['tipoObra_id'] == 2:
                        prs.slides[islide].shapes[k + 1].text = str('{0:,}'.format(tipo['numero_obras']))

                    if tipo['tipoObra_id'] == 1:
                        prs.slides[islide].shapes[k + 2].text = str('{0:,}'.format(tipo['numero_obras']))

                k = 22

            islide += 1

        output = StringIO.StringIO()
        prs.save(output)
        response = StreamingHttpResponse(FileWrapper(output),
                                         content_type='application/vnd.openxmlformats-officedocument.presentationml.presentation')
        response['Content-Disposition'] = 'attachment; filename="presentacionObrasAPF.pptx"'
        response['Content-Length'] = output.tell()

        output.seek(0)

        return response
