# Create your views here.
import os, sys
from django.template import RequestContext, loader
from django.http import HttpResponse
from django.db.models import Sum,Count
from django.contrib.auth.decorators import login_required, user_passes_test
from decimal import *
import itertools
from obras.tools import *
from obras.BuscarObras import BuscaObra
from pptx.util import Pt
import json
# from pptx import Presentation
from obras.models import *
from obras.models import Obra
import datetime
import tempfile
#from flask import send_file

from obras.BuscarObras import BuscarObras
from django.shortcuts import render_to_response, redirect
from oauth2_provider.models import AccessToken

from pptx import Presentation
from django.core.servers.basehttp import FileWrapper
import mimetypes
from django.http import StreamingHttpResponse
from pptx.util import Inches
from pptx.dml.color import RGBColor
import PIL  # MODULO PARA PROCESAR IMAGENES
from PIL import Image


def get_user_for_token(token):
    if token:
        return AccessToken.objects.get(token=token).user
    else:
        return None

def register_by_access_token(request):

    #del request.session['access_token']

    if request.session.get('access_token'):
        token = {
        'access_token': request.session.get('access_token'),
        'token_type': 'Bearer'
    }
        return JsonResponse(token)
    else:
        #user = get_user_for_token('3DVteYz9OIH6gvQDyYX78GOpHKXgPy'
        user = request.user
        return get_access_token(user,request)

def ayuda(request):
    return render_to_response('admin/obras/ayuda/c_ayuda.html', locals(),
                              context_instance=RequestContext(request))
def videos(request):
    return render_to_response('admin/obras/videos/videos_lista.html', locals(),
                              context_instance=RequestContext(request))
def manualesPdf(request):
    return render_to_response('admin/obras/manuales/manuales_lista.html', locals(),
                              context_instance=RequestContext(request))

def catalogo(request):
    return render_to_response('admin/obras/catalogo.html', locals(),
                              context_instance=RequestContext(request))


def c_clasificacion(request):
    return render_to_response('admin/obras/c_clasificacion.html', locals(),
                              context_instance=RequestContext(request))


def c_dependencia(request):
    return render_to_response('admin/obras/c_dependencia.html', locals(),
                              context_instance=RequestContext(request))

def c_subdependencia(request):
    return render_to_response('admin/obras/c_subdependencia.html', locals(),
                              context_instance=RequestContext(request))


def c_impacto(request):
    return render_to_response('admin/obras/c_impacto.html', locals(),
                              context_instance=RequestContext(request))


def c_inaugurador(request):
    return render_to_response('admin/obras/c_inaugurador.html', locals(),
                              context_instance=RequestContext(request))


def c_inversion(request):
    return render_to_response('admin/obras/c_inversion.html', locals(),
                              context_instance=RequestContext(request))


def movimientos(request):
    return render_to_response('admin/obras/movimientos.html', locals(),
                              context_instance=RequestContext(request))


def modifica(request):
    return render_to_response('admin/obras/obra/modifica.html', locals(),
                              context_instance=RequestContext(request))


def consultas(request):
    return render_to_response('admin/obras/consultas.html', locals(),
                              context_instance=RequestContext(request))


def c_filtro(request):
    return render_to_response('admin/obras/c_filtro.html', locals(),
                              context_instance=RequestContext(request))


def c_guardada(request):
    return render_to_response('admin/obras/c_guardada.html', locals(),
                              context_instance=RequestContext(request))


def c_predefinida(request):
    return render_to_response('admin/obras/c_predefinida.html', locals(),
                              context_instance=RequestContext(request))


def usuarios(request):
    return render_to_response('admin/obras/usuarios.html', locals(),
                              context_instance=RequestContext(request))


def is_super_admin(user):
    return user.usuario.rol == 'SA'


@login_required()
@user_passes_test(is_super_admin)
def balance_general(request):
    start_date = datetime.date(2012, 12, 01)
    end_date = datetime.date(2013, 12, 31)
    obras2013 = Obra.objects.filter(
        Q(fechaTermino__range=(start_date, end_date)),
        Q(tipoObra=3)
    )

    total_obras_2013 = obras2013.count()
    total_invertido_2013 = obras2013.aggregate(Sum('inversionTotal'))

    # print 'Total Obras 2013: ' + str(total_obras_2013)
    # print 'Total Invertido 2013: ' + str(total_invertido_2013)

    start_date = datetime.date(2014, 01, 01)
    end_date = datetime.date(2014, 12, 31)
    obras2014 = Obra.objects.filter(
        Q(fechaTermino__range=(start_date, end_date)),
        Q(tipoObra=3)
    )

    total_obras_2014 = obras2014.count()
    total_invertido_2014 = obras2014.aggregate(Sum('inversionTotal'))

    # print 'Total Obras 2014: ' + str(total_obras_2014)
    # print 'Total Invertido 2014: ' + str(total_invertido_2014)

    obras_proceso = Obra.objects.filter(
        Q(tipoObra=2)
    )

    total_obras_proceso = obras_proceso.count()
    total_invertido_proceso = obras_proceso.aggregate(Sum('inversionTotal'))

    # print 'Total Obras Proceso: ' + str(total_obras_proceso)
    # print 'Total Invertido Proceso: ' + str(total_invertido_proceso)

    obras_proyectadas = Obra.objects.filter(
        Q(tipoObra=1)
    )

    total_obras_proyectadas = obras_proyectadas.count()
    total_invertido_proyectadas = obras_proyectadas.aggregate(Sum('inversionTotal'))

    # print 'Total Obras Proyectadas: ' + str(total_obras_proyectadas)
    # print 'Total Invertido Proyectadas: ' + str(total_invertido_proyectadas)

    total_obras = total_obras_2013 + total_obras_2014 + total_obras_proceso + total_obras_proyectadas
    total_invertido = total_invertido_2013.get('inversionTotal__sum') + total_invertido_2014.get(
        'inversionTotal__sum') + total_invertido_proceso.get('inversionTotal__sum') + total_invertido_proyectadas.get(
        'inversionTotal__sum')

    template = loader.get_template('reportes/balance_general.html')
    context = RequestContext(request, {
        'total_obras_2013': total_obras_2013,
        'total_invertido_2013': total_invertido_2013,
        'total_obras_2014': total_obras_2014,
        'total_invertido_2014': total_invertido_2014,
        'total_obras_proceso': total_obras_proceso,
        'total_invertido_proceso': total_invertido_proceso,
        'total_obras_proyectadas': total_obras_proyectadas,
        'total_invertido_proyectadas': total_invertido_proyectadas,
        'total_obras': total_obras,
        'total_invertido': total_invertido,
    })
    return HttpResponse(template.render(context))


@login_required()
@user_passes_test(is_super_admin)
def hipervinculo_informacion_general(request):
    obras_concluidas = Obra.objects.filter(
        Q(tipoObra=3)
    )

    obras_proceso = Obra.objects.filter(
        Q(tipoObra=2)
    )

    obras_proyectadas = Obra.objects.filter(
        Q(tipoObra=1)
    )

    total_obras_proyectadas = obras_proyectadas.count()
    total_obras_proceso = obras_proceso.count()
    total_obras_concluidas = obras_concluidas.count()

    total_invertido_proyectadas = obras_proyectadas.aggregate(Sum('inversionTotal'))
    total_invertido_proceso = obras_proceso.aggregate(Sum('inversionTotal'))
    total_invertido_concluidas = obras_concluidas.aggregate(Sum('inversionTotal'))

    template = loader.get_template('reportes/hipervinculo_informacion_general.html')
    context = RequestContext(request, {
        'total_obras_proyectadas': total_obras_proyectadas,
        'total_obras_proceso': total_obras_proceso,
        'total_obras_concluidas': total_obras_concluidas,
        'total_invertido_proyectadas': total_invertido_proyectadas,
        'total_invertido_proceso': total_invertido_proceso,
        'total_invertido_concluidas': total_invertido_concluidas,
    })
    return HttpResponse(template.render(context))


@login_required()
@user_passes_test(is_super_admin)
def hipervinculo_sector(request):
    start_date_2013 = datetime.date(2012, 12, 01)
    end_date_2013 = datetime.date(2013, 12, 31)

    start_date_2014 = datetime.date(2014, 01, 01)
    end_date_2014 = datetime.date(2014, 12, 31)
    dependencias = {}

    for dependencia in Dependencia.objects.all():
        print dependencia.nombreDependencia

        obras_2013_concluidas = Obra.objects.filter(
            Q(fechaTermino__range=(start_date_2013, end_date_2013)),
            Q(tipoObra=3),
            Q(dependencia=dependencia),
        )

        obras_2014_concluidas = Obra.objects.filter(
            Q(fechaTermino__range=(start_date_2014, end_date_2014)),
            Q(tipoObra=3),
            Q(dependencia=dependencia),
        )

        obras_2014_proceso = Obra.objects.filter(
            Q(fechaTermino__range=(start_date_2014, end_date_2014)),
            Q(tipoObra=2),
            Q(dependencia=dependencia),
        )

        obras_2014_proyectadas = Obra.objects.filter(
            Q(fechaTermino__range=(start_date_2014, end_date_2014)),
            Q(tipoObra=1),
            Q(dependencia=dependencia),
        )

        total_obras_concluidas_2013 = obras_2013_concluidas.count()
        total_obras_concluidas_2014 = obras_2014_concluidas.count()
        total_obras_proceso = obras_2014_proceso.count()
        total_obras_proyectadas = obras_2014_proyectadas.count()

        total_invertido_2013 = obras_2013_concluidas.aggregate(Sum('inversionTotal'))
        total_invertido_2014 = obras_2014_concluidas.aggregate(Sum('inversionTotal'))
        total_invertido_proceso = obras_2014_proceso.aggregate(Sum('inversionTotal'))
        total_invertido_proyectadas = obras_2014_proyectadas.aggregate(Sum('inversionTotal'))

        dependencias[dependencia.nombreDependencia] = {}
        dependencias[dependencia.nombreDependencia]['total_obras_concluidas_2013'] = total_obras_concluidas_2013
        dependencias[dependencia.nombreDependencia]['total_obras_concluidas_2014'] = total_obras_concluidas_2014
        dependencias[dependencia.nombreDependencia]['total_obras_proceso'] = total_obras_proceso
        dependencias[dependencia.nombreDependencia]['total_obras_proyectadas'] = total_obras_proyectadas

        dependencias[dependencia.nombreDependencia]['total_invertido_2013'] = total_invertido_2013
        dependencias[dependencia.nombreDependencia]['total_invertido_2014'] = total_invertido_2014
        dependencias[dependencia.nombreDependencia]['total_invertido_proceso'] = total_invertido_proceso
        dependencias[dependencia.nombreDependencia]['total_invertido_proyectadas'] = total_invertido_proyectadas

        print 'total_obras_concluidas_2013: ' + str(total_obras_concluidas_2013)
        print 'total_obras_concluidas_2014: ' + str(total_obras_concluidas_2014)
        print 'total_obras_proceso: ' + str(total_obras_proceso)
        print 'total_obras_proyectadas: ' + str(total_obras_proyectadas)

        print 'total_invertido_2013: ' + str(total_invertido_2013)
        print 'total_invertido_2014: ' + str(total_invertido_2014)
        print 'total_invertido_proceso: ' + str(total_invertido_proceso)
        print 'total_invertido_proyectadas: ' + str(total_invertido_proyectadas)

    template = loader.get_template('reportes/hipervinculo_informacion_sector.html')
    context = RequestContext(request, {
        'dependencias': dependencias,
    })
    print(dependencias)
    return HttpResponse(template.render(context))


@login_required()
@user_passes_test(is_super_admin)
def hipervinculo_entidad(request):
    start_date_2013 = datetime.date(2012, 12, 01)
    end_date_2013 = datetime.date(2013, 12, 31)

    start_date_2014 = datetime.date(2014, 01, 01)
    end_date_2014 = datetime.date(2014, 12, 31)

    start_date_2015 = datetime.date(2015, 01, 01)
    end_date_2015 = datetime.date(2015, 12, 31)

    start_date_2016 = datetime.date(2016, 01, 01)
    end_date_2016 = datetime.date(2016, 12, 31)

    start_date_2017 = datetime.date(2017, 01, 01)
    end_date_2017 = datetime.date(2017, 12, 31)

    start_date_2018 = datetime.date(2018, 01, 01)
    end_date_2018 = datetime.date(2018, 12, 31)
    estados = {}

    for estado in Estado.objects.all():
        print estado.nombreEstado

        obras_2013_concluidas = Obra.objects.filter(
            Q(fechaTermino__range=(start_date_2013, end_date_2013)),
            Q(tipoObra=3),
            Q(estado=estado),
        )

        obras_2014_concluidas = Obra.objects.filter(
            Q(fechaTermino__range=(start_date_2014, end_date_2014)),
            Q(tipoObra=3),
            Q(estado=estado),
        )

        obras_proceso = Obra.objects.filter(
            Q(tipoObra=2),
            Q(estado=estado),
        )

        obras_2014_proyectadas = Obra.objects.filter(
            Q(fechaTermino__range=(start_date_2014, end_date_2014)),
            Q(tipoObra=1),
            Q(estado=estado),
        )

        obras_2015_proyectadas = Obra.objects.filter(
            Q(fechaTermino__range=(start_date_2015, end_date_2015)),
            Q(tipoObra=1),
            Q(estado=estado),
        )

        obras_2016_proyectadas = Obra.objects.filter(
            Q(fechaTermino__range=(start_date_2016, end_date_2016)),
            Q(tipoObra=1),
            Q(estado=estado),
        )

        obras_2017_proyectadas = Obra.objects.filter(
            Q(fechaTermino__range=(start_date_2017, end_date_2017)),
            Q(tipoObra=1),
            Q(estado=estado),
        )

        obras_2018_proyectadas = Obra.objects.filter(
            Q(fechaTermino__range=(start_date_2018, end_date_2018)),
            Q(tipoObra=1),
            Q(estado=estado),
        )

        total_obras_concluidas_2013 = obras_2013_concluidas.count()
        total_obras_concluidas_2014 = obras_2014_concluidas.count()
        total_obras_proceso = obras_proceso.count()
        total_obras_proyectadas_2014 = obras_2014_proyectadas.count()
        total_obras_proyectadas_2015 = obras_2015_proyectadas.count()
        total_obras_proyectadas_2016 = obras_2016_proyectadas.count()
        total_obras_proyectadas_2017 = obras_2017_proyectadas.count()
        total_obras_proyectadas_2018 = obras_2018_proyectadas.count()

        total_invertido_2013_concluidas = obras_2013_concluidas.aggregate(Sum('inversionTotal'))
        total_invertido_2014_concluidas = obras_2014_concluidas.aggregate(Sum('inversionTotal'))
        total_invertido_proceso = obras_proceso.aggregate(Sum('inversionTotal'))

        total_invertido_proyectadas_2014 = obras_2014_proyectadas.aggregate(Sum('inversionTotal'))
        total_invertido_proyectadas_2015 = obras_2015_proyectadas.aggregate(Sum('inversionTotal'))
        total_invertido_proyectadas_2016 = obras_2016_proyectadas.aggregate(Sum('inversionTotal'))
        total_invertido_proyectadas_2017 = obras_2017_proyectadas.aggregate(Sum('inversionTotal'))
        total_invertido_proyectadas_2018 = obras_2018_proyectadas.aggregate(Sum('inversionTotal'))

        estados[estado.nombreEstado] = {}
        estados[estado.nombreEstado]['total_obras_concluidas_2013'] = total_obras_concluidas_2013
        estados[estado.nombreEstado]['total_obras_concluidas_2014'] = total_obras_concluidas_2014
        estados[estado.nombreEstado]['total_obras_proceso'] = total_obras_proceso
        estados[estado.nombreEstado]['total_obras_proyectadas_2014'] = total_obras_proyectadas_2014
        estados[estado.nombreEstado]['total_obras_proyectadas_2015'] = total_obras_proyectadas_2015
        estados[estado.nombreEstado]['total_obras_proyectadas_2016'] = total_obras_proyectadas_2016
        estados[estado.nombreEstado]['total_obras_proyectadas_2017'] = total_obras_proyectadas_2017
        estados[estado.nombreEstado]['total_obras_proyectadas_2018'] = total_obras_proyectadas_2018

        estados[estado.nombreEstado]['total_invertido_2013_concluidas'] = total_invertido_2013_concluidas
        estados[estado.nombreEstado]['total_invertido_2014_concluidas'] = total_invertido_2014_concluidas
        estados[estado.nombreEstado]['total_invertido_proceso'] = total_invertido_proceso

        estados[estado.nombreEstado]['total_invertido_proyectadas_2014'] = total_invertido_proyectadas_2014
        estados[estado.nombreEstado]['total_invertido_proyectadas_2015'] = total_invertido_proyectadas_2015
        estados[estado.nombreEstado]['total_invertido_proyectadas_2016'] = total_invertido_proyectadas_2016
        estados[estado.nombreEstado]['total_invertido_proyectadas_2017'] = total_invertido_proyectadas_2017
        estados[estado.nombreEstado]['total_invertido_proyectadas_2018'] = total_invertido_proyectadas_2018

    template = loader.get_template('reportes/hipervinculo_entidad.html')
    context = RequestContext(request, {
        'estados': estados,
    })
    return HttpResponse(template.render(context))


@login_required()
@user_passes_test(is_super_admin)
def hipervinculo_concluidas_proceso_proyectadas(request):
    start_date_2013 = datetime.date(2012, 12, 01)
    end_date_2013 = datetime.date(2013, 12, 31)

    start_date_2014 = datetime.date(2014, 01, 01)
    end_date_2014 = datetime.date(2014, 12, 31)

    obras_2013_concluidas = Obra.objects.filter(
        Q(fechaTermino__range=(start_date_2013, end_date_2013)),
        Q(tipoObra=3),
    )

    obras_2014_concluidas = Obra.objects.filter(
        Q(fechaTermino__range=(start_date_2014, end_date_2014)),
        Q(tipoObra=3),
    )

    total_obras_concluidas_2013 = obras_2013_concluidas.count()
    total_invertido_2013_concluidas = obras_2013_concluidas.aggregate(Sum('inversionTotal'))
    total_obras_concluidas_2014 = obras_2014_concluidas.count()
    total_invertido_2014_concluidas = obras_2014_concluidas.aggregate(Sum('inversionTotal'))

    print 'total_obras_concluidas_2013: ' + str(total_obras_concluidas_2013)
    print 'total_invertido_2013_concluidas: ' + str(total_invertido_2013_concluidas)
    print 'total_obras_concluidas_2014: ' + str(total_obras_concluidas_2014)
    print 'total_invertido_2014_concluidas: ' + str(total_invertido_2014_concluidas)

    obras_proceso = Obra.objects.filter(
        Q(tipoObra=2),
    )

    obras_totales_proceso = obras_proceso.count()
    total_invertido_proceso = obras_proceso.aggregate(Sum('inversionTotal'))

    print 'obras_totales_proceso: ' + str(obras_totales_proceso)
    print 'total_invertido_proceso: ' + str(total_invertido_proceso)

    obras_proyectadas = Obra.objects.filter(
        Q(tipoObra=1),
    )

    obras_totales_proyectadas = obras_proyectadas.count()
    total_invertido_proyectadas = obras_proyectadas.aggregate(Sum('inversionTotal'))

    print 'obras_totales_proyectadas: ' + str(obras_totales_proyectadas)
    print 'total_invertido_proyectadas: ' + str(total_invertido_proyectadas)

    template = loader.get_template('reportes/hipervinculo_concluidas_proceso_proyectadas.html')
    context = RequestContext(request, {
        'total_obras_concluidas_2013': total_obras_concluidas_2013,
        'total_invertido_2013_concluidas': total_invertido_2013_concluidas,
        'total_obras_concluidas_2014': total_obras_concluidas_2014,
        'total_invertido_2014_concluidas': total_invertido_2014_concluidas,
        'obras_totales_proceso': obras_totales_proceso,
        'total_invertido_proceso': total_invertido_proceso,
        'obras_totales_proyectadas': obras_totales_proyectadas,
        'total_invertido_proyectadas': total_invertido_proyectadas,
    })
    return HttpResponse(template.render(context))


@login_required()
def consulta_web(request):

    print request.user.usuario.rol

    if request.user.usuario.rol == 'SA':
        dependencias = Dependencia.objects.filter(
            Q(obraoprograma='O') &
            Q(dependienteDe__isnull=True)
        )
        subdependencias = Dependencia.objects.filter(
            Q(obraoprograma='O') &
            Q(dependienteDe__isnull=False)
        )
    else:
        dependencias = Dependencia.objects.filter(
            Q(id__in=request.user.usuario.dependencia.all())
            & Q(dependienteDe__isnull=True)
        )
        subdependencias = Dependencia.objects.filter(Q(dependienteDe__in=dependencias))

    #templates = loader.get_template('consultas/busqueda_general.html')
    template = loader.get_template('admin/obras/consulta_filtros/consulta-filtros.html')
    context = RequestContext(request, {
        'estatusObra': TipoObra.objects.all(),
        'dependencias': dependencias,
        'subdependencias': subdependencias,
        'estados': Estado.objects.all(),
        'tipo_inversiones': TipoInversion.objects.all(),
        'impactos': Impacto.objects.all(),
        'clasificacion': TipoClasificacion.objects.all(),
        'inaugurador': Inaugurador.objects.all(),
        'InstanciaEjecutora': InstanciaEjecutora.objects.all(),
    })
    return HttpResponse(template.render(context))

@login_required()
def ver_video(request):
    cualVideo=request.GET.get('cualVideo', None),
    print(str(cualVideo[0]))
    if str(cualVideo[0]) =='creacionObra.mp4':
        tituloVideo='Crear una Obra',
    elif str(cualVideo[0]) =='modificacionObra.mp4':
        tituloVideo='Modificar una Obra',

    elif str(cualVideo[0]) =='consultaFiltros.mp4':
        tituloVideo='Consulta Mediante Filtros',
    elif str(cualVideo[0]) =='consultaPredefinida.mp4':
        tituloVideo='Consulta Predefinidos',

    elif str(cualVideo[0]) =='agregarClasificacion.mp4':
        tituloVideo='Agregar una Clasificacion',
    elif str(cualVideo[0]) =='buscarClasificacion.mp4':
        tituloVideo='Buscar una Clasificacion',
    elif str(cualVideo[0]) =='modificarClasificacion.mp4':
        tituloVideo='Modificar una Clasificacion',
    elif str(cualVideo[0]) =='eliminarClasificacion.mp4':
        tituloVideo='Eliminar una Clasificacion',

    elif str(cualVideo[0]) =='crearUsuario.mp4':
        tituloVideo='Agregar un Usuario',
    elif str(cualVideo[0]) =='modificarUsuario.mp4':
        tituloVideo='Modificar un Usuario',
    elif str(cualVideo[0]) =='buscarUsuario.mp4':
        tituloVideo='Buscar un Usuario',
    elif str(cualVideo[0]) =='eliminarUsuario.mp4':
        tituloVideo='Eliminar un Usuario',

    elif str(cualVideo[0]) =='agregarDependencia.mp4':
        tituloVideo='Agregar una Dependencia',
    elif str(cualVideo[0]) =='buscarDependencia.mp4':
        tituloVideo='Buscar una Dependencia',
    elif str(cualVideo[0]) =='modificarDependencia.mp4':
        tituloVideo='Modificar una Dependencia',
    elif str(cualVideo[0]) =='eliminarDependencia.mp4':
        tituloVideo='Eliminar una Dependencia',

    elif str(cualVideo[0]) =='agregarSubDependencia.mp4':
        tituloVideo='Agregar una SubDependencia',
    elif str(cualVideo[0]) =='buscarSubDependencia.mp4':
        tituloVideo='Buscar una SubDependencia',
    elif str(cualVideo[0]) =='modificarSubDependencia.mp4':
        tituloVideo='Modificar una SubDependencia',
    elif str(cualVideo[0]) =='eliminarSubDependencia.mp4':
        tituloVideo='Eliminar una SubDependencia',


    elif str(cualVideo[0]) =='agregarInversion.mp4':
        tituloVideo='Agregar un tipo de Inversion',
    elif str(cualVideo[0]) =='buscarInversion.mp4':
        tituloVideo='Buscar un tipo de Inversion',
    elif str(cualVideo[0]) =='modificarInversion.mp4':
        tituloVideo='Modificar un tipo de Inversion',
    elif str(cualVideo[0]) =='eliminarInversion.mp4':
        tituloVideo='Eliminar un tipo de Inversion',

    elif str(cualVideo[0]) =='agregarImpacto.mp4':
        tituloVideo='Agregar un tipo de Impacto',
    elif str(cualVideo[0]) =='buscarImpacto.mp4':
        tituloVideo='Buscar un un tipo de Impacto',
    elif str(cualVideo[0]) =='modificarImpacto.mp4':
        tituloVideo='Modificar un un tipo de Impacto',
    elif str(cualVideo[0]) =='eliminarImpacto.mp4':
        tituloVideo='Eliminar un un tipo de Impacto',

    elif str(cualVideo[0]) =='agregarInaugurador.mp4':
        tituloVideo='Agregar un Inaugurador',
    elif str(cualVideo[0]) =='buscarInaugurador.mp4':
        tituloVideo='Buscar un Inaugurador',
    elif str(cualVideo[0]) =='modificarInaugurador.mp4':
        tituloVideo='Modificar un Inaugurador',
    elif str(cualVideo[0]) =='eliminarInaugurador.mp4':
        tituloVideo='Eliminar un Inaugurador',

    template = loader.get_template('admin/obras/videos/videos_lista.html')
    context = RequestContext(request, {
        'cualVideo': cualVideo,
        'tituloVideo': tituloVideo,
    })
    return HttpResponse(template.render(context))

def get_array_or_none(the_string):
    if the_string is None:
        return None
    else:
        return map(int, the_string.split(','))

def get_string_or_none(the_string):
    if the_string is None:
        return None
    else:
        return the_string.split(',')

def buscar_obras_web(request):
    #TODO cambiar los parametros 'None' por get del request
    buscador = BuscarObras(
            idtipoobra=get_array_or_none(request.GET.get('tipoDeObra')),
            iddependencias=get_array_or_none(request.GET.get('dependencia')),
            estados=get_array_or_none(request.GET.get('estado')),
            clasificaciones=get_array_or_none(request.GET.get('clasificacion')),
            inversiones=get_array_or_none(request.GET.get('tipoDeInversion')),
            inauguradores=get_array_or_none(request.GET.get('inaugurador')),
            impactos=get_array_or_none(request.GET.get('impacto')),
            inaugurada=request.GET.get('inaugurada', None),
            inversion_minima=request.GET.get('inversionMinima', None),
            inversion_maxima=request.GET.get('inversionMaxima', None),
            fecha_inicio_primera=request.GET.get('fechaInicio', None),
            fecha_inicio_segunda=request.GET.get('fechaInicio', None),
            fecha_fin_primera=request.GET.get('fechaFin', None),
            fecha_fin_segunda=request.GET.get('fechaFinSegunda', None),
            denominacion=request.GET.get('denominacion', None),
            instancia_ejecutora=get_array_or_none(request.GET.get('instanciaEjecutora')),
            municipios=get_array_or_none(request.GET.get('municipios', None))
    )

    buscador.filtrar_dependencias(request.user)
    resultados = buscador.buscar()

    template = loader.get_template('admin/obras/consulta_filtros/consulta-filtros.html')
    context = RequestContext(request, {
        'resultados': resultados,
        'filtros': request.GET
    })
    return HttpResponse(template.render(context))


def obras_iniciadas(request):
    usuario = request.user.usuario

    query = Q(fechaInicio__lte= datetime.datetime.now().date())
    if not (usuario.rol == 'SA'):
        subdependencias = get_subdependencias_as_list_flat(usuario.dependencia.all())
        query = query & (Q(dependencia__in=subdependencias) | Q(subdependencia__in=subdependencias))
    obras = Obra.objects.filter(query)

    template = loader.get_template('admin/obras/consulta_predefinidos/consulta-predefinidos.html')
    context = RequestContext(request, {
        'obras_resultado': obras
    })
    return HttpResponse(template.render(context))


def obras_vencidas(request):
    usuario = request.user.usuario

    today = datetime.datetime.now().date()
    if usuario.rol == 'SA':
        obras = Obra.objects.filter(fechaTermino__lte=today)
    else:
        obras = Obra.objects.filter(Q(fechaTermino__lte=today) & Q(
            dependencia__in=get_subdependencias_as_list_flat(usuario.dependencia.all())))

    template = loader.get_template('admin/obras/consulta_predefinidos/consulta-predefinidos.html')
    context = RequestContext(request, {
        'obras_resultado': obras
    })
    return HttpResponse(template.render(context))


def obras_for_dependencia(request):
    usuario = request.user.usuario

    if usuario.rol == 'SA':
        obras = Obra.objects.all()
    else:
        obras = usuario.dependencia.get_obras()

    template = loader.get_template('admin/obras/consulta_predefinidos/consulta-predefinidos.html')
    context = RequestContext(request, {
        'obras_resultado': obras
    })
    return HttpResponse(template.render(context))


def ajax_prueba(request):
    template = loader.get_template('prueba.html')
    return HttpResponse(template.render(RequestContext(request)))




# reportes de power point ************************************************************************************

def fichaTecnica(request):
        prs = Presentation('/home/obrasapf/djangoObras/obras/static/ppt/FichaTecnicaObras.pptx')
        #prs = Presentation('obras/static/ppt/FichaTecnicaObras.pptx')
        usuario = request.user.usuario
        buscador = BuscaObra(
            identificador_unico=request.GET.get('identificador_unico', None)
        )
        resultados = buscador.busca()

        json_map = {}
        json_map['obras'] = []
        for obra in resultados['obras']:
            json_map['obras'].append(obra.to_serializable_dict())

        json_map['DInversion'] = []
        for DetalleInversion in resultados['DInversion']:
            json_map['DInversion'].append(DetalleInversion.to_serializable_dict())

        json_map['DClasificacion'] = []
        for DetalleClas in resultados['DClasificacion']:
            map = {}
            map['tipoClasificacion'] = DetalleClas['tipoClasificacion__id']
            if DetalleClas['subclasificacion__nombreTipoClasificacion'] is None:
                map['subClasificacion'] = ""
            else:
                map['subClasificacion'] = DetalleClas['subclasificacion__nombreTipoClasificacion']

            json_map['DClasificacion'].append(map)


        #generales
        prs.slides[0].shapes[8].text_frame.paragraphs[0].font.size = Pt(8)
        prs.slides[0].shapes[8].text = json_map['obras'][0]['identificador_unico']
        prs.slides[0].shapes[9].text_frame.paragraphs[0].font.size = Pt(8)
        prs.slides[0].shapes[9].text = json_map['obras'][0]['denominacion']
        prs.slides[0].shapes[10].text_frame.paragraphs[0].font.size = Pt(8)
        prs.slides[0].shapes[10].text = json_map['obras'][0]['dependencia']['nombreDependencia']
        prs.slides[0].shapes[11].text_frame.paragraphs[0].font.size = Pt(8)
        prs.slides[0].shapes[11].text = json_map['obras'][0]['estado']['nombreEstado']
        prs.slides[0].shapes[12].text_frame.paragraphs[0].font.size = Pt(8)
        try:
               prs.slides[0].shapes[12].text = json_map['obras'][0]['municipio']
        except Exception as e:
            prs.slides[0].shapes[12].text = ""
        #prs.slides[0].shapes[12].text = json_map['obras'][0]['municipio']
        prs.slides[0].shapes[13].text_frame.paragraphs[0].font.size = Pt(8)
        prs.slides[0].shapes[13].text = json_map['obras'][0]['fechaInicio']
        prs.slides[0].shapes[14].text_frame.paragraphs[0].font.size = Pt(8)
        prs.slides[0].shapes[14].text = json_map['obras'][0]['fechaTermino']
        prs.slides[0].shapes[15].text_frame.paragraphs[0].font.size = Pt(8)
        prs.slides[0].shapes[15].text = json_map['obras'][0]['tipoObra']['nombreTipoObra']

        prs.slides[0].shapes[16].text_frame.paragraphs[0].font.size = Pt(8)
        prs.slides[0].shapes[16].text = str(json_map['obras'][0]['porcentajeAvance'])
        prs.slides[0].shapes[17].text_frame.paragraphs[0].font.size = Pt(8)
        prs.slides[0].shapes[17].text = json_map['obras'][0]['fechaModificacion']
        #detalle inversion
        prs.slides[0].shapes[18].text_frame.paragraphs[0].font.size = Pt(8)
        prs.slides[0].shapes[18].text = "No"
        prs.slides[0].shapes[19].text_frame.paragraphs[0].font.size = Pt(8)
        prs.slides[0].shapes[19].text = "No"
        prs.slides[0].shapes[20].text_frame.paragraphs[0].font.size = Pt(8)
        prs.slides[0].shapes[20].text = "No"
        prs.slides[0].shapes[21].text_frame.paragraphs[0].font.size = Pt(8)
        prs.slides[0].shapes[21].text = "No"
        prs.slides[0].shapes[22].text_frame.paragraphs[0].font.size = Pt(8)
        prs.slides[0].shapes[22].text = "No"
        prs.slides[0].shapes[23].text_frame.paragraphs[0].font.size = Pt(8)
        prs.slides[0].shapes[23].text = "No"
        prs.slides[0].shapes[24].text_frame.paragraphs[0].font.size = Pt(8)
        prs.slides[0].shapes[24].text = "N/A"

        for DI in json_map['DInversion']:
            if not (DI['tipoInversion'] is None):
                if DI['tipoInversion'] == 1:
                    prs.slides[0].shapes[18].text_frame.paragraphs[0].font.size = Pt(7)
                    prs.slides[0].shapes[18].text = str(DI['monto'])
                if DI['tipoInversion'] == 2:
                    prs.slides[0].shapes[19].text_frame.paragraphs[0].font.size = Pt(7)
                    prs.slides[0].shapes[19].text = str(DI['monto'])
                if DI['tipoInversion'] == 3:
                    prs.slides[0].shapes[20].text_frame.paragraphs[0].font.size = Pt(7)
                    prs.slides[0].shapes[20].text = str(DI['monto'])
                if DI['tipoInversion'] == 4:
                    prs.slides[0].shapes[21].text_frame.paragraphs[0].font.size = Pt(7)
                    prs.slides[0].shapes[21].text = str(DI['monto'])
                if DI['tipoInversion'] == 5:
                    prs.slides[0].shapes[22].text_frame.paragraphs[0].font.size = Pt(7)
                    prs.slides[0].shapes[22].text = str(DI['monto'])
                if DI['tipoInversion'] == 6:
                    prs.slides[0].shapes[23].text_frame.paragraphs[0].font.size = Pt(7)
                    prs.slides[0].shapes[23].text = str(DI['monto'])

        prs.slides[0].shapes[25].text_frame.paragraphs[0].font.size = Pt(8)
        prs.slides[0].shapes[25].text = str(json_map['obras'][0]['inversionTotal'])
        prs.slides[0].shapes[26].text_frame.paragraphs[0].font.size = Pt(8)
        prs.slides[0].shapes[26].text = json_map['obras'][0]['tipoMoneda']['nombreTipoDeMoneda']
        prs.slides[0].shapes[42].text_frame.paragraphs[0].font.size = Pt(8)
        prs.slides[0].shapes[42].text = json_map['obras'][0]['registroHacendario']
        #poblacion
        prs.slides[0].shapes[43].text_frame.paragraphs[0].font.size = Pt(8)
        prs.slides[0].shapes[43].text = str(json_map['obras'][0]['totalBeneficiarios'])

        prs.slides[0].shapes[27].text_frame.paragraphs[0].font.size = Pt(8)
        prs.slides[0].shapes[27].text = json_map['obras'][0]['poblacionObjetivo']
        prs.slides[0].shapes[28].text_frame.paragraphs[0].font.size = Pt(8)
        try:
               prs.slides[0].shapes[28].text = json_map['obras'][0]['impacto']['nombreImpacto']
        except Exception as e:
            prs.slides[0].shapes[28].text = ""
        #prs.slides[0].shapes[28].text = json_map['obras'][0]['impacto']['nombreImpacto']
        prs.slides[0].shapes[29].text_frame.paragraphs[0].font.size = Pt(8)
        if json_map['obras'][0]['senalizacion']=="false":
            prs.slides[0].shapes[29].text = "No"
        else:
            prs.slides[0].shapes[29].text = "Si"

        #clasificacion y subclasificacion
        prs.slides[0].shapes[30].text_frame.paragraphs[0].font.size = Pt(8)
        prs.slides[0].shapes[30].text = "No"
        prs.slides[0].shapes[31].text_frame.paragraphs[0].font.size = Pt(8)
        prs.slides[0].shapes[31].text = ""
        prs.slides[0].shapes[32].text_frame.paragraphs[0].font.size = Pt(8)
        prs.slides[0].shapes[32].text = "No"
        prs.slides[0].shapes[33].text_frame.paragraphs[0].font.size = Pt(8)
        prs.slides[0].shapes[33].text = "No"
        prs.slides[0].shapes[34].text_frame.paragraphs[0].font.size = Pt(8)
        prs.slides[0].shapes[34].text = "No"
        prs.slides[0].shapes[41].text_frame.paragraphs[0].font.size = Pt(8)
        prs.slides[0].shapes[41].text = ""
        prs.slides[0].shapes[35].text_frame.paragraphs[0].font.size = Pt(8)
        prs.slides[0].shapes[35].text = "No"
        prs.slides[0].shapes[36].text_frame.paragraphs[0].font.size = Pt(8)
        prs.slides[0].shapes[36].text = "No"
        for DC in json_map['DClasificacion']:
            if DC['tipoClasificacion'] == 1:
                prs.slides[0].shapes[30].text_frame.paragraphs[0].font.size = Pt(8)
                prs.slides[0].shapes[30].text = "Si"
                prs.slides[0].shapes[31].text_frame.paragraphs[0].font.size = Pt(8)
                prs.slides[0].shapes[31].text = DC['subClasificacion']
            if DC['tipoClasificacion'] == 2:
                prs.slides[0].shapes[32].text_frame.paragraphs[0].font.size = Pt(8)
                prs.slides[0].shapes[32].text = "Si"
            if DC['tipoClasificacion'] == 3:
                prs.slides[0].shapes[33].text_frame.paragraphs[0].font.size = Pt(8)
                prs.slides[0].shapes[33].text = "Si"
            if DC['tipoClasificacion'] == 4:
                prs.slides[0].shapes[34].text_frame.paragraphs[0].font.size = Pt(8)
                prs.slides[0].shapes[34].text = "Si"
                prs.slides[0].shapes[41].text_frame.paragraphs[0].font.size = Pt(8)
                prs.slides[0].shapes[41].text = DC['subClasificacion']
            if DC['tipoClasificacion'] == 5:
                prs.slides[0].shapes[35].text_frame.paragraphs[0].font.size = Pt(8)
                prs.slides[0].shapes[35].text = "Si"
            if DC['tipoClasificacion'] == 6:
                prs.slides[0].shapes[36].text_frame.paragraphs[0].font.size = Pt(8)
                prs.slides[0].shapes[36].text = "Si"

        #descripcion
        prs.slides[0].shapes[37].text_frame.paragraphs[0].font.size = Pt(8)
        prs.slides[0].shapes[37].text = json_map['obras'][0]['descripcion']
        #observaciones
        prs.slides[0].shapes[38].text_frame.paragraphs[0].font.size = Pt(8)
        prs.slides[0].shapes[38].text = json_map['obras'][0]['observaciones']
        #inauguracion
        prs.slides[0].shapes[39].text_frame.paragraphs[0].font.size = Pt(8)
        if json_map['obras'][0]['inaugurada'] == False:
            prs.slides[0].shapes[39].text = "No"
        else:
            prs.slides[0].shapes[39].text = "Si"

        #prs.slides[0].shapes[40].text_frame.paragraphs[0].font.size = Pt(8)
        #if json_map['obras'][0]['susceptibleInauguracion'] == "false":
        #    prs.slides[0].shapes[40].text = "No"
        #else:
        #    prs.slides[0].shapes[40].text = "Si"

        try:
            prs.slides[0].shapes[40].text_frame.paragraphs[0].font.size = Pt(8)
            prs.slides[0].shapes[40].text = json_map['obras'][0]['inaugurador']['nombreCargoInaugura']
        except Exception as e:
            prs.slides[0].shapes[40].text = ""
        #logo dependencia
        top = Inches(1)
        left = Inches(0.4)

        pic = prs.slides[0].shapes.add_picture('/home/obrasapf/djangoObras/obras' + json_map['obras'][0]['dependencia']['imagenDependencia'], left, top)

        #imagenes de la obra
        left = Inches(7.08)
        top = Inches(5.93)
        widthP = Inches(0.76)
        heightP = Inches(0.78)


        if (json_map['obras'][0]['fotoAntes']) != "":
            scad=(json_map['obras'][0]['fotoAntes']).split('/',2)
            img = Image.open('/home/obrasapf/djangoObras/obras/media/' + json_map['obras'][0]['fotoAntes'])  # ABRIMOS LA IMAGEN PARA TRABAJAR SOBRE ELLA
            width = img.size[0] # CHEQUEAMOS EL ANCHO
            heigh = img.size[1] # CHEQUEAMOS EL ALTO
            if width > 800 or heigh > 800:
                cadena = redimensiona('/home/obrasapf/djangoObras/obras/media/' + scad[0]+ "/"+scad[1] + "/antes_resize.jpg",img,width,heigh)
                pic = prs.slides[0].shapes.add_picture('/home/obrasapf/djangoObras/obras/media/' + scad[0]+ "/"+scad[1] + "/antes_resize.jpg", left, top, widthP, heightP)
            else:
                pic = prs.slides[0].shapes.add_picture("/home/obrasapf/djangoObras/obras/media/" + json_map['obras'][0]['fotoAntes'], left, top, widthP, heightP)

        left = Inches(7.9291)
        top = Inches(5.93)
        if (json_map['obras'][0]['fotoDurante']) != "":
            scad=(json_map['obras'][0]['fotoDurante']).split('/',2)
            img = Image.open('/home/obrasapf/djangoObras/obras/media/' + json_map['obras'][0]['fotoDurante'])  # ABRIMOS LA IMAGEN PARA TRABAJAR SOBRE ELLA
            width = img.size[0] # CHEQUEAMOS EL ANCHO
            heigh = img.size[1] # CHEQUEAMOS EL ALTO
            if width > 800 or heigh > 800:
                cadena = redimensiona('/home/obrasapf/djangoObras/obras/media/' + scad[0]+ "/"+scad[1] + "/durante_resize.jpg",img,width,heigh)
                pic = prs.slides[0].shapes.add_picture('/home/obrasapf/djangoObras/obras/media/' + scad[0]+ "/"+scad[1] + "/durante_resize.jpg", left, top, widthP, heightP)
            else:
                pic = prs.slides[0].shapes.add_picture("/home/obrasapf/djangoObras/obras/media/" + json_map['obras'][0]['fotoDurante'], left, top, widthP, heightP)


        left = Inches(8.7677)
        top = Inches(5.93)
        if (json_map['obras'][0]['fotoDespues']) != "":
            scad=(json_map['obras'][0]['fotoDespues']).split('/',2)
            img = Image.open('/home/obrasapf/djangoObras/obras/media/' + json_map['obras'][0]['fotoDespues'])  # ABRIMOS LA IMAGEN PARA TRABAJAR SOBRE ELLA
            width = img.size[0] # CHEQUEAMOS EL ANCHO
            heigh = img.size[1] # CHEQUEAMOS EL ALTO
            if width > 800 or heigh > 800:
                cadena = redimensiona('/home/obrasapf/djangoObras/obras/media/' + scad[0]+ "/"+scad[1] + "/despues_resize.jpg",img,width,heigh)
                pic = prs.slides[0].shapes.add_picture('/home/obrasapf/djangoObras/obras/media/' + scad[0]+ "/"+scad[1] + "/despues_resize.jpg", left, top, widthP, heightP)
            else:
                pic = prs.slides[0].shapes.add_picture("/home/obrasapf/djangoObras/obras/media/" + json_map['obras'][0]['fotoDespues'], left, top, widthP, heightP)




        prs.save('/home/obrasapf/djangoObras/obras/static/ppt/ppt-generados/FichaTecnicaObras_' + str(usuario.user.id) + '.pptx')

        the_file = '/home/obrasapf/djangoObras/obras/static/ppt/ppt-generados/FichaTecnicaObras_' + str(usuario.user.id) + '.pptx'

        #prs.save('obras/static/ppt/ppt-generados/FichaTecnicaObras_' + str(usuario.user.id) + '.pptx')

        #the_file = 'obras/static/ppt/ppt-generados/FichaTecnicaObras_' + str(usuario.user.id) + '.pptx'



        filename = os.path.basename(the_file)
        chunk_size = 8192
        response = StreamingHttpResponse(FileWrapper(open(the_file,"rb"), chunk_size),
                               content_type=mimetypes.guess_type(the_file)[0])
        response['Content-Length'] = os.path.getsize(the_file)
        response['Content-Disposition'] = "attachment; filename=%s" % filename
        return response


        #print(json_map)
        #return HttpResponse(json.dumps(json_map), 'application/json')

def redimensiona(ruta,img,width,heigh):


    if width > heigh: # SI EL ANCHO ES MAYOR QUE EL ALTO (FOTO HORIZONTAL), LO TOMAMOS COMO REFERENCIA
                basewidth = 400
                wpercent = (basewidth / float(img.size[0]))
                hsize = int((float(img.size[1]) * float(wpercent)))
                img = img.resize((basewidth, hsize), PIL.Image.ANTIALIAS)
                img.save(ruta)  # SALVAMOS LA IMAGEN EN EL DIRECTORIO
    else: # SI EL ALTO ES MAYOR QUE EL ANCHO (FOTO VERTICAL) LO TOMAMOS COMO REFERENCIA
                baseheight = 400
                hpercent = (baseheight / float(img.size[1]))
                wsize = int((float(img.size[0]) * float(hpercent)))
                img = img.resize((wsize, baseheight), PIL.Image.ANTIALIAS)
                img.save(ruta) # SALVAMOS LA IMAGEN EN EL DIRECTORIO


@login_required()
def reportes_predefinidos(request):
    return render_to_response('admin/obras/consulta_predefinidos/consulta-predefinidos.html', {'clases': ''}, context_instance=RequestContext(request))

def abrir_pptx(archivo):
    f = os.popen(archivo,'r')
    prs = Presentation(f)
    f.close()

@login_required()
def balance_general_ppt(request):
    usuario = request.user.usuario
    dependencias = usuario.dependencia.all()
    subdependencias = usuario.subdependencia.all()
    query = Q()
    if dependencias and dependencias.count() > 0:
        if usuario.rol == 'US':
            query = Q(subdependencia__in=get_subdependencias_as_list_flat(subdependencias))
        else:

            query = Q(dependencia__in=get_subdependencias_as_list_flat(dependencias)) | Q(subdependencia__in=get_subdependencias_as_list_flat(dependencias)
        )



    prs = Presentation('/home/obrasapf/djangoObras/obras/static/ppt/PRINCIPAL_BALANCE_GENERAL_APF.pptx')

    # informacion para el 2013
    start_date = datetime.date(2012, 12, 01)
    end_date = datetime.date(2013, 12, 31)
    obras2013 = Obra.objects.filter(
        Q(fechaTermino__range=(start_date, end_date)),
        Q(tipoObra=3),query
    )

    total_obras_2013 = obras2013.count()
    total_invertido_2013 = obras2013.aggregate(Sum('inversionTotal'))

    # informacion para el 2014
    start_date = datetime.date(2014, 01, 01)
    end_date = datetime.date(2014, 12, 31)
    obras2014 = Obra.objects.filter(
        Q(fechaTermino__range=(start_date, end_date)),
        Q(tipoObra=3),query
    )

    total_obras_2014 = obras2014.count()
    total_invertido_2014 = obras2014.aggregate(Sum('inversionTotal'))

    # informacion para obras concluidas
    obras_concluidas = Obra.objects.filter(
        Q(tipoObra=3),query
    )

    total_obras_concluidas = obras_concluidas.count()
    total_invertido_concluidas = obras_concluidas.aggregate(Sum('inversionTotal'))

    # informacion para obras en proceso
    obras_proceso = Obra.objects.filter(
        Q(tipoObra=2),query
    )

    total_obras_proceso = obras_proceso.count()
    total_invertido_proceso = obras_proceso.aggregate(Sum('inversionTotal'))

    # informacion para obras proyectadas
    obras_proyectadas = Obra.objects.filter(
        Q(tipoObra=1),query
    )

    total_obras_proyectadas = obras_proyectadas.count()
    total_invertido_proyectadas = obras_proyectadas.aggregate(Sum('inversionTotal'))

    # informacion para obras totales
    total_obras = total_obras_concluidas + total_obras_proceso + total_obras_proyectadas


    total_invertido = 0
    if total_invertido_concluidas.get('inversionTotal__sum',0):
        total_invertido = total_invertido + total_invertido_concluidas.get('inversionTotal__sum',0)
    if total_invertido_proceso.get('inversionTotal__sum',0):
        total_invertido = total_invertido + total_invertido_proceso.get('inversionTotal__sum',0)
    if total_invertido_proyectadas.get('inversionTotal__sum',0):
        total_invertido = total_invertido + total_invertido_proyectadas.get('inversionTotal__sum',0)

    totalinvertidoconcluidas=0
    if total_invertido_concluidas.get('inversionTotal__sum',0):
        totalinvertidoconcluidas = total_invertido_concluidas.get('inversionTotal__sum',0)

    totalinvertidoproceso=0
    if total_invertido_proceso.get('inversionTotal__sum',0):
        totalinvertidoproceso = total_invertido_proceso.get('inversionTotal__sum',0)

    totalinvertidoproyectadas=0
    if total_invertido_proyectadas.get('inversionTotal__sum',0):
        totalinvertidoproyectadas = total_invertido_proyectadas.get('inversionTotal__sum',0)

    #total_invertido_2013.get('inversionTotal__sum',0) + total_invertido_2014.get(
    #'inversionTotal__sum',0) + total_invertido_proceso.get('inversionTotal__sum',0) + total_invertido_proyectadas.get(
    #'inversionTotal__sum',0)


    prs.slides[0].shapes[15].text= '{0:,}'.format(total_obras_concluidas)
    prs.slides[0].shapes[16].text= '$ {0:,.2f}'.format(totalinvertidoconcluidas)
    prs.slides[0].shapes[17].text= '{0:,}'.format(total_obras_proceso)
    prs.slides[0].shapes[18].text= '$ {0:,.2f}'.format(totalinvertidoproceso)
    prs.slides[0].shapes[19].text= '{0:,}'.format(total_obras_proyectadas)
    prs.slides[0].shapes[20].text= '$ {0:,.2f}'.format(totalinvertidoproyectadas)
    prs.slides[0].shapes[21].text= '{0:,}'.format(total_obras)
    prs.slides[0].shapes[22].text= '$ {0:,.2f}'.format(total_invertido)


    prs.save('/home/obrasapf/djangoObras/obras/static/ppt/ppt-generados/balance_general_' + str(usuario.user.id) + '.pptx')

    the_file = '/home/obrasapf/djangoObras/obras/static/ppt/ppt-generados/balance_general_' + str(usuario.user.id) + '.pptx'
    filename = os.path.basename(the_file)
    chunk_size = 8192
    response = StreamingHttpResponse(FileWrapper(open(the_file,"rb"), chunk_size),
                           content_type=mimetypes.guess_type(the_file)[0])
    response['Content-Length'] = os.path.getsize(the_file)
    response['Content-Disposition'] = "attachment; filename=%s" % filename
    return response



@login_required()
def hiper_info_general_ppt(request):
    prs = Presentation('/home/obrasapf/djangoObras/obras/static/ppt/HIPERVINCULO_INFORMACION_GENERAL.pptx')
    usuario = request.user.usuario
    dependencias = usuario.dependencia.all()
    subdependencias = usuario.subdependencia.all()
    query = Q()
    if dependencias and dependencias.count() > 0:
        if usuario.rol == 'US':
            query = Q(subdependencia__in=get_subdependencias_as_list_flat(subdependencias))
        else:

            query = Q(dependencia__in=get_subdependencias_as_list_flat(dependencias)) | Q(subdependencia__in=get_subdependencias_as_list_flat(dependencias)
        )

    obras_concluidas = Obra.objects.filter(
        Q(tipoObra=3),query

    )

    obras_proceso = Obra.objects.filter(
        Q(tipoObra=2),query
    )

    obras_proyectadas = Obra.objects.filter(
        Q(tipoObra=1),query
    )

    total_obras_proyectadas=0
    totalinvertidoproyectadas=0
    total_obras_proceso =0
    totalinvertidoproceso=0
    total_obras_concluidas =0
    totalinvertidoconcluidas =0

    if obras_proyectadas:
        total_obras_proyectadas = obras_proyectadas.count()
        total_invertido_proyectadas = obras_proyectadas.aggregate(Sum('inversionTotal'))
        totalinvertidoproyectadas = total_invertido_proyectadas.get('inversionTotal__sum',0)
    if obras_proceso:
        total_obras_proceso = obras_proceso.count()
        total_invertido_proceso = obras_proceso.aggregate(Sum('inversionTotal'))
        totalinvertidoproceso = total_invertido_proceso.get('inversionTotal__sum',0)

    if obras_concluidas:
        total_obras_concluidas = obras_concluidas.count()
        total_invertido_concluidas = obras_concluidas.aggregate(Sum('inversionTotal'))
        totalinvertidoconcluidas = total_invertido_concluidas.get('inversionTotal__sum',0)

    #total_obras_proyectadas = obras_proyectadas.count()
    #total_obras_proceso = obras_proceso.count()
    #total_obras_concluidas = obras_concluidas.count()

    #total_invertido_proyectadas = obras_proyectadas.aggregate(Sum('inversionTotal'))
    #total_invertido_proceso = obras_proceso.aggregate(Sum('inversionTotal'))
    #total_invertido_concluidas = obras_concluidas.aggregate(Sum('inversionTotal'))

    prs.slides[0].shapes[3].text_frame.paragraphs[0].font.size = Pt(9)
    prs.slides[0].shapes[3].text= '{0:,}'.format(total_obras_concluidas)
    prs.slides[0].shapes[4].text_frame.paragraphs[0].font.size = Pt(9)
    prs.slides[0].shapes[4].text= '$ {0:,.2f}'.format(totalinvertidoconcluidas)
    prs.slides[1].shapes[3].text_frame.paragraphs[0].font.size = Pt(9)
    prs.slides[1].shapes[3].text= '{0:,}'.format(total_obras_proceso)
    prs.slides[1].shapes[4].text_frame.paragraphs[0].font.size = Pt(9)
    prs.slides[1].shapes[4].text= '$ {0:,.2f}'.format(totalinvertidoproceso)
    prs.slides[2].shapes[3].text_frame.paragraphs[0].font.size = Pt(9)
    prs.slides[2].shapes[3].text= '{0:,}'.format(total_obras_proyectadas)
    prs.slides[2].shapes[4].text_frame.paragraphs[0].font.size = Pt(9)
    prs.slides[2].shapes[4].text= '$ {0:,.2f}'.format(totalinvertidoproyectadas)



    prs.save('/home/obrasapf/djangoObras/obras/static/ppt/ppt-generados/hiper_info_general_' + str(usuario.user.id) + '.pptx')

    the_file = '/home/obrasapf/djangoObras/obras/static/ppt/ppt-generados/hiper_info_general_' + str(usuario.user.id) + '.pptx'
    filename = os.path.basename(the_file)
    chunk_size = 8192
    response = StreamingHttpResponse(FileWrapper(open(the_file,"rb"), chunk_size),
                           content_type=mimetypes.guess_type(the_file)[0])
    response['Content-Length'] = os.path.getsize(the_file)
    response['Content-Disposition'] = "attachment; filename=%s" % filename

    return response
    #abrir_pptx('hiper_info_general.pptx')
    #return render_to_response('admin/obras/consulta_predefinidos/consulta-predefinidos.html', {'clases': ''}, context_instance=RequestContext(request))


@login_required()
@user_passes_test(is_super_admin)
def hiper_por_sector_ppt(request):
    #prs = Presentation('obras/static/ppt/HIPERVINCULO_POR_SECTOR.pptx')
    prs = Presentation('/home/obrasapf/djangoObras/obras/static/ppt/HIPERVINCULO_POR_SECTOR.pptx')
    usuario = request.user.usuario

    start_date = datetime.date(2000, 01, 01)
    end_date = datetime.date(2016, 12, 31)
    dependencias = {}

    for dependencia in Dependencia.objects.filter(
        Q(obraoprograma='O'),Q(dependienteDe=None)
    ):
        print dependencia.nombreDependencia


        obras_concluidas = Obra.objects.filter(
            Q(fechaTermino__range=(start_date, end_date)),
            Q(tipoObra=3),
            Q(dependencia=dependencia),
        )

        obras_proceso = Obra.objects.filter(
            Q(fechaTermino__range=(start_date, end_date)),
            Q(tipoObra=2),
            Q(dependencia=dependencia),
        )

        obras_proyectadas = Obra.objects.filter(
            Q(fechaTermino__range=(start_date, end_date)),
            Q(tipoObra=1),
            Q(dependencia=dependencia),
        )


        total_obras_concluidas = obras_concluidas.count()
        total_obras_proceso = obras_proceso.count()
        total_obras_proyectadas = obras_proyectadas.count()


        total_invertido_concluidas = obras_concluidas.aggregate(Sum('inversionTotal'))
        total_invertido_proceso = obras_proceso.aggregate(Sum('inversionTotal'))
        total_invertido_proyectadas = obras_proyectadas.aggregate(Sum('inversionTotal'))



        if dependencia.nombreDependencia =='SEGOB': indiceSlide =0
        elif dependencia.nombreDependencia =='SEDESOL': indiceSlide =2
        elif dependencia.nombreDependencia =='SEMARNAT': indiceSlide = 4
        elif dependencia.nombreDependencia =='SAGARPA': indiceSlide = 6
        elif dependencia.nombreDependencia =='SCT': indiceSlide = 8
        elif dependencia.nombreDependencia =='SEP': indiceSlide = 10
        elif dependencia.nombreDependencia =='SS': indiceSlide = 12
        elif dependencia.nombreDependencia =='SEDATU': indiceSlide = 14
        elif dependencia.nombreDependencia =='CULTURA': indiceSlide = 16
        elif dependencia.nombreDependencia =='SECTUR': indiceSlide = 18
        elif dependencia.nombreDependencia =='PEMEX': indiceSlide = 20
        elif dependencia.nombreDependencia =='CFE': indiceSlide = 22
        elif dependencia.nombreDependencia =='IMSS': indiceSlide = 24
        elif dependencia.nombreDependencia =='ISSSTE': indiceSlide = 26
        elif dependencia.nombreDependencia =='CONAGUA': indiceSlide = 28
        else: indiceSlide =30

        totalinvertidoproceso=0
        totalinvertidoconcluidas=0
        totalinvertidoproyectadas=0

        if str(total_invertido_concluidas.get('inversionTotal__sum',0)) != 'None': totalinvertidoconcluidas=total_invertido_concluidas.get('inversionTotal__sum',0)
        if str(total_invertido_proyectadas.get('inversionTotal__sum',0)) != 'None': totalinvertidoproyectadas=total_invertido_proyectadas.get('inversionTotal__sum',0)
        if str(total_invertido_proceso.get('inversionTotal__sum',0)) != 'None': totalinvertidoproceso=total_invertido_proceso.get('inversionTotal__sum',0)

        TOTAL_INVERTIDO=totalinvertidoconcluidas+totalinvertidoproceso+totalinvertidoproyectadas
        TOTAL_OBRAS=total_obras_concluidas+total_obras_proceso+total_obras_proyectadas

        for x in range(5,13):
            prs.slides[indiceSlide].shapes[x].text_frame.paragraphs[0].font.size = Pt(14)

        prs.slides[indiceSlide].shapes[5].text= '{0:,}'.format(total_obras_concluidas)
        prs.slides[indiceSlide].shapes[6].text= '$ {0:,.2f}'.format(totalinvertidoconcluidas)
        prs.slides[indiceSlide].shapes[7].text= '{0:,}'.format(total_obras_proceso)
        prs.slides[indiceSlide].shapes[8].text= '$ {0:,.2f}'.format(totalinvertidoproceso)
        prs.slides[indiceSlide].shapes[11].text= '{0:,}'.format(total_obras_proyectadas)
        prs.slides[indiceSlide].shapes[12].text= '$ {0:,.2f}'.format(totalinvertidoproyectadas)
        prs.slides[indiceSlide].shapes[9].text= '{0:,}'.format(TOTAL_OBRAS)
        prs.slides[indiceSlide].shapes[10].text= '$ {0:,.2f}'.format(TOTAL_INVERTIDO)



    #prs.save('obras/static/ppt/ppt-generados/hiper_por_sector_' + str(usuario.user.id) + '.pptx')
    #the_file = 'obras/static/ppt/ppt-generados/hiper_por_sector_' + str(usuario.user.id) + '.pptx'

    prs.save('/home/obrasapf/djangoObras/obras/static/ppt/ppt-generados/hiper_por_sector_' + str(usuario.user.id) + '.pptx')
    the_file = '/home/obrasapf/djangoObras/obras/static/ppt/ppt-generados/hiper_por_sector_' + str(usuario.user.id) + '.pptx'

    filename = os.path.basename(the_file)
    chunk_size = 8192
    response = StreamingHttpResponse(FileWrapper(open(the_file,"rb"), chunk_size),
                           content_type=mimetypes.guess_type(the_file)[0])
    response['Content-Length'] = os.path.getsize(the_file)
    response['Content-Disposition'] = "attachment; filename=%s" % filename

    return response
@login_required()
@user_passes_test(is_super_admin)
def hiper_por_entidad_ppt(request):

    #prs = Presentation('/home/obrasapf/djangoObras/obras/static/ppt/HIPERVINCULO_POR_ENTIDAD.pptx')
    prs = Presentation('obras/static/ppt/HIPERVINCULO_POR_ENTIDAD.pptx')

    usuario = request.user.usuario
    start_date_2012 = datetime.date(2012, 01, 01)
    end_date_2012 = datetime.date(2012, 12, 31)

    start_date_2013 = datetime.date(2013, 01, 01)
    end_date_2013 = datetime.date(2013, 12, 31)

    start_date_2014 = datetime.date(2014, 01, 01)
    end_date_2014 = datetime.date(2014, 12, 31)

    start_date_2015 = datetime.date(2015, 01, 01)
    end_date_2015 = datetime.date(2015, 12, 31)

    start_date_2016 = datetime.date(2016, 01, 01)
    end_date_2016 = datetime.date(2016, 12, 31)

    start_date_2017 = datetime.date(2017, 01, 01)
    end_date_2017 = datetime.date(2017, 12, 31)

    start_date_2018 = datetime.date(2018, 01, 01)
    end_date_2018 = datetime.date(2018, 12, 31)
    indiceSlide=0
    estados = {}
    listaEstados = Estado.objects.exclude(nombreEstado='INTERESTATAL').exclude(nombreEstado='NACIONAL').order_by('nombreEstado')
    listaEstados = listaEstados.order_by('nombreEstado')

    obras = Obra.objects.exclude(tipoObra__id=4)
    obras = obras.values('dependencia__nombreDependencia').annotate(totalinvertido=Sum('inversionTotal'))
    obrasMinimo = obras.order_by('totalinvertido')[:3]
    obrasMaximo = obras.order_by('totalinvertido').reverse()[:3]
    minimoDependencia1=obrasMinimo[2]['dependencia__nombreDependencia']
    minimoDependencia2=obrasMinimo[1]['dependencia__nombreDependencia']
    minimoDependencia3=obrasMinimo[0]['dependencia__nombreDependencia']
    minimoInvertido1=obrasMinimo[2]['totalinvertido']
    minimoInvertido2=obrasMinimo[1]['totalinvertido']
    minimoInvertido3=obrasMinimo[0]['totalinvertido']

    maximoDependencia1=obrasMaximo[0]['dependencia__nombreDependencia']
    maximoDependencia2=obrasMaximo[1]['dependencia__nombreDependencia']
    maximoDependencia3=obrasMaximo[2]['dependencia__nombreDependencia']
    maximoInvertido1=obrasMaximo[0]['totalinvertido']
    maximoInvertido2=obrasMaximo[1]['totalinvertido']
    maximoInvertido3=obrasMaximo[2]['totalinvertido']

    for obra in obras:
        print obra['totalinvertido']
        print obra['dependencia__nombreDependencia']


    for estado in listaEstados:
        print estado.nombreEstado
        obras_2012_concluidas = Obra.objects.filter(
            Q(fechaTermino__range=(start_date_2012, end_date_2012)),
            Q(tipoObra=3),
            Q(estado=estado),
        )
        obras_2013_concluidas = Obra.objects.filter(
            Q(fechaTermino__range=(start_date_2013, end_date_2013)),
            Q(tipoObra=3),
            Q(estado=estado),
        )
        obras_2014_concluidas = Obra.objects.filter(
            Q(fechaTermino__range=(start_date_2014, end_date_2014)),
            Q(tipoObra=3),
            Q(estado=estado),
        )
        obras_2015_concluidas = Obra.objects.filter(
            Q(fechaTermino__range=(start_date_2015, end_date_2015)),
            Q(tipoObra=3),
            Q(estado=estado),
        )
        obras_2016_concluidas = Obra.objects.filter(
            Q(fechaTermino__range=(start_date_2016, end_date_2016)),
            Q(tipoObra=3),
            Q(estado=estado),
        )

        obras_proceso = Obra.objects.filter(
            Q(tipoObra=2),
            Q(estado=estado),
        )

        obras_2014_proyectadas = Obra.objects.filter(
            Q(fechaTermino__range=(start_date_2014, end_date_2014)),
            Q(tipoObra=1),
            Q(estado=estado),
        )

        obras_2015_proyectadas = Obra.objects.filter(
            Q(fechaTermino__range=(start_date_2015, end_date_2015)),
            Q(tipoObra=1),
            Q(estado=estado),
        )

        obras_2016_proyectadas = Obra.objects.filter(
            Q(fechaTermino__range=(start_date_2016, end_date_2016)),
            Q(tipoObra=1),
            Q(estado=estado),
        )

        obras_2017_proyectadas = Obra.objects.filter(
            Q(fechaTermino__range=(start_date_2017, end_date_2017)),
            Q(tipoObra=1),
            Q(estado=estado),
        )

        obras_2018_proyectadas = Obra.objects.filter(
            Q(fechaTermino__range=(start_date_2018, end_date_2018)),
            Q(tipoObra=1),
            Q(estado=estado),
        )

        total_obras_concluidas_2012 = obras_2012_concluidas.count()
        total_obras_concluidas_2013 = obras_2013_concluidas.count()
        total_obras_concluidas_2014 = obras_2014_concluidas.count()
        total_obras_concluidas_2015 = obras_2015_concluidas.count()
        total_obras_concluidas_2016 = obras_2016_concluidas.count()
        total_obras_proceso = obras_proceso.count()
        total_obras_proyectadas_2014 = obras_2014_proyectadas.count()
        total_obras_proyectadas_2015 = obras_2015_proyectadas.count()
        total_obras_proyectadas_2016 = obras_2016_proyectadas.count()
        total_obras_proyectadas_2017 = obras_2017_proyectadas.count()
        total_obras_proyectadas_2018 = obras_2018_proyectadas.count()

        total_invertido_2012_concluidas = obras_2012_concluidas.aggregate(Sum('inversionTotal'))
        total_invertido_2013_concluidas = obras_2013_concluidas.aggregate(Sum('inversionTotal'))
        total_invertido_2014_concluidas = obras_2014_concluidas.aggregate(Sum('inversionTotal'))
        total_invertido_2015_concluidas = obras_2015_concluidas.aggregate(Sum('inversionTotal'))
        total_invertido_2016_concluidas = obras_2016_concluidas.aggregate(Sum('inversionTotal'))

        total_invertido_proceso = obras_proceso.aggregate(Sum('inversionTotal'))

        total_invertido_proyectadas_2014 = obras_2014_proyectadas.aggregate(Sum('inversionTotal'))
        total_invertido_proyectadas_2015 = obras_2015_proyectadas.aggregate(Sum('inversionTotal'))
        total_invertido_proyectadas_2016 = obras_2016_proyectadas.aggregate(Sum('inversionTotal'))
        total_invertido_proyectadas_2017 = obras_2017_proyectadas.aggregate(Sum('inversionTotal'))
        total_invertido_proyectadas_2018 = obras_2018_proyectadas.aggregate(Sum('inversionTotal'))

        totalinvertido2012=0
        totalinvertido2013=0
        totalinvertido2014=0
        totalinvertido2015=0
        totalinvertido2016=0
        totalinvertidoproceso=0
        totalinvertidoproyectadas2014=0
        totalinvertidoproyectadas2015=0
        totalinvertidoproyectadas2016=0
        totalinvertidoproyectadas2017=0
        totalinvertidoproyectadas2018=0

        if str(total_invertido_2012_concluidas.get('inversionTotal__sum',0)) != 'None': totalinvertido2012=total_invertido_2012_concluidas.get('inversionTotal__sum',0)
        if str(total_invertido_2013_concluidas.get('inversionTotal__sum',0)) != 'None': totalinvertido2013=total_invertido_2013_concluidas.get('inversionTotal__sum',0)
        if str(total_invertido_2014_concluidas.get('inversionTotal__sum',0)) != 'None': totalinvertido2014=total_invertido_2014_concluidas.get('inversionTotal__sum',0)
        if str(total_invertido_2015_concluidas.get('inversionTotal__sum',0)) != 'None': totalinvertido2015=total_invertido_2015_concluidas.get('inversionTotal__sum',0)
        if str(total_invertido_2016_concluidas.get('inversionTotal__sum',0)) != 'None': totalinvertido2016=total_invertido_2016_concluidas.get('inversionTotal__sum',0)
        if str(total_invertido_proceso.get('inversionTotal__sum',0)) != 'None': totalinvertidoproceso=total_invertido_proceso.get('inversionTotal__sum',0)
        if str(total_invertido_proyectadas_2014.get('inversionTotal__sum',0)) != 'None': totalinvertidoproyectadas2014=total_invertido_proyectadas_2014.get('inversionTotal__sum',0)
        if str(total_invertido_proyectadas_2015.get('inversionTotal__sum',0)) != 'None': totalinvertidoproyectadas2015=total_invertido_proyectadas_2015.get('inversionTotal__sum',0)
        if str(total_invertido_proyectadas_2016.get('inversionTotal__sum',0)) != 'None': totalinvertidoproyectadas2016=total_invertido_proyectadas_2016.get('inversionTotal__sum',0)
        if str(total_invertido_proyectadas_2017.get('inversionTotal__sum',0)) != 'None': totalinvertidoproyectadas2017=total_invertido_proyectadas_2017.get('inversionTotal__sum',0)
        if str(total_invertido_proyectadas_2018.get('inversionTotal__sum',0)) != 'None': totalinvertidoproyectadas2018=total_invertido_proyectadas_2018.get('inversionTotal__sum',0)

        totalObrasConcluidas= total_obras_concluidas_2012+total_obras_concluidas_2013+total_obras_concluidas_2014+total_obras_concluidas_2015+total_obras_concluidas_2016
        totalInvertidoConcluidas=totalinvertido2012+totalinvertido2013+totalinvertido2014+totalinvertido2015+totalinvertido2016


        totalObras15_18 = total_obras_proyectadas_2016+total_obras_proyectadas_2017+total_obras_proyectadas_2018
        totalInvertido15_18 = totalinvertidoproyectadas2016+totalinvertidoproyectadas2017+totalinvertidoproyectadas2018

        totalObrasGeneral = totalObras15_18+total_obras_concluidas_2012+total_obras_concluidas_2013+total_obras_concluidas_2014+total_obras_concluidas_2015+total_obras_concluidas_2016+total_obras_proceso
        totalInvertidoGeneral = totalInvertido15_18+totalinvertido2012+totalinvertido2013+totalinvertido2014+totalinvertido2015+totalinvertido2016+totalinvertidoproceso

        for x in range(12,36):
            prs.slides[indiceSlide].shapes[x].text_frame.paragraphs[0].font.size = Pt(8)

        #concluidas
        prs.slides[indiceSlide].shapes[12].text= '{0:,}'.format(total_obras_concluidas_2012)
        prs.slides[indiceSlide].shapes[13].text= '{0:,.2f}'.format(totalinvertido2012)
        prs.slides[indiceSlide].shapes[14].text= '{0:,}'.format(total_obras_concluidas_2013)
        prs.slides[indiceSlide].shapes[15].text= '{0:,.2f}'.format(totalinvertido2013)
        prs.slides[indiceSlide].shapes[16].text= '{0:,}'.format(total_obras_concluidas_2014)
        prs.slides[indiceSlide].shapes[17].text= '{0:,.2f}'.format(totalinvertido2014)
        prs.slides[indiceSlide].shapes[18].text= '{0:,}'.format(total_obras_concluidas_2015)
        prs.slides[indiceSlide].shapes[19].text= '{0:,.2f}'.format(totalinvertido2015)
        prs.slides[indiceSlide].shapes[20].text= '{0:,}'.format(total_obras_concluidas_2016)
        prs.slides[indiceSlide].shapes[21].text= '{0:,.2f}'.format(totalinvertido2016)
        prs.slides[indiceSlide].shapes[22].text= '{0:,}'.format(totalObrasConcluidas)
        prs.slides[indiceSlide].shapes[23].text= '{0:,.2f}'.format(totalInvertidoConcluidas)

        #proceso
        prs.slides[indiceSlide].shapes[24].text= '{0:,}'.format(total_obras_proceso)
        prs.slides[indiceSlide].shapes[25].text= '{0:,.2f}'.format(totalinvertidoproceso)

        #proyectadas
        prs.slides[indiceSlide].shapes[26].text= '{0:,}'.format(total_obras_proyectadas_2016)
        prs.slides[indiceSlide].shapes[27].text= '{0:,.2f}'.format(totalinvertidoproyectadas2016)
        prs.slides[indiceSlide].shapes[28].text= '{0:,}'.format(total_obras_proyectadas_2017)
        prs.slides[indiceSlide].shapes[29].text= '{0:,.2f}'.format(totalinvertidoproyectadas2017)
        prs.slides[indiceSlide].shapes[30].text= '{0:,}'.format(total_obras_proyectadas_2018)
        prs.slides[indiceSlide].shapes[31].text= '{0:,.2f}'.format(totalinvertidoproyectadas2018)
        prs.slides[indiceSlide].shapes[32].text= '{0:,}'.format(totalObras15_18)
        prs.slides[indiceSlide].shapes[33].text= '{0:,.2f}'.format(totalInvertido15_18)

        #total general
        prs.slides[indiceSlide].shapes[34].text= '{0:,}'.format(totalObrasGeneral)
        prs.slides[indiceSlide].shapes[35].text= '{0:,.2f}'.format(totalInvertidoGeneral)

        for x in range(36,47):
            prs.slides[indiceSlide].shapes[x].text_frame.paragraphs[0].font.size = Pt(7)
        #minimos y maximos
        prs.slides[indiceSlide].shapes[36].text= maximoDependencia1
        prs.slides[indiceSlide].shapes[37].text= '{0:,.2f}'.format(maximoInvertido1)
        prs.slides[indiceSlide].shapes[38].text= maximoDependencia2
        prs.slides[indiceSlide].shapes[39].text= '{0:,.2f}'.format(maximoInvertido2)
        prs.slides[indiceSlide].shapes[40].text= maximoDependencia3
        prs.slides[indiceSlide].shapes[41].text= '{0:,.2f}'.format(maximoInvertido3)

        prs.slides[indiceSlide].shapes[42].text= minimoDependencia1
        prs.slides[indiceSlide].shapes[43].text= '{0:,.2f}'.format(minimoInvertido1)
        prs.slides[indiceSlide].shapes[44].text= minimoDependencia2
        prs.slides[indiceSlide].shapes[45].text= '{0:,.2f}'.format(minimoInvertido2)
        prs.slides[indiceSlide].shapes[46].text= minimoDependencia3
        prs.slides[indiceSlide].shapes[47].text= '{0:,.2f}'.format(minimoInvertido3)




        indiceSlide=indiceSlide+1


    #prs.save('/home/obrasapf/djangoObras/obras/static/ppt/ppt-generados/hiper_por_entidad_' + str(usuario.user.id) + '.pptx')
    #the_file = '/home/obrasapf/djangoObras/obras/static/ppt/ppt-generados/hiper_por_entidad_' + str(usuario.user.id) + '.pptx'

    prs.save('obras/static/ppt/ppt-generados/hiper_por_entidad_' + str(usuario.user.id) + '.pptx')
    the_file = 'obras/static/ppt/ppt-generados/hiper_por_entidad_' + str(usuario.user.id) + '.pptx'

    filename = os.path.basename(the_file)
    chunk_size = 8192
    response = StreamingHttpResponse(FileWrapper(open(the_file,"rb"), chunk_size),
                           content_type=mimetypes.guess_type(the_file)[0])
    response['Content-Length'] = os.path.getsize(the_file)
    response['Content-Disposition'] = "attachment; filename=%s" % filename

    return response


@login_required()
def hiper_concluidas_ppt(request):
    #prs = Presentation('obras/static/ppt/HIPERVINCULO_CONCLUIDAS.pptx')
    prs = Presentation('/home/obrasapf/djangoObras/obras/static/ppt/HIPERVINCULO_CONCLUIDAS.pptx')
    usuario = request.user.usuario
    dependencias = usuario.dependencia.all()
    subdependencias = usuario.subdependencia.all()
    query = Q()
    if dependencias and dependencias.count() > 0:
        if usuario.rol == 'US':
            query = Q(subdependencia__in=get_subdependencias_as_list_flat(subdependencias))
        else:

            query = Q(dependencia__in=get_subdependencias_as_list_flat(dependencias)) | Q(subdependencia__in=get_subdependencias_as_list_flat(dependencias)
        )

    start_date_2012 = datetime.date(2012, 01, 01)
    end_date_2012 = datetime.date(2012, 12, 31)

    start_date_2013 = datetime.date(2013, 01, 01)
    end_date_2013 = datetime.date(2013, 12, 31)

    start_date_2014 = datetime.date(2014, 01, 01)
    end_date_2014 = datetime.date(2014, 12, 31)

    start_date_2015 = datetime.date(2015, 01, 01)
    end_date_2015 = datetime.date(2015, 12, 31)

    start_date_2016 = datetime.date(2016, 01, 01)
    end_date_2016 = datetime.date(2016, 12, 31)

    obras_2012_concluidas = Obra.objects.filter(
        Q(fechaTermino__range=(start_date_2012, end_date_2012)),
        Q(tipoObra=3),query,
    )
    obras_2013_concluidas = Obra.objects.filter(
        Q(fechaTermino__range=(start_date_2013, end_date_2013)),
        Q(tipoObra=3),query,
    )
    obras_2014_concluidas = Obra.objects.filter(
        Q(fechaTermino__range=(start_date_2014, end_date_2014)),
        Q(tipoObra=3),query,
    )
    obras_2015_concluidas = Obra.objects.filter(
        Q(fechaTermino__range=(start_date_2015, end_date_2015)),
        Q(tipoObra=3),query,
    )
    obras_2016_concluidas = Obra.objects.filter(
        Q(fechaTermino__range=(start_date_2016, end_date_2016)),
        Q(tipoObra=3),query,
    )

    obras_proceso = Obra.objects.filter(
        Q(tipoObra=2),query,
    )

    obras_proyectadas = Obra.objects.filter(
        Q(tipoObra=1),query,
    )

    total_obras_concluidas_2012 = obras_2012_concluidas.count()
    total_obras_concluidas_2013 = obras_2013_concluidas.count()
    total_obras_concluidas_2014 = obras_2014_concluidas.count()
    total_obras_concluidas_2015 = obras_2015_concluidas.count()
    total_obras_concluidas_2016 = obras_2016_concluidas.count()
    total_obras_proceso = obras_proceso.count()
    total_obras_proyectadas = obras_proyectadas.count()


    total_invertido_2012_concluidas = obras_2012_concluidas.aggregate(Sum('inversionTotal'))
    total_invertido_2013_concluidas = obras_2013_concluidas.aggregate(Sum('inversionTotal'))
    total_invertido_2014_concluidas = obras_2014_concluidas.aggregate(Sum('inversionTotal'))
    total_invertido_2015_concluidas = obras_2015_concluidas.aggregate(Sum('inversionTotal'))
    total_invertido_2016_concluidas = obras_2016_concluidas.aggregate(Sum('inversionTotal'))

    total_invertido_proceso = obras_proceso.aggregate(Sum('inversionTotal'))
    total_invertido_proyectadas = obras_proyectadas.aggregate(Sum('inversionTotal'))

    totalinvertido2012=0
    totalinvertido2013=0
    totalinvertido2014=0
    totalinvertido2015=0
    totalinvertido2016=0
    totalinvertidoproceso=0
    totalinvertidoproyectadas=0

    if str(total_invertido_2012_concluidas.get('inversionTotal__sum',0)) != 'None': totalinvertido2012=total_invertido_2012_concluidas.get('inversionTotal__sum',0)
    if str(total_invertido_2013_concluidas.get('inversionTotal__sum',0)) != 'None': totalinvertido2013=total_invertido_2013_concluidas.get('inversionTotal__sum',0)
    if str(total_invertido_2014_concluidas.get('inversionTotal__sum',0)) != 'None': totalinvertido2014=total_invertido_2014_concluidas.get('inversionTotal__sum',0)
    if str(total_invertido_2015_concluidas.get('inversionTotal__sum',0)) != 'None': totalinvertido2015=total_invertido_2015_concluidas.get('inversionTotal__sum',0)
    if str(total_invertido_2016_concluidas.get('inversionTotal__sum',0)) != 'None': totalinvertido2016=total_invertido_2016_concluidas.get('inversionTotal__sum',0)
    if str(total_invertido_proceso.get('inversionTotal__sum',0)) != 'None': totalinvertidoproceso=total_invertido_proceso.get('inversionTotal__sum',0)
    if str(total_invertido_proyectadas.get('inversionTotal__sum',0)) != 'None': totalinvertidoproyectadas=total_invertido_proyectadas.get('inversionTotal__sum',0)

    totalObrasConcluidas= total_obras_concluidas_2012+total_obras_concluidas_2013+total_obras_concluidas_2014+total_obras_concluidas_2015+total_obras_concluidas_2016
    totalInvertidoConcluidas=totalinvertido2012+totalinvertido2013+totalinvertido2014+totalinvertido2015+totalinvertido2016

    for x in range(2,12):
        prs.slides[0].shapes[x].text_frame.paragraphs[0].font.size = Pt(14)

    for x in range(2,4):
        prs.slides[1].shapes[x].text_frame.paragraphs[0].font.size = Pt(14)

    for x in range(2,4):
        prs.slides[2].shapes[x].text_frame.paragraphs[0].font.size = Pt(14)

    #concluidas
    prs.slides[0].shapes[2].text= '{0:,}'.format(totalObrasConcluidas)
    prs.slides[0].shapes[3].text= '{0:,.2f}'.format(totalInvertidoConcluidas)
    prs.slides[0].shapes[4].text= '{0:,}'.format(total_obras_concluidas_2012)
    prs.slides[0].shapes[5].text= '{0:,.2f}'.format(totalinvertido2012)
    prs.slides[0].shapes[6].text= '{0:,}'.format(total_obras_concluidas_2013)
    prs.slides[0].shapes[7].text= '{0:,.2f}'.format(totalinvertido2013)
    prs.slides[0].shapes[8].text= '{0:,}'.format(total_obras_concluidas_2014)
    prs.slides[0].shapes[9].text= '{0:,.2f}'.format(totalinvertido2014)
    prs.slides[0].shapes[10].text= '{0:,}'.format(total_obras_concluidas_2015)
    prs.slides[0].shapes[11].text= '{0:,.2f}'.format(totalinvertido2015)
    prs.slides[0].shapes[12].text= '{0:,}'.format(total_obras_concluidas_2016)
    prs.slides[0].shapes[13].text= '{0:,.2f}'.format(totalinvertido2016)

    #proceso
    prs.slides[1].shapes[2].text= '{0:,}'.format(total_obras_proceso)
    prs.slides[1].shapes[3].text= '{0:,.2f}'.format(totalinvertidoproceso)

    #proyectadas
    prs.slides[2].shapes[2].text= '{0:,}'.format(total_obras_proyectadas)
    prs.slides[2].shapes[3].text= '{0:,.2f}'.format(totalinvertidoproyectadas)

    prs.save('/home/obrasapf/djangoObras/obras/static/ppt/ppt-generados/hiper_concluidas_' + str(usuario.user.id) + '.pptx')
    the_file = '/home/obrasapf/djangoObras/obras/static/ppt/ppt-generados/hiper_concluidas_' + str(usuario.user.id) + '.pptx'

    #prs.save('obras/static/ppt/ppt-generados/hiper_concluidas_' + str(usuario.user.id) + '.pptx')
    #the_file = 'obras/static/ppt/ppt-generados/hiper_concluidas_' + str(usuario.user.id) + '.pptx'


    filename = os.path.basename(the_file)
    chunk_size = 8192
    response = StreamingHttpResponse(FileWrapper(open(the_file,"rb"), chunk_size),
                           content_type=mimetypes.guess_type(the_file)[0])
    response['Content-Length'] = os.path.getsize(the_file)
    response['Content-Disposition'] = "attachment; filename=%s" % filename

    return response

@login_required()
def hiper_inauguradas_ppt(request):
    prs = Presentation('/home/obrasapf/djangoObras/obras/static/ppt/HIPERVINCULO_INAUGURADAS.pptx')
    #prs = Presentation('obras/static/ppt/HIPERVINCULO_INAUGURADAS.pptx')
    usuario = request.user.usuario
    dependencias = usuario.dependencia.all()
    subdependencias = usuario.subdependencia.all()
    query = Q()
    if dependencias and dependencias.count() > 0:
        if usuario.rol == 'US':
            query = Q(subdependencia__in=get_subdependencias_as_list_flat(subdependencias))
        else:

            query = Q(dependencia__in=get_subdependencias_as_list_flat(dependencias)) | Q(subdependencia__in=get_subdependencias_as_list_flat(dependencias)
        )

    start_date_2012 = datetime.date(2012, 01, 01)
    end_date_2012 = datetime.date(2012, 12, 31)

    start_date_2013 = datetime.date(2013, 01, 01)
    end_date_2013 = datetime.date(2013, 12, 31)

    start_date_2014 = datetime.date(2014, 01, 01)
    end_date_2014 = datetime.date(2014, 12, 31)

    start_date_2015 = datetime.date(2015, 01, 01)
    end_date_2015 = datetime.date(2015, 12, 31)

    start_date_2016 = datetime.date(2016, 01, 01)
    end_date_2016 = datetime.date(2016, 12, 31)

    #inauguradas presidente
    obras_2012_inauguradas_P = Obra.objects.filter(
        Q(fechaTermino__range=(start_date_2012, end_date_2012)),
        Q(inaugurador_id=1),query,
    )
    obras_2013_inauguradas_P = Obra.objects.filter(
        Q(fechaTermino__range=(start_date_2013, end_date_2013)),
        Q(inaugurador_id=1),query,
    )
    obras_2014_inauguradas_P = Obra.objects.filter(
        Q(fechaTermino__range=(start_date_2014, end_date_2014)),
        Q(inaugurador_id=1),query,
    )
    obras_2015_inauguradas_P = Obra.objects.filter(
        Q(fechaTermino__range=(start_date_2015, end_date_2015)),
        Q(inaugurador_id=1),query,
    )
    obras_2016_inauguradas_P = Obra.objects.filter(
        Q(fechaTermino__range=(start_date_2016, end_date_2016)),
        Q(inaugurador_id=1),query,
    )
    #inauguradas otros
    obras_2012_inauguradas_O = Obra.objects.filter(
        Q(fechaTermino__range=(start_date_2012, end_date_2012)),
        Q(inaugurador_id=7),query,
    )
    obras_2013_inauguradas_O = Obra.objects.filter(
        Q(fechaTermino__range=(start_date_2013, end_date_2013)),
        Q(inaugurador_id=7),query,
    )
    obras_2014_inauguradas_O = Obra.objects.filter(
        Q(fechaTermino__range=(start_date_2014, end_date_2014)),
        Q(inaugurador_id=7),query,
    )
    obras_2015_inauguradas_O = Obra.objects.filter(
        Q(fechaTermino__range=(start_date_2015, end_date_2015)),
        Q(inaugurador_id=7),query,
    )
    obras_2016_inauguradas_O = Obra.objects.filter(
        Q(fechaTermino__range=(start_date_2016, end_date_2016)),
        Q(inaugurador_id=7),query,
    )

    #SENIALIZADAS concluidas
    obras_2012_senializadas = Obra.objects.filter(
        Q(fechaTermino__range=(start_date_2012, end_date_2012)),
        Q(tipoObra=3),Q(senalizacion=1),query,
    )
    obras_2013_senializadas = Obra.objects.filter(
        Q(fechaTermino__range=(start_date_2013, end_date_2013)),
        Q(tipoObra=3),Q(senalizacion=1),query,
    )
    obras_2014_senializadas = Obra.objects.filter(
        Q(fechaTermino__range=(start_date_2014, end_date_2014)),
        Q(tipoObra=3),Q(senalizacion=1),query,
    )
    obras_2015_senializadas = Obra.objects.filter(
        Q(fechaTermino__range=(start_date_2015, end_date_2015)),
        Q(tipoObra=3),Q(senalizacion=1),query,
    )
    obras_2016_senializadas = Obra.objects.filter(
        Q(fechaTermino__range=(start_date_2016, end_date_2016)),
        Q(tipoObra=3),Q(senalizacion=1),query,
    )


    obras_proceso = Obra.objects.filter(
        Q(tipoObra=2),Q(senalizacion=1),query,
    )

    obras_proyectadas = Obra.objects.filter(
        Q(tipoObra=1),Q(senalizacion=1),query,
    )

    total_obras_inauguradas_P_2012 = obras_2012_inauguradas_P.count()
    total_obras_inauguradas_P_2013 = obras_2013_inauguradas_P.count()
    total_obras_inauguradas_P_2014 = obras_2014_inauguradas_P.count()
    total_obras_inauguradas_P_2015 = obras_2015_inauguradas_P.count()
    total_obras_inauguradas_P_2016 = obras_2016_inauguradas_P.count()

    total_obras_inauguradas_O_2012 = obras_2012_inauguradas_O.count()
    total_obras_inauguradas_O_2013 = obras_2013_inauguradas_O.count()
    total_obras_inauguradas_O_2014 = obras_2014_inauguradas_O.count()
    total_obras_inauguradas_O_2015 = obras_2015_inauguradas_O.count()
    total_obras_inauguradas_O_2016 = obras_2016_inauguradas_O.count()

    total_obras_inauguradas_P= total_obras_inauguradas_P_2012+total_obras_inauguradas_P_2013+total_obras_inauguradas_P_2014+total_obras_inauguradas_P_2015+total_obras_inauguradas_P_2016
    total_obras_inauguradas_O= total_obras_inauguradas_O_2012+total_obras_inauguradas_O_2013+total_obras_inauguradas_O_2014+total_obras_inauguradas_O_2015+total_obras_inauguradas_O_2016

    total_obras_senializadas_2012 = obras_2012_senializadas.count()
    total_obras_senializadas_2013 = obras_2013_senializadas.count()
    total_obras_senializadas_2014 = obras_2014_senializadas.count()
    total_obras_senializadas_2015 = obras_2015_senializadas.count()
    total_obras_senializadas_2016 = obras_2016_senializadas.count()

    total_obras_proceso = obras_proceso.count()
    total_obras_proyectadas = obras_proyectadas.count()

    total_obras_senializadas =  total_obras_senializadas_2012+total_obras_senializadas_2013+total_obras_senializadas_2014+total_obras_senializadas_2015+total_obras_senializadas_2016+total_obras_proceso+total_obras_proyectadas


    for x in range(8,25):
        prs.slides[0].shapes[x].text_frame.paragraphs[0].font.size = Pt(10)

    #inauguradas
    prs.slides[0].shapes[8].text= '{0:,}'.format(total_obras_inauguradas_P_2012)
    prs.slides[0].shapes[9].text= '{0:,}'.format(total_obras_inauguradas_O_2012)
    prs.slides[0].shapes[10].text= '{0:,}'.format(total_obras_inauguradas_P_2013)
    prs.slides[0].shapes[11].text= '{0:,}'.format(total_obras_inauguradas_O_2013)
    prs.slides[0].shapes[12].text= '{0:,}'.format(total_obras_inauguradas_P_2014)
    prs.slides[0].shapes[13].text= '{0:,}'.format(total_obras_inauguradas_O_2014)
    prs.slides[0].shapes[14].text= '{0:,}'.format(total_obras_inauguradas_P_2015)
    prs.slides[0].shapes[15].text= '{0:,}'.format(total_obras_inauguradas_O_2015)
    prs.slides[0].shapes[16].text= '{0:,}'.format(total_obras_inauguradas_P_2016)
    prs.slides[0].shapes[17].text= '{0:,}'.format(total_obras_inauguradas_O_2016)
    prs.slides[0].shapes[18].text= '{0:,}'.format(total_obras_inauguradas_P)
    prs.slides[0].shapes[19].text= '{0:,}'.format(total_obras_inauguradas_O)
    #SENIALIZADAS
    prs.slides[0].shapes[20].text= '{0:,}'.format(total_obras_senializadas_2012)
    prs.slides[0].shapes[21].text= '{0:,}'.format(total_obras_senializadas_2013)
    prs.slides[0].shapes[22].text= '{0:,}'.format(total_obras_senializadas_2014)
    prs.slides[0].shapes[23].text= '{0:,}'.format(total_obras_senializadas_2015)
    prs.slides[0].shapes[24].text= '{0:,}'.format(total_obras_senializadas_2016)
    prs.slides[0].shapes[25].text= '{0:,}'.format(total_obras_proceso)
    prs.slides[0].shapes[26].text= '{0:,}'.format(total_obras_proyectadas)
    prs.slides[0].shapes[27].text= '{0:,}'.format(total_obras_senializadas)


    prs.save('/home/obrasapf/djangoObras/obras/static/ppt/ppt-generados/hiper_inauguradas_senalizadas_' + str(usuario.user.id) + '.pptx')
    the_file = '/home/obrasapf/djangoObras/obras/static/ppt/ppt-generados/hiper_inauguradas_senalizadas_' + str(usuario.user.id) + '.pptx'

    #prs.save('obras/static/ppt/ppt-generados/hiper_inauguradas_senalizadas_' + str(usuario.user.id) + '.pptx')
    #the_file = 'obras/static/ppt/ppt-generados/hiper_inauguradas_senalizadas_' + str(usuario.user.id) + '.pptx'

    filename = os.path.basename(the_file)
    chunk_size = 8192
    response = StreamingHttpResponse(FileWrapper(open(the_file,"rb"), chunk_size),
                           content_type=mimetypes.guess_type(the_file)[0])
    response['Content-Length'] = os.path.getsize(the_file)
    response['Content-Disposition'] = "attachment; filename=%s" % filename

    return response

@login_required
def hiper_interestatal_ppt(request):
    prs = Presentation('/home/obrasapf/djangoObras/obras/static/ppt/HIPERVINCULO_INTERESTATAL.pptx')
    #prs = Presentation('obras/static/ppt/HIPERVINCULO_INTERESTATAL.pptx')
    usuario = request.user.usuario
    dependencias = usuario.dependencia.all()
    subdependencias = usuario.subdependencia.all()
    query = Q()
    if dependencias and dependencias.count() > 0:
        if usuario.rol == 'US':
            query = Q(subdependencia__in=get_subdependencias_as_list_flat(subdependencias))
        else:

            query = Q(dependencia__in=get_subdependencias_as_list_flat(dependencias)) | Q(subdependencia__in=get_subdependencias_as_list_flat(dependencias)
        )

    start_date_2012 = datetime.date(2012, 01, 01)
    end_date_2012 = datetime.date(2012, 12, 31)

    start_date_2013 = datetime.date(2013, 01, 01)
    end_date_2013 = datetime.date(2013, 12, 31)

    start_date_2014 = datetime.date(2014, 01, 01)
    end_date_2014 = datetime.date(2014, 12, 31)

    start_date_2015 = datetime.date(2015, 01, 01)
    end_date_2015 = datetime.date(2015, 12, 31)

    start_date_2016 = datetime.date(2016, 01, 01)
    end_date_2016 = datetime.date(2016, 12, 31)

    start_date_2017 = datetime.date(2017, 01, 01)
    end_date_2017 = datetime.date(2017, 12, 31)

    start_date_2018 = datetime.date(2018, 01, 01)
    end_date_2018 = datetime.date(2018, 12, 31)

    obras_2012_concluidas = Obra.objects.filter(
        Q(fechaTermino__range=(start_date_2012, end_date_2012)),
        Q(estado_id=34),Q(tipoObra=3),query,
    )
    obras_2013_concluidas = Obra.objects.filter(
        Q(fechaTermino__range=(start_date_2013, end_date_2013)),
        Q(estado_id=34),Q(tipoObra=3),query,
    )
    obras_2014_concluidas = Obra.objects.filter(
        Q(fechaTermino__range=(start_date_2014, end_date_2014)),
        Q(estado_id=34),Q(tipoObra=3),query,
    )
    obras_2015_concluidas = Obra.objects.filter(
        Q(fechaTermino__range=(start_date_2015, end_date_2015)),
        Q(estado_id=34),Q(tipoObra=3),query,
    )
    obras_2016_concluidas = Obra.objects.filter(
        Q(fechaTermino__range=(start_date_2016, end_date_2016)),
        Q(estado_id=34),Q(tipoObra=3),query,
    )

    obras_proceso = Obra.objects.filter(
        Q(estado_id=34),Q(tipoObra=2),query,
    )

    obras_2015_proyectadas = Obra.objects.filter(
        Q(fechaTermino__range=(start_date_2015, end_date_2015)),
        Q(estado_id=34),Q(tipoObra=1),query,
    )
    obras_2016_proyectadas = Obra.objects.filter(
        Q(fechaTermino__range=(start_date_2016, end_date_2016)),
        Q(estado_id=34),Q(tipoObra=1),query,
    )
    obras_2017_proyectadas = Obra.objects.filter(
        Q(fechaTermino__range=(start_date_2017, end_date_2017)),
        Q(estado_id=34),Q(tipoObra=1),query,
    )
    obras_2018_proyectadas = Obra.objects.filter(
        Q(fechaTermino__range=(start_date_2018, end_date_2018)),
        Q(estado_id=34),Q(tipoObra=1),query,
    )

    total_obras_concluidas_2012 = obras_2012_concluidas.count()
    total_obras_concluidas_2013 = obras_2013_concluidas.count()
    total_obras_concluidas_2014 = obras_2014_concluidas.count()
    total_obras_concluidas_2015 = obras_2015_concluidas.count()
    total_obras_concluidas_2016 = obras_2016_concluidas.count()
    total_obras_concluidas = total_obras_concluidas_2012+total_obras_concluidas_2013+total_obras_concluidas_2014+total_obras_concluidas_2015+total_obras_concluidas_2016
    total_obras_proceso = obras_proceso.count()
    total_obras_proyectadas_2015 = obras_2015_proyectadas.count()
    total_obras_proyectadas_2016 = obras_2016_proyectadas.count()
    total_obras_proyectadas_2017 = obras_2017_proyectadas.count()
    total_obras_proyectadas_2018 = obras_2018_proyectadas.count()
    total_obras_proyectadas = total_obras_proyectadas_2016+total_obras_proyectadas_2017+total_obras_proyectadas_2018
    total_obras = total_obras_concluidas+total_obras_proceso+total_obras_proyectadas

    total_invertido_2012_concluidas = obras_2012_concluidas.aggregate(Sum('inversionTotal'))
    total_invertido_2013_concluidas = obras_2013_concluidas.aggregate(Sum('inversionTotal'))
    total_invertido_2014_concluidas = obras_2014_concluidas.aggregate(Sum('inversionTotal'))
    total_invertido_2015_concluidas = obras_2015_concluidas.aggregate(Sum('inversionTotal'))
    total_invertido_2016_concluidas = obras_2016_concluidas.aggregate(Sum('inversionTotal'))

    total_invertido_proceso = obras_proceso.aggregate(Sum('inversionTotal'))

    total_invertido_2015_proyectadas = obras_2015_proyectadas.aggregate(Sum('inversionTotal'))
    total_invertido_2016_proyectadas = obras_2016_proyectadas.aggregate(Sum('inversionTotal'))
    total_invertido_2017_proyectadas = obras_2017_proyectadas.aggregate(Sum('inversionTotal'))
    total_invertido_2018_proyectadas = obras_2018_proyectadas.aggregate(Sum('inversionTotal'))

    totalinvertido2012=0
    totalinvertido2013=0
    totalinvertido2014=0
    totalinvertido2015=0
    totalinvertido2016=0
    totalinvertidoproceso=0
    totalproyectadas2015=0
    totalproyectadas2016=0
    totalproyectadas2017=0
    totalproyectadas2018=0

    if str(total_invertido_2012_concluidas.get('inversionTotal__sum',0)) != 'None': totalinvertido2012=total_invertido_2012_concluidas.get('inversionTotal__sum',0)
    if str(total_invertido_2013_concluidas.get('inversionTotal__sum',0)) != 'None': totalinvertido2013=total_invertido_2013_concluidas.get('inversionTotal__sum',0)
    if str(total_invertido_2014_concluidas.get('inversionTotal__sum',0)) != 'None': totalinvertido2014=total_invertido_2014_concluidas.get('inversionTotal__sum',0)
    if str(total_invertido_2015_concluidas.get('inversionTotal__sum',0)) != 'None': totalinvertido2015=total_invertido_2015_concluidas.get('inversionTotal__sum',0)
    if str(total_invertido_2016_concluidas.get('inversionTotal__sum',0)) != 'None': totalinvertido2016=total_invertido_2016_concluidas.get('inversionTotal__sum',0)
    if str(total_invertido_proceso.get('inversionTotal__sum',0)) != 'None': totalinvertidoproceso=total_invertido_proceso.get('inversionTotal__sum',0)
    if str(total_invertido_2015_proyectadas.get('inversionTotal__sum',0)) != 'None': totalproyectadas2015=total_invertido_2015_proyectadas.get('inversionTotal__sum',0)
    if str(total_invertido_2016_proyectadas.get('inversionTotal__sum',0)) != 'None': totalproyectadas2016=total_invertido_2016_proyectadas.get('inversionTotal__sum',0)
    if str(total_invertido_2017_proyectadas.get('inversionTotal__sum',0)) != 'None': totalproyectadas2017=total_invertido_2017_proyectadas.get('inversionTotal__sum',0)
    if str(total_invertido_2018_proyectadas.get('inversionTotal__sum',0)) != 'None': totalproyectadas2018=total_invertido_2018_proyectadas.get('inversionTotal__sum',0)

    totalInvertidoConcluidas=totalinvertido2012+totalinvertido2013+totalinvertido2014+totalinvertido2015+totalinvertido2016
    totalInvertidoProyectadas=totalproyectadas2016+totalproyectadas2017+totalproyectadas2018
    totalInvertido = totalInvertidoConcluidas+totalinvertidoproceso+totalInvertidoProyectadas

    for x in range(9,31):
        prs.slides[0].shapes[x].text_frame.paragraphs[0].font.size = Pt(8)

    for x in range(31,33):
        prs.slides[0].shapes[x].text_frame.paragraphs[0].font.size = Pt(14)


    #concluidas
    prs.slides[0].shapes[9].text= '{0:,}'.format(total_obras_concluidas_2012)
    prs.slides[0].shapes[10].text= '{0:,.2f}'.format(totalinvertido2012)
    prs.slides[0].shapes[11].text= '{0:,}'.format(total_obras_concluidas_2013)
    prs.slides[0].shapes[12].text= '{0:,.2f}'.format(totalinvertido2013)
    prs.slides[0].shapes[13].text= '{0:,}'.format(total_obras_concluidas_2014)
    prs.slides[0].shapes[14].text= '{0:,.2f}'.format(totalinvertido2014)
    prs.slides[0].shapes[15].text= '{0:,}'.format(total_obras_concluidas_2015)
    prs.slides[0].shapes[16].text= '{0:,.2f}'.format(totalinvertido2015)
    prs.slides[0].shapes[17].text= '{0:,}'.format(total_obras_concluidas_2016)
    prs.slides[0].shapes[18].text= '{0:,.2f}'.format(totalinvertido2016)
    prs.slides[0].shapes[19].text= '{0:,}'.format(total_obras_concluidas)
    prs.slides[0].shapes[20].text= '{0:,.2f}'.format(totalInvertidoConcluidas)
    #proceso
    prs.slides[0].shapes[21].text= '{0:,}'.format(total_obras_proceso)
    prs.slides[0].shapes[22].text= '{0:,.2f}'.format(totalinvertidoproceso)
    #proyectadas
    prs.slides[0].shapes[23].text= '{0:,}'.format(total_obras_proyectadas_2016)
    prs.slides[0].shapes[24].text= '{0:,.2f}'.format(totalproyectadas2016)
    prs.slides[0].shapes[25].text= '{0:,}'.format(total_obras_proyectadas_2017)
    prs.slides[0].shapes[26].text= '{0:,.2f}'.format(totalproyectadas2017)
    prs.slides[0].shapes[27].text= '{0:,}'.format(total_obras_proyectadas_2018)
    prs.slides[0].shapes[28].text= '{0:,.2f}'.format(totalproyectadas2018)

    prs.slides[0].shapes[29].text= '{0:,}'.format(total_obras_proyectadas)
    prs.slides[0].shapes[30].text= '{0:,.2f}'.format(totalInvertidoProyectadas)

    #totales
    prs.slides[0].shapes[31].text= '{0:,}'.format(total_obras)
    prs.slides[0].shapes[32].text= '{0:,.2f}'.format(totalInvertido)



    prs.save('/home/obrasapf/djangoObras/obras/static/ppt/ppt-generados/hiper_interestatal_' + str(usuario.user.id) + '.pptx')
    the_file = '/home/obrasapf/djangoObras/obras/static/ppt/ppt-generados/hiper_interestatal_' + str(usuario.user.id) + '.pptx'

    #prs.save('obras/static/ppt/ppt-generados/hiper_interestatal_' + str(usuario.user.id) + '.pptx')
    #the_file = 'obras/static/ppt/ppt-generados/hiper_interestatal_' + str(usuario.user.id) + '.pptx'

    filename = os.path.basename(the_file)
    chunk_size = 8192
    response = StreamingHttpResponse(FileWrapper(open(the_file,"rb"), chunk_size),
                           content_type=mimetypes.guess_type(the_file)[0])
    response['Content-Length'] = os.path.getsize(the_file)
    response['Content-Disposition'] = "attachment; filename=%s" % filename

    return response

@login_required()
def hiper_nacional_ppt(request):
    prs = Presentation('/home/obrasapf/djangoObras/obras/static/ppt/HIPERVINCULO_NACIONAL.pptx')
    #prs = Presentation('obras/static/ppt/HIPERVINCULO_NACIONAL.pptx')
    usuario = request.user.usuario
    dependencias = usuario.dependencia.all()
    subdependencias = usuario.subdependencia.all()
    query = Q()
    if dependencias and dependencias.count() > 0:
        if usuario.rol == 'US':
            query = Q(subdependencia__in=get_subdependencias_as_list_flat(subdependencias))
        else:

            query = Q(dependencia__in=get_subdependencias_as_list_flat(dependencias)) | Q(subdependencia__in=get_subdependencias_as_list_flat(dependencias)
        )

    start_date_2012 = datetime.date(2012, 01, 01)
    end_date_2012 = datetime.date(2012, 12, 31)

    start_date_2013 = datetime.date(2013, 01, 01)
    end_date_2013 = datetime.date(2013, 12, 31)

    start_date_2014 = datetime.date(2014, 01, 01)
    end_date_2014 = datetime.date(2014, 12, 31)

    start_date_2015 = datetime.date(2015, 01, 01)
    end_date_2015 = datetime.date(2015, 12, 31)

    start_date_2016 = datetime.date(2016, 01, 01)
    end_date_2016 = datetime.date(2016, 12, 31)

    start_date_2017 = datetime.date(2017, 01, 01)
    end_date_2017 = datetime.date(2017, 12, 31)

    start_date_2018 = datetime.date(2018, 01, 01)
    end_date_2018 = datetime.date(2018, 12, 31)

    obras_2012_concluidas = Obra.objects.filter(
        Q(fechaTermino__range=(start_date_2012, end_date_2012)),
        Q(estado_id=33),Q(tipoObra=3),query,
    )
    obras_2013_concluidas = Obra.objects.filter(
        Q(fechaTermino__range=(start_date_2013, end_date_2013)),
        Q(estado_id=33),Q(tipoObra=3),query,
    )
    obras_2014_concluidas = Obra.objects.filter(
        Q(fechaTermino__range=(start_date_2014, end_date_2014)),
        Q(estado_id=33),Q(tipoObra=3),query,
    )
    obras_2015_concluidas = Obra.objects.filter(
        Q(fechaTermino__range=(start_date_2015, end_date_2015)),
        Q(estado_id=33),Q(tipoObra=3),query,
    )
    obras_2016_concluidas = Obra.objects.filter(
        Q(fechaTermino__range=(start_date_2016, end_date_2016)),
        Q(estado_id=33),Q(tipoObra=3),query,
    )

    obras_proceso = Obra.objects.filter(
        Q(estado_id=33),Q(tipoObra=2),query,
    )

    obras_2015_proyectadas = Obra.objects.filter(
        Q(fechaTermino__range=(start_date_2015, end_date_2015)),
        Q(estado_id=33),Q(tipoObra=1),query,
    )
    obras_2016_proyectadas = Obra.objects.filter(
        Q(fechaTermino__range=(start_date_2016, end_date_2016)),
        Q(estado_id=33),Q(tipoObra=1),query,
    )
    obras_2017_proyectadas = Obra.objects.filter(
        Q(fechaTermino__range=(start_date_2017, end_date_2017)),
        Q(estado_id=33),Q(tipoObra=1),query,
    )
    obras_2018_proyectadas = Obra.objects.filter(
        Q(fechaTermino__range=(start_date_2018, end_date_2018)),
        Q(estado_id=33),Q(tipoObra=1),query,
    )

    total_obras_concluidas_2012 = obras_2012_concluidas.count()
    total_obras_concluidas_2013 = obras_2013_concluidas.count()
    total_obras_concluidas_2014 = obras_2014_concluidas.count()
    total_obras_concluidas_2015 = obras_2015_concluidas.count()
    total_obras_concluidas_2016 = obras_2016_concluidas.count()
    total_obras_concluidas = total_obras_concluidas_2012+total_obras_concluidas_2013+total_obras_concluidas_2014+total_obras_concluidas_2015+total_obras_concluidas_2016
    total_obras_proceso = obras_proceso.count()
    total_obras_proyectadas_2015 = obras_2015_proyectadas.count()
    total_obras_proyectadas_2016 = obras_2016_proyectadas.count()
    total_obras_proyectadas_2017 = obras_2017_proyectadas.count()
    total_obras_proyectadas_2018 = obras_2018_proyectadas.count()
    total_obras_proyectadas = total_obras_proyectadas_2016+total_obras_proyectadas_2017+total_obras_proyectadas_2018
    total_obras = total_obras_concluidas+total_obras_proceso+total_obras_proyectadas

    total_invertido_2012_concluidas = obras_2012_concluidas.aggregate(Sum('inversionTotal'))
    total_invertido_2013_concluidas = obras_2013_concluidas.aggregate(Sum('inversionTotal'))
    total_invertido_2014_concluidas = obras_2014_concluidas.aggregate(Sum('inversionTotal'))
    total_invertido_2015_concluidas = obras_2015_concluidas.aggregate(Sum('inversionTotal'))
    total_invertido_2016_concluidas = obras_2016_concluidas.aggregate(Sum('inversionTotal'))

    total_invertido_proceso = obras_proceso.aggregate(Sum('inversionTotal'))

    total_invertido_2015_proyectadas = obras_2015_proyectadas.aggregate(Sum('inversionTotal'))
    total_invertido_2016_proyectadas = obras_2016_proyectadas.aggregate(Sum('inversionTotal'))
    total_invertido_2017_proyectadas = obras_2017_proyectadas.aggregate(Sum('inversionTotal'))
    total_invertido_2018_proyectadas = obras_2018_proyectadas.aggregate(Sum('inversionTotal'))

    totalinvertido2012=0
    totalinvertido2013=0
    totalinvertido2014=0
    totalinvertido2015=0
    totalinvertido2016=0
    totalinvertidoproceso=0
    totalproyectadas2015=0
    totalproyectadas2016=0
    totalproyectadas2017=0
    totalproyectadas2018=0

    if str(total_invertido_2012_concluidas.get('inversionTotal__sum',0)) != 'None': totalinvertido2012=total_invertido_2012_concluidas.get('inversionTotal__sum',0)
    if str(total_invertido_2013_concluidas.get('inversionTotal__sum',0)) != 'None': totalinvertido2013=total_invertido_2013_concluidas.get('inversionTotal__sum',0)
    if str(total_invertido_2014_concluidas.get('inversionTotal__sum',0)) != 'None': totalinvertido2014=total_invertido_2014_concluidas.get('inversionTotal__sum',0)
    if str(total_invertido_2015_concluidas.get('inversionTotal__sum',0)) != 'None': totalinvertido2015=total_invertido_2015_concluidas.get('inversionTotal__sum',0)
    if str(total_invertido_2016_concluidas.get('inversionTotal__sum',0)) != 'None': totalinvertido2016=total_invertido_2016_concluidas.get('inversionTotal__sum',0)
    if str(total_invertido_proceso.get('inversionTotal__sum',0)) != 'None': totalinvertidoproceso=total_invertido_proceso.get('inversionTotal__sum',0)
    if str(total_invertido_2015_proyectadas.get('inversionTotal__sum',0)) != 'None': totalproyectadas2015=total_invertido_2015_proyectadas.get('inversionTotal__sum',0)
    if str(total_invertido_2016_proyectadas.get('inversionTotal__sum',0)) != 'None': totalproyectadas2016=total_invertido_2016_proyectadas.get('inversionTotal__sum',0)
    if str(total_invertido_2017_proyectadas.get('inversionTotal__sum',0)) != 'None': totalproyectadas2017=total_invertido_2017_proyectadas.get('inversionTotal__sum',0)
    if str(total_invertido_2018_proyectadas.get('inversionTotal__sum',0)) != 'None': totalproyectadas2018=total_invertido_2018_proyectadas.get('inversionTotal__sum',0)

    totalInvertidoConcluidas=totalinvertido2012+totalinvertido2013+totalinvertido2014+totalinvertido2015+totalinvertido2016
    totalInvertidoProyectadas=totalproyectadas2016+totalproyectadas2017+totalproyectadas2018
    totalInvertido = totalInvertidoConcluidas+totalinvertidoproceso+totalInvertidoProyectadas

    for x in range(9,31):
        prs.slides[0].shapes[x].text_frame.paragraphs[0].font.size = Pt(8)

    for x in range(31,33):
        prs.slides[0].shapes[x].text_frame.paragraphs[0].font.size = Pt(14)


    #concluidas
    prs.slides[0].shapes[9].text= '{0:,}'.format(total_obras_concluidas_2012)
    prs.slides[0].shapes[10].text= '{0:,.2f}'.format(totalinvertido2012)
    prs.slides[0].shapes[11].text= '{0:,}'.format(total_obras_concluidas_2013)
    prs.slides[0].shapes[12].text= '{0:,.2f}'.format(totalinvertido2013)
    prs.slides[0].shapes[13].text= '{0:,}'.format(total_obras_concluidas_2014)
    prs.slides[0].shapes[14].text= '{0:,.2f}'.format(totalinvertido2014)
    prs.slides[0].shapes[15].text= '{0:,}'.format(total_obras_concluidas_2015)
    prs.slides[0].shapes[16].text= '{0:,.2f}'.format(totalinvertido2015)
    prs.slides[0].shapes[17].text= '{0:,}'.format(total_obras_concluidas_2016)
    prs.slides[0].shapes[18].text= '{0:,.2f}'.format(totalinvertido2016)
    prs.slides[0].shapes[19].text= '{0:,}'.format(total_obras_concluidas)
    prs.slides[0].shapes[20].text= '{0:,.2f}'.format(totalInvertidoConcluidas)
    #proceso
    prs.slides[0].shapes[21].text= '{0:,}'.format(total_obras_proceso)
    prs.slides[0].shapes[22].text= '{0:,.2f}'.format(totalinvertidoproceso)
    #proyectadas

    prs.slides[0].shapes[23].text= '{0:,}'.format(total_obras_proyectadas_2016)
    prs.slides[0].shapes[24].text= '{0:,.2f}'.format(totalproyectadas2016)
    prs.slides[0].shapes[25].text= '{0:,}'.format(total_obras_proyectadas_2017)
    prs.slides[0].shapes[26].text= '{0:,.2f}'.format(totalproyectadas2017)
    prs.slides[0].shapes[27].text= '{0:,}'.format(total_obras_proyectadas_2018)
    prs.slides[0].shapes[28].text= '{0:,.2f}'.format(totalproyectadas2017)

    prs.slides[0].shapes[29].text= '{0:,}'.format(total_obras_proyectadas)
    prs.slides[0].shapes[30].text= '{0:,.2f}'.format(totalInvertidoProyectadas)

    #totales
    prs.slides[0].shapes[31].text= '{0:,}'.format(total_obras)
    prs.slides[0].shapes[32].text= '{0:,.2f}'.format(totalInvertido)



    prs.save('/home/obrasapf/djangoObras/obras/static/ppt/ppt-generados/hiper_nacional_' + str(usuario.user.id) + '.pptx')
    the_file = '/home/obrasapf/djangoObras/obras/static/ppt/ppt-generados/hiper_nacional_' + str(usuario.user.id) + '.pptx'

    prs.save('obras/static/ppt/ppt-generados/hiper_nacional_' + str(usuario.user.id) + '.pptx')
    the_file = 'obras/static/ppt/ppt-generados/hiper_nacional_' + str(usuario.user.id) + '.pptx'

    filename = os.path.basename(the_file)
    chunk_size = 8192
    response = StreamingHttpResponse(FileWrapper(open(the_file,"rb"), chunk_size),
                           content_type=mimetypes.guess_type(the_file)[0])
    response['Content-Length'] = os.path.getsize(the_file)
    response['Content-Disposition'] = "attachment; filename=%s" % filename

    return response

@login_required()
def hiper_rangos_ppt(request):
    prs = Presentation('/home/obrasapf/djangoObras/obras/static/ppt/HIPERVINCULO_RANGOS.pptx')
    usuario = request.user.usuario
    dependencias = usuario.dependencia.all()
    subdependencias = usuario.subdependencia.all()
    query = Q()
    if dependencias and dependencias.count() > 0:
        if usuario.rol == 'US':
            query = Q(subdependencia__in=get_subdependencias_as_list_flat(subdependencias))
        else:

            query = Q(dependencia__in=get_subdependencias_as_list_flat(dependencias)) | Q(subdependencia__in=get_subdependencias_as_list_flat(dependencias)
        )

    #concluidas
    obras_concluidas_1 = Obra.objects.filter(
        Q(tipoObra=3),Q(inversionTotal__range=(0, 9.99)),query,
    )
    obras_concluidas_2 = Obra.objects.filter(
        Q(tipoObra=3),Q(inversionTotal__range=(10, 49.99)),query,
    )
    obras_concluidas_3 = Obra.objects.filter(
        Q(tipoObra=3),Q(inversionTotal__range=(50, 99.99)),query,
    )
    obras_concluidas_4 = Obra.objects.filter(
        Q(tipoObra=3),Q(inversionTotal__range=(100, 199.99)),query,
    )
    obras_concluidas_5 = Obra.objects.filter(
        Q(tipoObra=3),Q(inversionTotal__range=(200, 400.99)),query,
    )
    obras_concluidas_6 = Obra.objects.filter(
        Q(tipoObra=3),Q(inversionTotal__range=(500, 10000000000.99)),query,
    )

    #proceso
    obras_proceso_1 = Obra.objects.filter(
        Q(tipoObra=2),Q(inversionTotal__range=(0, 9.99)),query,
    )
    obras_proceso_2 = Obra.objects.filter(
        Q(tipoObra=2),Q(inversionTotal__range=(10, 49.99)),query,
    )
    obras_proceso_3 = Obra.objects.filter(
        Q(tipoObra=2),Q(inversionTotal__range=(50, 99.99)),query,
    )
    obras_proceso_4 = Obra.objects.filter(
        Q(tipoObra=2),Q(inversionTotal__range=(100, 199.99)),query,
    )
    obras_proceso_5 = Obra.objects.filter(
        Q(tipoObra=2),Q(inversionTotal__range=(200, 400.99)),query,
    )
    obras_proceso_6 = Obra.objects.filter(
        Q(tipoObra=2),Q(inversionTotal__range=(500, 10000000000.99)),query,
    )
    #proyectadas
    obras_proyectadas_1 = Obra.objects.filter(
        Q(tipoObra=1),Q(inversionTotal__range=(0, 9.99)),query,
    )
    obras_proyectadas_2 = Obra.objects.filter(
        Q(tipoObra=1),Q(inversionTotal__range=(10, 49.99)),query,
    )
    obras_proyectadas_3 = Obra.objects.filter(
        Q(tipoObra=1),Q(inversionTotal__range=(50, 99.99)),query,
    )
    obras_proyectadas_4 = Obra.objects.filter(
        Q(tipoObra=1),Q(inversionTotal__range=(100, 199.99)),query,
    )
    obras_proyectadas_5 = Obra.objects.filter(
        Q(tipoObra=1),Q(inversionTotal__range=(200, 400.99)),query,
    )
    obras_proyectadas_6 = Obra.objects.filter(
        Q(tipoObra=1),Q(inversionTotal__range=(500, 10000000000.99)),query,
    )

    #concluidas
    total_obras_concluidas_1 = obras_concluidas_1.count()
    total_obras_concluidas_2 = obras_concluidas_2.count()
    total_obras_concluidas_3 = obras_concluidas_3.count()
    total_obras_concluidas_4 = obras_concluidas_4.count()
    total_obras_concluidas_5 = obras_concluidas_5.count()
    total_obras_concluidas_6 = obras_concluidas_6.count()
    total_concluidas = total_obras_concluidas_1+total_obras_concluidas_2+total_obras_concluidas_3+total_obras_concluidas_4+total_obras_concluidas_5+total_obras_concluidas_6
    print('division' + str(Decimal(total_obras_concluidas_1)/Decimal(total_concluidas)))
    ciento_concluidas_1 = (Decimal(total_obras_concluidas_1)/Decimal(total_concluidas))*100
    ciento_concluidas_2 = (Decimal(total_obras_concluidas_2)/Decimal(total_concluidas))*100
    ciento_concluidas_3 = (Decimal(total_obras_concluidas_3)/Decimal(total_concluidas))*100
    ciento_concluidas_4 = (Decimal(total_obras_concluidas_4)/Decimal(total_concluidas))*100
    ciento_concluidas_5 = (Decimal(total_obras_concluidas_5)/Decimal(total_concluidas))*100
    ciento_concluidas_6 = (Decimal(total_obras_concluidas_6)/Decimal(total_concluidas))*100
    total_ciento_concluidas = ciento_concluidas_1+ciento_concluidas_2+ciento_concluidas_3+ciento_concluidas_4+ciento_concluidas_5+ciento_concluidas_6
    #proceso
    total_obras_proceso_1 = obras_proceso_1.count()
    total_obras_proceso_2 = obras_proceso_2.count()
    total_obras_proceso_3 = obras_proceso_3.count()
    total_obras_proceso_4 = obras_proceso_4.count()
    total_obras_proceso_5 = obras_proceso_5.count()
    total_obras_proceso_6 = obras_proceso_6.count()
    total_proceso = total_obras_proceso_1+total_obras_proceso_2+total_obras_proceso_3+total_obras_proceso_4+total_obras_proceso_5+total_obras_proceso_6
    ciento_proceso_1 = (Decimal(total_obras_proceso_1)/Decimal(total_proceso))*100
    ciento_proceso_2 = (Decimal(total_obras_proceso_2)/Decimal(total_proceso))*100
    ciento_proceso_3 = (Decimal(total_obras_proceso_3)/Decimal(total_proceso))*100
    ciento_proceso_4 = (Decimal(total_obras_proceso_4)/Decimal(total_proceso))*100
    ciento_proceso_5 = (Decimal(total_obras_proceso_5)/Decimal(total_proceso))*100
    ciento_proceso_6 = (Decimal(total_obras_proceso_6)/Decimal(total_proceso))*100
    total_ciento_proceso = ciento_proceso_1+ciento_proceso_2+ciento_proceso_3+ciento_proceso_4+ciento_proceso_5+ciento_proceso_6
    #proyectadas
    total_obras_proyectadas_1 = obras_proyectadas_1.count()
    total_obras_proyectadas_2 = obras_proyectadas_2.count()
    total_obras_proyectadas_3 = obras_proyectadas_3.count()
    total_obras_proyectadas_4 = obras_proyectadas_4.count()
    total_obras_proyectadas_5 = obras_proyectadas_5.count()
    total_obras_proyectadas_6 = obras_proyectadas_6.count()
    total_proyectadas = total_obras_proyectadas_1+total_obras_proyectadas_2+total_obras_proyectadas_3+total_obras_proyectadas_4+total_obras_proyectadas_5+total_obras_proyectadas_6
    ciento_proyectadas_1 = (Decimal(total_obras_proyectadas_1)/Decimal(total_proyectadas))*100
    ciento_proyectadas_2 = (Decimal(total_obras_proyectadas_2)/Decimal(total_proyectadas))*100
    ciento_proyectadas_3 = (Decimal(total_obras_proyectadas_3)/Decimal(total_proyectadas))*100
    ciento_proyectadas_4 = (Decimal(total_obras_proyectadas_4)/Decimal(total_proyectadas))*100
    ciento_proyectadas_5 = (Decimal(total_obras_proyectadas_5)/Decimal(total_proyectadas))*100
    ciento_proyectadas_6 = (Decimal(total_obras_proyectadas_6)/Decimal(total_proyectadas))*100
    total_ciento_proyectadas=ciento_proyectadas_1+ciento_proyectadas_2+ciento_proyectadas_3+ciento_proyectadas_4+ciento_proyectadas_5+ciento_proyectadas_6
    #totales
    total_obras_1 = total_obras_concluidas_1+total_obras_proceso_1+total_obras_proyectadas_1
    total_obras_2 = total_obras_concluidas_2+total_obras_proceso_2+total_obras_proyectadas_2
    total_obras_3 = total_obras_concluidas_3+total_obras_proceso_3+total_obras_proyectadas_3
    total_obras_4 = total_obras_concluidas_4+total_obras_proceso_4+total_obras_proyectadas_4
    total_obras_5 = total_obras_concluidas_5+total_obras_proceso_5+total_obras_proyectadas_5
    total_obras_6 = total_obras_concluidas_6+total_obras_proceso_6+total_obras_proyectadas_6
    total_obras = total_obras_1+total_obras_2+total_obras_3+total_obras_4+total_obras_5+total_obras_6
    ciento_obras_1 = (Decimal(total_obras_1)/Decimal(total_obras))*100
    ciento_obras_2 = (Decimal(total_obras_2)/Decimal(total_obras))*100
    ciento_obras_3 = (Decimal(total_obras_3)/Decimal(total_obras))*100
    ciento_obras_4 = (Decimal(total_obras_4)/Decimal(total_obras))*100
    ciento_obras_5 = (Decimal(total_obras_5)/Decimal(total_obras))*100
    ciento_obras_6 = (Decimal(total_obras_6)/Decimal(total_obras))*100
    total_ciento_obras = ciento_obras_1+ciento_obras_2+ciento_obras_3+ciento_obras_4+ciento_obras_5+ciento_obras_6

    for x in range(13,69):
        prs.slides[0].shapes[x].text_frame.paragraphs[0].font.size = Pt(10)

    #concluidas
    prs.slides[0].shapes[13].text= '{0:,}'.format(total_obras_concluidas_1)
    prs.slides[0].shapes[14].text= '{0:,}'.format(total_obras_concluidas_2)
    prs.slides[0].shapes[15].text= '{0:,}'.format(total_obras_concluidas_3)
    prs.slides[0].shapes[16].text= '{0:,}'.format(total_obras_concluidas_4)
    prs.slides[0].shapes[17].text= '{0:,}'.format(total_obras_concluidas_5)
    prs.slides[0].shapes[18].text= '{0:,}'.format(total_obras_concluidas_6)
    prs.slides[0].shapes[19].text= '{0:,}'.format(total_concluidas)
    prs.slides[0].shapes[20].text= '{0:,.2f}'.format(ciento_concluidas_1)
    prs.slides[0].shapes[21].text= '{0:,.2f}'.format(ciento_concluidas_2)
    prs.slides[0].shapes[22].text= '{0:,.2f}'.format(ciento_concluidas_3)
    prs.slides[0].shapes[23].text= '{0:,.2f}'.format(ciento_concluidas_4)
    prs.slides[0].shapes[24].text= '{0:,.2f}'.format(ciento_concluidas_5)
    prs.slides[0].shapes[25].text= '{0:,.2f}'.format(ciento_concluidas_6)
    prs.slides[0].shapes[26].text= '{0:,.2f}'.format(total_ciento_concluidas)
    #proceso
    prs.slides[0].shapes[27].text= '{0:,}'.format(total_obras_proceso_1)
    prs.slides[0].shapes[28].text= '{0:,}'.format(total_obras_proceso_2)
    prs.slides[0].shapes[29].text= '{0:,}'.format(total_obras_proceso_3)
    prs.slides[0].shapes[30].text= '{0:,}'.format(total_obras_proceso_4)
    prs.slides[0].shapes[31].text= '{0:,}'.format(total_obras_proceso_5)
    prs.slides[0].shapes[32].text= '{0:,}'.format(total_obras_proceso_6)
    prs.slides[0].shapes[33].text= '{0:,}'.format(total_proceso)
    prs.slides[0].shapes[34].text= '{0:,.2f}'.format(ciento_proceso_1)
    prs.slides[0].shapes[35].text= '{0:,.2f}'.format(ciento_proceso_2)
    prs.slides[0].shapes[36].text= '{0:,.2f}'.format(ciento_proceso_3)
    prs.slides[0].shapes[37].text= '{0:,.2f}'.format(ciento_proceso_4)
    prs.slides[0].shapes[38].text= '{0:,.2f}'.format(ciento_proceso_5)
    prs.slides[0].shapes[39].text= '{0:,.2f}'.format(ciento_proceso_6)
    prs.slides[0].shapes[40].text= '{0:,.2f}'.format(total_ciento_proceso)
    #proyectadas
    prs.slides[0].shapes[41].text= '{0:,}'.format(total_obras_proyectadas_1)
    prs.slides[0].shapes[42].text= '{0:,}'.format(total_obras_proyectadas_2)
    prs.slides[0].shapes[43].text= '{0:,}'.format(total_obras_proyectadas_3)
    prs.slides[0].shapes[44].text= '{0:,}'.format(total_obras_proyectadas_4)
    prs.slides[0].shapes[45].text= '{0:,}'.format(total_obras_proyectadas_5)
    prs.slides[0].shapes[46].text= '{0:,}'.format(total_obras_proyectadas_6)
    prs.slides[0].shapes[47].text= '{0:,}'.format(total_proyectadas)
    prs.slides[0].shapes[48].text= '{0:,.2f}'.format(ciento_proyectadas_1)
    prs.slides[0].shapes[49].text= '{0:,.2f}'.format(ciento_proyectadas_2)
    prs.slides[0].shapes[50].text= '{0:,.2f}'.format(ciento_proyectadas_3)
    prs.slides[0].shapes[51].text= '{0:,.2f}'.format(ciento_proyectadas_4)
    prs.slides[0].shapes[52].text= '{0:,.2f}'.format(ciento_proyectadas_5)
    prs.slides[0].shapes[53].text= '{0:,.2f}'.format(ciento_proyectadas_6)
    prs.slides[0].shapes[54].text= '{0:,.2f}'.format(total_ciento_proyectadas)
    #totales
    prs.slides[0].shapes[55].text= '{0:,}'.format(total_obras_1)
    prs.slides[0].shapes[56].text= '{0:,}'.format(total_obras_2)
    prs.slides[0].shapes[57].text= '{0:,}'.format(total_obras_3)
    prs.slides[0].shapes[58].text= '{0:,}'.format(total_obras_4)
    prs.slides[0].shapes[59].text= '{0:,}'.format(total_obras_5)
    prs.slides[0].shapes[60].text= '{0:,}'.format(total_obras_6)
    prs.slides[0].shapes[61].text= '{0:,}'.format(total_obras)
    prs.slides[0].shapes[62].text= '{0:,.2f}'.format(ciento_obras_1)
    prs.slides[0].shapes[63].text= '{0:,.2f}'.format(ciento_obras_2)
    prs.slides[0].shapes[64].text= '{0:,.2f}'.format(ciento_obras_3)
    prs.slides[0].shapes[65].text= '{0:,.2f}'.format(ciento_obras_4)
    prs.slides[0].shapes[66].text= '{0:,.2f}'.format(ciento_obras_5)
    prs.slides[0].shapes[67].text= '{0:,.2f}'.format(ciento_obras_6)
    prs.slides[0].shapes[68].text= '{0:,.2f}'.format(total_ciento_obras)



    #f = tempfile.TemporaryFile()
    #prs.save(f)
    #f.seek(0)
    #response = send_file(f, as_attachment=True, attachment_filename='nuevo.pptx', add_etags=False)

    #f.seek(0, os.SEEK_END)
    #size = f.tell()
    #f.seek(0)
    #response.headers.extend({'Content-Length': size,'Cache-Control': 'no-cache' })
    prs.save('/home/obrasapf/djangoObras/obras/static/ppt/ppt-generados/hiper_rangos_' + str(usuario.user.id) + '.pptx')

    the_file = '/home/obrasapf/djangoObras/obras/static/ppt/ppt-generados/hiper_rangos_' + str(usuario.user.id) + '.pptx'
    filename = os.path.basename(the_file)
    chunk_size = 8192
    response = StreamingHttpResponse(FileWrapper(open(the_file,"rb"), chunk_size),
                           content_type=mimetypes.guess_type(the_file)[0])
    response['Content-Length'] = os.path.getsize(the_file)
    response['Content-Disposition'] = "attachment; filename=%s" % filename

    return response

@login_required()
@user_passes_test(is_super_admin)
def logros_resultados_ppt(request):
    prs = Presentation('/home/obrasapf/djangoObras/obras/static/ppt/LOGROS_RESULTADOS.pptx')
    #prs = Presentation('obras/static/ppt/LOGROS_RESULTADOS.pptx')
    usuario = request.user.usuario
    dependencias = usuario.dependencia.all()
    subdependencias = usuario.subdependencia.all()
    query = Q()
    if dependencias and dependencias.count() > 0:
        if usuario.rol == 'US':
            query = Q(subdependencia__in=get_subdependencias_as_list_flat(subdependencias))
        else:

            query = Q(dependencia__in=get_subdependencias_as_list_flat(dependencias)) | Q(subdependencia__in=get_subdependencias_as_list_flat(dependencias)
        )

    obras_concluidas = Obra.objects.filter(
        Q(tipoObra=3),query
    )

    obras_proceso = Obra.objects.filter(
        Q(tipoObra=2),query
    )

    #obras_proyectadas = Obra.objects.filter(
    #    Q(tipoObra=1),Q(dependencia__isnull=True),query
    #)


    obras_proyectadas = Obra.objects.filter(
        Q(tipoObra=1),query
    )

    obras_totales = Obra.objects.filter(
        Q(tipoObra__in=[1,2,3]),query
    )

    total_obras_proyectadas=0
    totalinvertidoproyectadas=0
    total_obras_proceso =0
    totalinvertidoproceso=0
    total_obras_concluidas =0
    totalinvertidoconcluidas =0
    totalObras=totalInvertido=0

    if obras_totales.count()>0:
        totalObras = obras_totales.count()
        total_invertido = obras_totales.aggregate(Sum('inversionTotal'))
        totalInvertido = total_invertido.get('inversionTotal__sum',0)

    if obras_proyectadas.count()>0:
        total_obras_proyectadas = obras_proyectadas.count()
        total_invertido_proyectadas = obras_proyectadas.aggregate(Sum('inversionTotal'))
        totalinvertidoproyectadas = total_invertido_proyectadas.get('inversionTotal__sum',0)
    if obras_proceso.count()>0:
        total_obras_proceso = obras_proceso.count()
        total_invertido_proceso = obras_proceso.aggregate(Sum('inversionTotal'))
        totalinvertidoproceso = total_invertido_proceso.get('inversionTotal__sum',0)

    if obras_concluidas.count()>0:
        total_obras_concluidas = obras_concluidas.count()
        total_invertido_concluidas = obras_concluidas.aggregate(Sum('inversionTotal'))
        totalinvertidoconcluidas = total_invertido_concluidas.get('inversionTotal__sum',0)

    #para resultados por dependencia

    table = prs.slides[0].shapes[0].table
    i=1

    for obra in obras_totales.values('dependencia__nombreDependencia')\
            .annotate(numero_obras=Count('id')).annotate(sumatotal=Sum('inversionTotal')).order_by('dependencia__orden_secretaria'):
        table.cell(i, 1).text_frame.paragraphs[0].font.size = Pt(11)
        table.cell(i, 2).text_frame.paragraphs[0].font.size = Pt(11)
        table.cell(i, 1).text = '{0:,}'.format(obra['numero_obras'])
        table.cell(i, 2).text = '{0:,.2f}'.format(obra['sumatotal'])
        print (obra['dependencia__nombreDependencia'])
        i+=1

    table.cell(i, 1).text_frame.paragraphs[0].font.size = Pt(11)
    table.cell(i, 2).text_frame.paragraphs[0].font.size = Pt(11)
    table.cell(i, 1).text_frame.paragraphs[0].font.color.rgb = RGBColor(0x00, 0x80, 0x00)
    table.cell(i, 2).text_frame.paragraphs[0].font.color.rgb = RGBColor(0x00, 0x80, 0x00)
    table.cell(i, 1).text = '{0:,}'.format(totalObras)
    table.cell(i, 2).text = '{0:,.2f}'.format(totalInvertido)

    tableX = prs.slides[1].shapes[0].table
    tableY = prs.slides[1].shapes[1].table
    i=0
    j=1
    x=0
    y=1
    tableX.cell(0, 0).text_frame.paragraphs[0].font.color.rgb = RGBColor(0x00, 0x00, 0x00)
    tableY.cell(0, 1).text_frame.paragraphs[0].font.color.rgb = RGBColor(0x00, 0x00, 0x00)
    claveSecretaria=1
    for obra in obras_concluidas.values('dependencia__nombreDependencia','dependencia__orden_secretaria')\
            .annotate(numero_obras=Count('id')).annotate(sumatotal=Sum('inversionTotal')).order_by('dependencia__orden_secretaria'):

        #si hay alguna dependencia sin obras en proceso se salta
        while claveSecretaria != obra['dependencia__orden_secretaria']:
            claveSecretaria+=1
            if i<=14:
                i+=2
                j+=2
            else:
                x+=2
                y+=2

        if i<=14:
            tableX.cell(i, 0).text_frame.paragraphs[0].font.size = Pt(8)
            tableX.cell(j, 0).text_frame.paragraphs[0].font.size = Pt(8)
            tableX.cell(i, 0).text = '{0:,}'.format(obra['numero_obras']) + ' Obras'
            tableX.cell(j, 0).text = '{0:,.2f}'.format(obra['sumatotal'])+' mdp'
            i+=2
            j+=2
        else:
            tableY.cell(x, 3).text_frame.paragraphs[0].font.size = Pt(8)
            tableY.cell(y, 3).text_frame.paragraphs[0].font.size = Pt(8)
            tableY.cell(x, 3).text = '{0:,}'.format(obra['numero_obras']) + ' Obras'
            tableY.cell(y, 3).text = '{0:,.2f}'.format(obra['sumatotal'])+' mdp'
            x+=2
            y+=2
        claveSecretaria+=1

    #tableY.cell(x, 2).text_frame.paragraphs[0].font.size = Pt(8)
    #tableY.cell(y, 2).text_frame.paragraphs[0].font.size = Pt(8)
    #tableY.cell(x, 2).text = str('{0:,}'.format(total_obras_concluidas)) + ' Obras'
    #tableY.cell(y, 2).text = str('{0:,.2f}'.format(totalinvertidoconcluidas))+' mdp'


    tableX = prs.slides[2].shapes[0].table
    tableY = prs.slides[2].shapes[1].table
    i=0
    j=1
    x=0
    y=1
    tableX.cell(0, 0).text_frame.paragraphs[0].font.color.rgb = RGBColor(0x00, 0x00, 0x00)
    tableY.cell(0, 1).text_frame.paragraphs[0].font.color.rgb = RGBColor(0x00, 0x00, 0x00)
    claveSecretaria=1
    for obra in obras_proceso.values('dependencia__nombreDependencia','dependencia__orden_secretaria')\
            .annotate(numero_obras=Count('id')).annotate(sumatotal=Sum('inversionTotal')).order_by('dependencia__orden_secretaria'):

        #si hay alguna dependencia sin obras en proceso se salta
        while claveSecretaria != obra['dependencia__orden_secretaria']:
            claveSecretaria+=1
            if i<=14:
                i+=2
                j+=2
            else:
                x+=2
                y+=2

        if i<=14:
            tableX.cell(i, 0).text_frame.paragraphs[0].font.size = Pt(8)
            tableX.cell(j, 0).text_frame.paragraphs[0].font.size = Pt(8)
            tableX.cell(i, 0).text = '{0:,}'.format(obra['numero_obras']) + ' Obras'
            tableX.cell(j, 0).text = '{0:,.2f}'.format(obra['sumatotal'])+' mdp'
            i+=2
            j+=2
        else:
            tableY.cell(x, 3).text_frame.paragraphs[0].font.size = Pt(8)
            tableY.cell(y, 3).text_frame.paragraphs[0].font.size = Pt(8)
            tableY.cell(x, 3).text = '{0:,}'.format(obra['numero_obras']) + ' Obras'
            tableY.cell(y, 3).text = '{0:,.2f}'.format(obra['sumatotal'])+' mdp'
            x+=2
            y+=2
        claveSecretaria+=1

    if x<15: # si Cultura no tiene obras en proceso
        x+=2
        y+=2
    #tableY.cell(x, 2).text_frame.paragraphs[0].font.size = Pt(8)
    #tableY.cell(y, 2).text_frame.paragraphs[0].font.size = Pt(8)
    #tableY.cell(x, 2).text = str('{0:,}'.format(total_obras_proceso)) + ' Obras'
    #tableY.cell(y, 2).text = str('{0:,.2f}'.format(totalinvertidoproceso))+' mdp'


    tableX = prs.slides[3].shapes[0].table
    tableY = prs.slides[3].shapes[1].table
    i=0
    j=1
    x=0
    y=1
    tableX.cell(0, 0).text_frame.paragraphs[0].font.color.rgb = RGBColor(0x00, 0x00, 0x00)
    tableY.cell(0, 1).text_frame.paragraphs[0].font.color.rgb = RGBColor(0x00, 0x00, 0x00)
    claveSecretaria=1
    for obra in obras_proyectadas.values('dependencia__nombreDependencia','dependencia__orden_secretaria')\
            .annotate(numero_obras=Count('id')).annotate(sumatotal=Sum('inversionTotal')).order_by('dependencia__orden_secretaria'):
        print(obra)
        #si hay alguna dependencia sin obras en proyecto se salta
        while claveSecretaria != obra['dependencia__orden_secretaria']:
            claveSecretaria+=1
            if i<=14:
                i+=2
                j+=2
            else:
                x+=2
                y+=2


        if i<=14:
                tableX.cell(i, 0).text_frame.paragraphs[0].font.size = Pt(8)
                tableX.cell(j, 0).text_frame.paragraphs[0].font.size = Pt(8)
                tableX.cell(i, 0).text = '{0:,}'.format(obra['numero_obras']) + ' Obras'
                tableX.cell(j, 0).text = '{0:,.2f}'.format(obra['sumatotal'])+' mdp'
                i+=2
                j+=2
        else:
                tableY.cell(x, 3).text_frame.paragraphs[0].font.size = Pt(8)
                tableY.cell(y, 3).text_frame.paragraphs[0].font.size = Pt(8)
                tableY.cell(x, 3).text = '{0:,}'.format(obra['numero_obras']) + ' Obras'
                tableY.cell(y, 3).text = '{0:,.2f}'.format(obra['sumatotal'])+' mdp'
                x+=2
                y+=2


        claveSecretaria+=1

    if x<15: # si Cultura no tiene obras proyectadas
        x+=2
        y+=2
    #tableY.cell(x, 2).text_frame.paragraphs[0].font.size = Pt(8)
    #tableY.cell(y, 2).text_frame.paragraphs[0].font.size = Pt(8)
    #tableY.cell(x, 2).text = '{0:,}'.format(total_obras_proyectadas) + ' Obras'
    #tableY.cell(y, 2).text = '{0:,.2f}'.format(totalinvertidoproyectadas)+' mdp'



    #para resultados por Estado

    table = prs.slides[4].shapes[0].table
    table2 = prs.slides[4].shapes[1].table
    i=2
    j=2
    obras1=obras_totales.values('estado__nombreEstado')\
            .annotate(numero_obras=Count('id')).annotate(sumatotal=Sum('inversionTotal')).exclude(estado__id=33).exclude(estado__id=34)\
        .order_by('estado__nombreEstado')
    obras2 = obras_totales.filter(Q(estado__id__in=[33,34])).values('estado__nombreEstado')\
            .annotate(numero_obras=Count('id')).annotate(sumatotal=Sum('inversionTotal'))\
        .order_by('estado__nombreEstado')

    obrasU=itertools.chain(obras1,obras2)

    for obra in obrasU:
        if i<20:
            table.cell(i, 1).text_frame.paragraphs[0].font.size = Pt(11)
            table.cell(i, 2).text_frame.paragraphs[0].font.size = Pt(11)
            table.cell(i, 1).text = '{0:,}'.format(obra['numero_obras'])
            table.cell(i, 2).text = '{0:,.2f}'.format(obra['sumatotal'])
            i+=1
        else:
            table2.cell(j, 1).text_frame.paragraphs[0].font.size = Pt(11)
            table2.cell(j, 2).text_frame.paragraphs[0].font.size = Pt(11)
            table2.cell(j, 1).text = '{0:,}'.format(obra['numero_obras'])
            table2.cell(j, 2).text = '{0:,.2f}'.format(obra['sumatotal'])
            j+=1


    table2.cell(j, 1).text_frame.paragraphs[0].font.size = Pt(11)
    table2.cell(j, 2).text_frame.paragraphs[0].font.size = Pt(11)
    table2.cell(j, 1).text_frame.paragraphs[0].font.color.rgb = RGBColor(0x00, 0x80, 0x00)
    table2.cell(j, 2).text_frame.paragraphs[0].font.color.rgb = RGBColor(0x00, 0x80, 0x00)
    table2.cell(j, 1).text = '{0:,}'.format(totalObras)
    table2.cell(j, 2).text = '{0:,.2f}'.format(totalInvertido)



    tableX = prs.slides[5].shapes[0].table
    tableY = prs.slides[5].shapes[1].table
    tableZ = prs.slides[5].shapes[2].table
    i=0
    j=1
    x=0
    y=1
    tableX.cell(0, 0).text_frame.paragraphs[0].font.color.rgb = RGBColor(0x00, 0x00, 0x00)
    tableY.cell(0, 1).text_frame.paragraphs[0].font.color.rgb = RGBColor(0x00, 0x00, 0x00)
    tableZ.cell(0, 1).text_frame.paragraphs[0].font.color.rgb = RGBColor(0x00, 0x00, 0x00)

    obras1=obras_concluidas.values('estado__nombreEstado')\
            .annotate(numero_obras=Count('id')).annotate(sumatotal=Sum('inversionTotal')).exclude(estado__id=33).exclude(estado__id=34)\
        .order_by('estado__nombreEstado')
    obras2 = obras_concluidas.filter(Q(estado__id__in=[33,34])).values('estado__nombreEstado')\
            .annotate(numero_obras=Count('id')).annotate(sumatotal=Sum('inversionTotal'))\
        .order_by('estado__nombreEstado')

    obrasU=itertools.chain(obras1,obras2)

    #for obra in obras_concluidas.values('estado__nombreEstado')\
     #       .annotate(numero_obras=Count('id')).annotate(sumatotal=Sum('inversionTotal')).order_by('estado__nombreEstado'):
    for obra in obrasU:
        if i<=30:
            tableX.cell(i, 0).text_frame.paragraphs[0].font.size = Pt(8)
            tableX.cell(j, 0).text_frame.paragraphs[0].font.size = Pt(8)
            tableX.cell(i, 0).text = '{0:,}'.format(obra['numero_obras']) + ' Obras'
            tableX.cell(j, 0).text = '{0:,.2f}'.format(obra['sumatotal'])+' mdp'
            i+=2
            j+=2
        else:
            if x<=30:
                tableY.cell(x, 1).text_frame.paragraphs[0].font.size = Pt(8)
                tableY.cell(y, 1).text_frame.paragraphs[0].font.size = Pt(8)
                tableY.cell(x, 1).text = '{0:,}'.format(obra['numero_obras']) + ' Obras'
                tableY.cell(y, 1).text = '{0:,.2f}'.format(obra['sumatotal'])+' mdp'
                x+=2
                y+=2
            else:
                tableZ.cell(x-32, 1).text_frame.paragraphs[0].font.size = Pt(8)
                tableZ.cell(y-32, 1).text_frame.paragraphs[0].font.size = Pt(8)
                tableZ.cell(x-32, 1).text = '{0:,}'.format(obra['numero_obras']) + ' Obras'
                tableZ.cell(y-32, 1).text = '{0:,.2f}'.format(obra['sumatotal'])+' mdp'
                x+=2
                y+=2

    tableX = prs.slides[6].shapes[0].table
    tableY = prs.slides[6].shapes[1].table
    tableZ = prs.slides[6].shapes[2].table
    i=0
    j=1
    x=0
    y=1
    tableX.cell(0, 0).text_frame.paragraphs[0].font.color.rgb = RGBColor(0x00, 0x00, 0x00)
    tableY.cell(0, 1).text_frame.paragraphs[0].font.color.rgb = RGBColor(0x00, 0x00, 0x00)
    tableZ.cell(0, 1).text_frame.paragraphs[0].font.color.rgb = RGBColor(0x00, 0x00, 0x00)
    obras1=obras_proceso.values('estado__nombreEstado')\
            .annotate(numero_obras=Count('id')).annotate(sumatotal=Sum('inversionTotal')).exclude(estado__id=33).exclude(estado__id=34)\
        .order_by('estado__nombreEstado')
    obras2 = obras_proceso.filter(Q(estado__id__in=[33,34])).values('estado__nombreEstado')\
            .annotate(numero_obras=Count('id')).annotate(sumatotal=Sum('inversionTotal'))\
        .order_by('estado__nombreEstado')

    obrasU=itertools.chain(obras1,obras2)

    for obra in obrasU:
        if i<=30:
            tableX.cell(i, 0).text_frame.paragraphs[0].font.size = Pt(8)
            tableX.cell(j, 0).text_frame.paragraphs[0].font.size = Pt(8)
            tableX.cell(i, 0).text = '{0:,}'.format(obra['numero_obras']) + ' Obras'
            tableX.cell(j, 0).text = '{0:,.2f}'.format(obra['sumatotal'])+' mdp'
            i+=2
            j+=2
        else:
            if x<=30:
                tableY.cell(x, 1).text_frame.paragraphs[0].font.size = Pt(8)
                tableY.cell(y, 1).text_frame.paragraphs[0].font.size = Pt(8)
                tableY.cell(x, 1).text = '{0:,}'.format(obra['numero_obras']) + ' Obras'
                tableY.cell(y, 1).text = '{0:,.2f}'.format(obra['sumatotal'])+' mdp'
                x+=2
                y+=2
            else:
                tableZ.cell(x-32, 1).text_frame.paragraphs[0].font.size = Pt(8)
                tableZ.cell(y-32, 1).text_frame.paragraphs[0].font.size = Pt(8)
                tableZ.cell(x-32, 1).text = '{0:,}'.format(obra['numero_obras']) + ' Obras'
                tableZ.cell(y-32, 1).text = '{0:,.2f}'.format(obra['sumatotal'])+' mdp'
                x+=2
                y+=2

    tableX = prs.slides[7].shapes[0].table
    tableY = prs.slides[7].shapes[1].table
    tableZ = prs.slides[7].shapes[2].table
    i=0
    j=1
    x=0
    y=1
    tableX.cell(0, 0).text_frame.paragraphs[0].font.color.rgb = RGBColor(0x00, 0x00, 0x00)
    tableY.cell(0, 1).text_frame.paragraphs[0].font.color.rgb = RGBColor(0x00, 0x00, 0x00)
    tableZ.cell(0, 1).text_frame.paragraphs[0].font.color.rgb = RGBColor(0x00, 0x00, 0x00)
    obras1=obras_proyectadas.values('estado__nombreEstado')\
            .annotate(numero_obras=Count('id')).annotate(sumatotal=Sum('inversionTotal')).exclude(estado__id=33).exclude(estado__id=34)\
        .order_by('estado__nombreEstado')
    obras2 = obras_proyectadas.filter(Q(estado__id__in=[33,34])).values('estado__nombreEstado')\
            .annotate(numero_obras=Count('id')).annotate(sumatotal=Sum('inversionTotal'))\
        .order_by('estado__nombreEstado')

    obrasU=itertools.chain(obras1,obras2)

    for obra in obrasU:
        if i<=30:
            tableX.cell(i, 0).text_frame.paragraphs[0].font.size = Pt(8)
            tableX.cell(j, 0).text_frame.paragraphs[0].font.size = Pt(8)
            tableX.cell(i, 0).text = '{0:,}'.format(obra['numero_obras']) + ' Obras'
            tableX.cell(j, 0).text = '{0:,.2f}'.format(obra['sumatotal'])+' mdp'
            i+=2
            j+=2
        else:
            if x<=30:
                tableY.cell(x, 1).text_frame.paragraphs[0].font.size = Pt(8)
                tableY.cell(y, 1).text_frame.paragraphs[0].font.size = Pt(8)
                tableY.cell(x, 1).text = '{0:,}'.format(obra['numero_obras']) + ' Obras'
                tableY.cell(y, 1).text = '{0:,.2f}'.format(obra['sumatotal'])+' mdp'
                x+=2
                y+=2
            else:
                tableZ.cell(x-32, 1).text_frame.paragraphs[0].font.size = Pt(8)
                tableZ.cell(y-32, 1).text_frame.paragraphs[0].font.size = Pt(8)
                tableZ.cell(x-32, 1).text = '{0:,}'.format(obra['numero_obras']) + ' Obras'
                tableZ.cell(y-32, 1).text = '{0:,.2f}'.format(obra['sumatotal'])+' mdp'
                x+=2
                y+=2

    prs.save('/home/obrasapf/djangoObras/obras/static/ppt/ppt-generados/logros_resultados_' + str(usuario.user.id) + '.pptx')
    the_file = '/home/obrasapf/djangoObras/obras/static/ppt/ppt-generados/logros_resultados_' + str(usuario.user.id) + '.pptx'

    #prs.save('obras/static/ppt/ppt-generados/logros_resultados_' + str(usuario.user.id) + '.pptx')
    #the_file = 'obras/static/ppt/ppt-generados/logros_resultados_' + str(usuario.user.id) + '.pptx'

    filename = os.path.basename(the_file)
    chunk_size = 8192
    response = StreamingHttpResponse(FileWrapper(open(the_file,"rb"), chunk_size),
                           content_type=mimetypes.guess_type(the_file)[0])
    response['Content-Length'] = os.path.getsize(the_file)
    response['Content-Disposition'] = "attachment; filename=%s" % filename

    return response






def redirect_admin(request):
    return redirect('admin/')

@login_required
def ios_view(request):
    return render_to_response('ios.html')
    #return HttpResponse('<a href="itms-services://?action=download-manifest&'
    #                    'url=https://obrasapf.mx/static/'
    #                    'manifest.plist">'
    #                    'Instalar App</a>')